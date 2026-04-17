#!/usr/bin/env python3
"""Build detailed PPTX for 15_tacit_knowledge_lifecycle (~8 slides).

Dense handout version with all tables, full detail, and comprehensive explanations.
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TEMPLATE = "project-rikyu-sales-proposals-poc/untracked/original_document/report/\u63d0\u6848\u4eee\u8aac\u69cb\u7bc9AI_\u5b9a\u4f8b\u4f1a1_20260219_v1.pptx"
WORKDIR = "output/15_tacit_lifecycle"

# === ACES Color Palette ===
GREEN_DARK = RGBColor(0x00, 0x68, 0x43)
GREEN_MEDIUM = RGBColor(0x19, 0x7A, 0x56)
GREEN_BRIGHT = RGBColor(0x00, 0x9A, 0x62)
GREEN_MINT = RGBColor(0xC2, 0xE7, 0xD9)
GREEN_PALE = RGBColor(0xDF, 0xE6, 0xE0)
CHARCOAL = RGBColor(0x39, 0x39, 0x39)
TEXT_MAIN = RGBColor(0x12, 0x12, 0x12)
TABLE_TEXT = RGBColor(0x45, 0x45, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x52, 0x98, 0xBA)
ACCENT_RED = RGBColor(0xDB, 0x5F, 0x5F)
ACCENT_ORANGE = RGBColor(0xFB, 0xAE, 0x40)
ACCENT_GOLD = RGBColor(0xC0, 0xB7, 0x6A)
GRAY_LIGHT = RGBColor(0x99, 0x99, 0x99)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

TITLE_FONT = "Yu Gothic"
SUB_FONT = "Yu Gothic Medium"
BODY_FONT = "Century Gothic"


# ── Helper functions ──────────────────────────────────

def find_shapes(slide):
    shapes_with_y = []
    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            y = s.top / 914400 if s.top else 0
            shapes_with_y.append((y, s))
    shapes_with_y.sort(key=lambda x: x[0])
    title = subtitle = body = None
    for y, s in shapes_with_y:
        if y < 0.7 and title is None: title = s
        elif 0.7 <= y < 1.5 and subtitle is None: subtitle = s
        elif y >= 1.5 and body is None: body = s
    return title, subtitle, body


def prep(slide):
    title, subtitle, body = find_shapes(slide)
    if subtitle:
        subtitle.height = Inches(0.32)
        if subtitle.width < Inches(1):
            subtitle.width = Inches(12.46)
            subtitle.left = Inches(0.42)
    if body:
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = " "
        p.font.size = Pt(1)
        p.font.color.rgb = WHITE
    return title, subtitle, body


def add_table(slide, data, left=Inches(0.42), top=Inches(1.90), width=Inches(12.46),
              col_widths=None, header_font_pt=8, cell_font_pt=7, row_height=0.28):
    rows, cols = len(data), len(data[0])
    height = Inches(row_height * rows)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.name = BODY_FONT
                if row_idx == 0:
                    para.font.size = Pt(header_font_pt)
                    para.font.color.rgb = WHITE
                    para.font.bold = True
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.font.size = Pt(cell_font_pt)
                    para.font.color.rgb = TABLE_TEXT
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_DARK
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_PALE
    return table


def add_box(slide, text, left, top, width, height, bg_color=GREEN_MINT, text_color=GREEN_DARK,
            font_size=9, bold=False, align=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_arrow_down(slide, left, top, width=Inches(0.4), height=Inches(0.18)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u25bc"
    p.font.size = Pt(12)
    p.font.color.rgb = GREEN_DARK
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.25)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u2192"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_note(slide, text, left, top, width, h=Inches(0.30)):
    txBox = slide.shapes.add_textbox(left, top, width, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(6)
    p.font.italic = True
    p.font.color.rgb = GRAY_LIGHT
    p.font.name = BODY_FONT


def set_title(slide, title_text, subtitle_text):
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = title_text
        p.font.name = TITLE_FONT
        p.font.size = Pt(20)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = subtitle_text
        p.font.name = SUB_FONT
        p.font.size = Pt(10)


# ── Build ──────────────────────────────────────────

def build():
    """Build detailed version: Title + 7 dense content slides."""
    prs = Presentation(TEMPLATE)

    layout_cover_data = prs.slide_layouts[1]
    layout_content = prs.slide_layouts[10]

    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    # 8 slides: Title + 7 content
    layouts = [layout_cover_data] + [layout_content] * 7
    for layout in layouts:
        new_prs.slides.add_slide(layout)

    slides = new_prs.slides

    # ── Slide 0: Title ──
    slide = slides[0]
    shapes_with_y = []
    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            y = s.top / 914400 if s.top else 0
            shapes_with_y.append((y, s))
    shapes_with_y.sort(key=lambda x: x[0])
    shape_list = [s for _, s in shapes_with_y]
    if len(shape_list) >= 1:
        tf = shape_list[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "\u682a\u5f0f\u4f1a\u793e\u308a\u305d\u306a\u9280\u884c \u5fa1\u4e2d"
        p.font.size = Pt(32)
        p2 = tf.add_paragraph()
        p2.text = "\u6697\u9ed9\u77e5\u306e\u30c7\u30fc\u30bf\u5316\u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb\u3010\u8a73\u7d30\u7248\u3011"
        p2.font.size = Pt(28)
    if len(shape_list) >= 2:
        tf = shape_list[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "2026/3/10 | rikyu \u30bd\u30ea\u30e5\u30fc\u30b7\u30e7\u30f3\u30bb\u30fc\u30eb\u30b9\u4f34\u8d70AI PoC"
        p.font.size = Pt(14)

    # ── Slide 1: \u6697\u9ed9\u77e5\u3068\u306f + \u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb\u5168\u4f53\u50cf ──
    slide = slides[1]
    set_title(slide, "\u6697\u9ed9\u77e5\u3068\u306f\u4f55\u304b + \u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb\u5168\u4f53\u50cf",
              "\u300c\u72b6\u614b\u300d\u3068\u3057\u3066\u306e\u6697\u9ed9\u77e5\u3068\u30015\u6bb5\u968e\u306e\u5909\u63db\u30d5\u30ed\u30fc")

    # Left: 3 examples (compact)
    add_box(slide, "\u6697\u9ed9\u77e5 = \u300c\u72b6\u614b\u300d\uff08\u672a\u8a18\u9332\u30fb\u672a\u8a00\u8a9e\u5316\uff09",
            Inches(0.42), Inches(1.70), Inches(5.8), Inches(0.22),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    examples_data = [
        ["\u6697\u9ed9\u77e5\u306e\u4f8b", "\u672a\u8a18\u9332\u306a\u3089", "\u8a18\u9332/\u8a00\u8a9e\u5316\u3057\u305f\u3089"],
        ["\u300cA\u793e\u306e\u7dcf\u52d9\u90e8\u9577\u3068\u5c02\u52d9\u306f\u4e0d\u4ef2\u300d", "\u6697\u9ed9\u77e5", "Data (KP\u30de\u30c3\u30d7)"],
        ["\u300c\u7269\u6d41\u696d\u306f\u4eba\u4ef6\u8cbb\u304cKey Driver\u300d", "\u6697\u9ed9\u77e5", "Knowledge (\u30eb\u30fc\u30eb)"],
        ["\u300c\u3053\u306e\u51fa\u529b\u306f\u7684\u5916\u308c\u3060\u300d\u3068\u3044\u3046\u76f4\u611f", "\u6697\u9ed9\u77e5", "Knowledge (\u4eee\u8aac\u4fee\u6b63\u30eb\u30fc\u30eb)"],
    ]
    add_table(slide, examples_data, left=Inches(0.42), top=Inches(1.95), width=Inches(5.8),
              col_widths=[2.5, 1.2, 2.1],
              header_font_pt=7, cell_font_pt=6.5, row_height=0.24)

    add_note(slide, "\u6697\u9ed9\u77e5\u306f\u300c\u5225\u306e\u60c5\u5831\u300d\u3067\u306f\u306a\u304f\u3001\u307e\u3060Data\u306b\u3082Knowledge\u306b\u3082\u306a\u3063\u3066\u3044\u306a\u3044\u300c\u72b6\u614b\u300d\u3002\u8a18\u9332\u30fb\u8a00\u8a9e\u5316\u3059\u308c\u3070\u6d3b\u7528\u53ef\u80fd\u306b\u306a\u308b\u3002",
             Inches(0.42), Inches(2.95), Inches(5.8))

    # Right: 5-stage lifecycle (compact vertical)
    add_box(slide, "5\u6bb5\u968e\u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb",
            Inches(6.50), Inches(1.70), Inches(6.3), Inches(0.22),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    stages = [
        ("\u2460 \u30ad\u30e3\u30d7\u30c1\u30e3", "\u8ffd\u52a0\u64cd\u4f5c\u306a\u3057\u306b\u6697\u9ed9\u77e5\u3092\u8a18\u9332", GREEN_MINT),
        ("\u2461 \u84c4\u7a4d", "\u64cd\u4f5c\u30fb\u5224\u65ad\u306e\u30d1\u30bf\u30fc\u30f3\u691c\u51fa", GREEN_PALE),
        ("\u2462 \u6607\u683c", "Data\u307e\u305f\u306fKnowledge\u306b\u5909\u63db", GREEN_MINT),
        ("\u2463 \u5c55\u958b", "\u500b\u4eba\u2192\u55b6\u696d\u5e97\u2192\u696d\u754c\u2192\u5168\u793e\u3067\u5171\u6709", GREEN_PALE),
        ("\u2464 \u9032\u5316", "Knowledge\u304cAI\u306e\u5224\u65ad\u3092\u6539\u5584", GREEN_MINT),
    ]
    y_st = 1.95
    for i, (label, desc, bg) in enumerate(stages):
        add_box(slide, f"{label}: {desc}",
                Inches(6.50), Inches(y_st), Inches(6.3), Inches(0.22),
                bg_color=bg, text_color=CHARCOAL, font_size=7)
        y_st += 0.22
        if i < len(stages) - 1:
            add_arrow_down(slide, Inches(9.4), Inches(y_st - 0.02), Inches(0.4), Inches(0.14))
            y_st += 0.11

    add_box(slide, "\u21b3 \u2464\u9032\u5316 \u2192 \u2460\u30ad\u30e3\u30d7\u30c1\u30e3\u306b\u623b\u308b\uff08\u6b63\u306e\u30b9\u30d1\u30a4\u30e9\u30eb\uff09",
            Inches(6.50), Inches(y_st + 0.05), Inches(6.3), Inches(0.20),
            bg_color=ACCENT_GOLD, text_color=CHARCOAL, font_size=7, bold=True)

    # ── Slide 2: \u2460\u30ad\u30e3\u30d7\u30c1\u30e3 \u8a73\u7d30 ──
    slide = slides[2]
    set_title(slide, "\u2460\u30ad\u30e3\u30d7\u30c1\u30e3: 5\u3064\u306e\u30bf\u30c3\u30c1\u30dd\u30a4\u30f3\u30c8 + \u300c\u8ffd\u52a0\u64cd\u4f5c\u306a\u3057\u300d\u539f\u5247",
              "RM\u304c\u666e\u6bb5\u901a\u308a\u306e\u696d\u52d9\u3092\u3059\u308b\u3060\u3051\u3067\u6697\u9ed9\u77e5\u304c\u81ea\u52d5\u7684\u306b\u8a18\u9332\u3055\u308c\u308b\u4ed5\u7d44\u307f")

    tp_data = [
        ["#", "\u30bf\u30c3\u30c1\u30dd\u30a4\u30f3\u30c8", "RM\u306e\u64cd\u4f5c", "\u30ad\u30e3\u30d7\u30c1\u30e3\u3055\u308c\u308b\u6697\u9ed9\u77e5", "\u8a18\u9332\u5148"],
        ["C1", "\u9762\u8ac7\u8a18\u9332\u5165\u529b", "\u9762\u8ac7\u5f8c\u306b\u8a18\u9332\u3092\u5165\u529b", "\u4ed6\u884c\u6761\u4ef6, KP\u95a2\u4fc2\u6027, \u8a2d\u5099\u72b6\u6cc1,\n\u30aa\u30fc\u30ca\u30fc\u60c5\u5831, \u7af6\u5408\u52d5\u5411", "\u9762\u8ac7DB"],
        ["C2", "AI\u66f4\u65b0\u63d0\u6848\u30ec\u30d3\u30e5\u30fc", "\u627f\u8a8d/\u4fee\u6b63/\u68c4\u5374+\u7406\u7531", "\u300c\u66f4\u65b0\u3059\u3079\u304d\u300d\u306e\u52d8\u6240,\n\u300c\u4fe1\u983c\u5ea6\u304c\u9055\u3046\u300d\u306e\u611f\u899a", "\u627f\u8a8d\u30ed\u30b0DB"],
        ["C3", "AI\u30c1\u30e3\u30c3\u30c8\u5229\u7528", "\u8cea\u554f+FB(\u7684\u5916\u308c\u30dc\u30bf\u30f3)", "\u300c\u7684\u5916\u308c\u300d\u306e\u7406\u7531,\n\u3088\u304f\u805e\u304b\u308c\u308b\u8cea\u554f\u30d1\u30bf\u30fc\u30f3", "\u30c1\u30e3\u30c3\u30c8\u30ed\u30b0DB"],
        ["C4", "\u6210\u7d04/\u5931\u6ce8\u8a18\u9332", "\u7d50\u679c+\u8981\u56e0\u3092\u8a18\u9332", "\u300c\u523a\u3055\u3063\u305f/\u523a\u3055\u3089\u306a\u304b\u3063\u305f\u300d\u8a18\u61b6,\n\u5546\u54c1\u9078\u5b9a\u7406\u7531", "\u63d0\u6848DB"],
        ["C5", "\u30a8\u30ad\u30b9\u30d1\u30fc\u30c8\u30d2\u30a2\u30ea\u30f3\u30b0", "\u30d2\u30a2\u30ea\u30f3\u30b0\u306b\u56de\u7b54", "\u696d\u754c\u30d1\u30bf\u30fc\u30f3, \u5546\u8ac7\u9032\u884c\u30b3\u30c4,\nNG\u4e8b\u9805", "Knowledge Master"],
    ]
    add_table(slide, tp_data, top=Inches(1.70),
              col_widths=[0.40, 1.80, 1.80, 4.50, 1.30],
              header_font_pt=7, cell_font_pt=6, row_height=0.42)

    # Zero effort comparison (bottom)
    add_box(slide, "\u2717 \u60aa\u3044\u8a2d\u8a08: \u300c\u6697\u9ed9\u77e5\u3092\u8a18\u9332\u3057\u3066\u304f\u3060\u3055\u3044\u300d \u2192 RM\u8ca0\u8377\u589e \u2192 \u4f7f\u308f\u308c\u306a\u3044",
            Inches(0.42), Inches(4.40), Inches(5.8), Inches(0.28),
            bg_color=ACCENT_RED, text_color=WHITE, font_size=7, bold=True, align=PP_ALIGN.LEFT)

    add_box(slide, "\u2713 \u826f\u3044\u8a2d\u8a08: \u666e\u6bb5\u306e\u64cd\u4f5c\uff08\u9762\u8ac7\u8a18\u9332\u3001\u627f\u8a8d\u3001\u30c1\u30e3\u30c3\u30c8\uff09\u306e\u4e2d\u306b\u69cb\u9020\u5316\u9805\u76ee\u3092\u57cb\u3081\u8fbc\u3080",
            Inches(6.50), Inches(4.40), Inches(6.3), Inches(0.28),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=7, bold=True, align=PP_ALIGN.LEFT)

    add_note(slide, "\u4f8b: \u9762\u8ac7\u8a18\u9332\u753b\u9762\u306e\u300c\u4ed6\u884c\u60c5\u5831\u300d\u300cKP\u95a2\u4fc2\u6027\u300d\u9805\u76ee\u306f\u3001RM\u304c\u9762\u8ac7\u3067\u805e\u3044\u305f\u60c5\u5831\u3092\u81ea\u7136\u306b\u8a18\u9332\u3059\u308b\u5c0e\u7dda\u3002\u300c\u6697\u9ed9\u77e5\u5165\u529b\u300d\u5c02\u7528\u753b\u9762\u306f\u4f5c\u3089\u306a\u3044\u3002",
             Inches(0.42), Inches(4.75), Inches(12.46))

    # ── Slide 3: \u2461\u84c4\u7a4d\u2192\u2462\u6607\u683c Path A ──
    slide = slides[3]
    set_title(slide, "\u2461\u84c4\u7a4d\u2192\u2462\u6607\u683c: \u30d1\u30b9A\uff08\u5373Data\uff09+ \u30d1\u30b9B\uff08Knowledge\u5316\uff09",
              "2\u3064\u306e\u6607\u683c\u30d1\u30b9: \u5373\u6642Data\u5316\uff085\u9805\u76ee\uff09\u3068\u3001\u30d1\u30bf\u30fc\u30f3\u84c4\u7a4d\u306b\u3088\u308bKnowledge\u5316\uff088\u9805\u76ee\uff09")

    # Path A
    add_box(slide, "\u30d1\u30b9A: \u8a18\u9332\u3059\u308c\u3070\u5373Data",
            Inches(0.42), Inches(1.65), Inches(12.46), Inches(0.20),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=7, bold=True)

    path_a = [
        ["\u6697\u9ed9\u77e5", "\u8a18\u9332\u65b9\u6cd5", "\u6607\u683c\u5148", "\u6d3b\u7528\u30bf\u30a4\u30df\u30f3\u30b0"],
        ["\u9762\u8ac7\u3067\u805e\u3044\u305f\u4ed6\u884c\u6761\u4ef6", "\u300c\u4ed6\u884c\u60c5\u5831\u300d\u9805\u76ee", "Data: \u4ed6\u884c\u63d0\u793a\u6761\u4ef6", "\u6b21\u56de\u306e\u63d0\u6848\u6642"],
        ["KP\u9593\u306e\u4eba\u9593\u95a2\u4fc2\u30fb\u529b\u5b66", "\u300cKP\u95a2\u4fc2\u6027\u300d\u9805\u76ee", "Data: \u5f79\u54e1\u30fbKP", "KP\u30de\u30c3\u30d7\u306b\u5373\u53cd\u6620"],
        ["\u8a2d\u5099\u30fb\u6295\u8cc7\u306e\u73fe\u6cc1", "\u300c\u8a2d\u5099\u30fb\u6295\u8cc7\u300d\u9805\u76ee", "Data: \u4e8b\u696d\u69cb\u9020", "\u30a2\u30b8\u30a7\u30f3\u30c0\u66f4\u65b0\u63d0\u6848"],
        ["\u7d4c\u55b6\u8005\u306e\u500b\u4eba\u7684\u72b6\u6cc1", "\u300c\u30aa\u30fc\u30ca\u30fc\u60c5\u5831\u300d\u9805\u76ee", "Data: \u7d4c\u55b6\u8005\u57fa\u672c\u60c5\u5831", "\u4e8b\u696d\u627f\u7d99\u5224\u5b9a"],
        ["\u7af6\u5408\u653b\u52e2\u306e\u60c5\u5831", "\u300c\u7af6\u5408\u52d5\u5411\u300d\u9805\u76ee", "Data: \u7af6\u5408\u653b\u52e2\u60c5\u5831", "\u63d0\u6848\u6226\u7565"],
    ]
    add_table(slide, path_a, left=Inches(0.42), top=Inches(1.88), width=Inches(12.46),
              col_widths=[3.0, 2.5, 3.0, 3.0],
              header_font_pt=7, cell_font_pt=6, row_height=0.22)

    # Path B
    y_b = 1.88 + 0.22 * 7 + 0.08
    add_box(slide, "\u30d1\u30b9B: \u30d1\u30bf\u30fc\u30f3\u84c4\u7a4d \u2192 Knowledge\u5316",
            Inches(0.42), Inches(y_b), Inches(12.46), Inches(0.20),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=7, bold=True)

    path_b = [
        ["\u6697\u9ed9\u77e5", "\u84c4\u7a4d\u65b9\u6cd5", "\u6607\u683c\u6761\u4ef6", "\u6607\u683c\u5148Knowledge"],
        ["\u300c\u3053\u306e\u767a\u8a00\u306a\u3089\u66f4\u65b0\u3059\u3079\u304d\u300d", "\u627f\u8a8d/\u4fee\u6b63/\u68c4\u5374\u30d1\u30bf\u30fc\u30f3", "\u540c\u4e003\u56de", "K2: \u66f4\u65b0\u30eb\u30fc\u30eb"],
        ["\u300c\u3053\u306e\u4fe1\u983c\u5ea6\u306f\u9055\u3046\u300d", "\u4fe1\u983c\u5ea6\u4e0a\u66f8\u304d\u4fee\u6b63", "\u540c\u4e00\u4fee\u6b63\u84c4\u7a4d", "K3: \u4fe1\u983c\u5ea6\u30de\u30c8\u30ea\u30af\u30b9"],
        ["\u300c\u3053\u3046\u805e\u304f\u3068\u524d\u306b\u9032\u3080\u300d", "\u6210\u529f\u30a2\u30b8\u30a7\u30f3\u30c0\u5206\u6790", "\u6210\u529f\u30d1\u30bf\u30fc\u30f3\u62bd\u51fa", "K6: \u5546\u8ac7\u9032\u884c\u30d1\u30bf\u30fc\u30f3"],
        ["\u300c\u3053\u306e\u9867\u5ba2\u306b\u306f\u3053\u3046\u5165\u308b\u300d", "\u30c1\u30e3\u30c3\u30c8FB", "PJ\u30c1\u30fc\u30e0\u30ec\u30d3\u30e5\u30fc", "K6: \u9867\u5ba2\u30bf\u30a4\u30d7\u5225"],
        ["\u300c\u3053\u306e\u696d\u754c\u306f\u4eca\u3053\u3046\u3044\u3046\u7a7a\u6c17\u300d", "\u30c1\u30e3\u30c3\u30c8FB", "\u8907\u6570RM\u540c\u50be\u5411", "K1: \u696d\u754c\u5225\u8ab2\u984c"],
        ["\u300c\u3053\u306e\u8ab2\u984c\u306b\u306f\u3053\u306e\u5546\u54c1\u300d", "\u6210\u7d04\u6642\u306e\u9078\u5b9a\u7406\u7531", "\u6210\u529f\u30d1\u30bf\u30fc\u30f3\u84c4\u7a4d", "K4: \u8ab2\u984c\u00d7\u5546\u54c1"],
        ["\u300c\u523a\u3055\u3063\u305f/\u523a\u3055\u3089\u306a\u304b\u3063\u305f\u300d", "\u6210\u7d04/\u5931\u6ce8\u8981\u56e0", "\u4e8b\u4f8b\u69cb\u9020\u5316", "K5: \u6210\u7d04/\u5931\u6ce8\u30d1\u30bf\u30fc\u30f3"],
        ["\u300c\u3053\u306e\u51fa\u529b\u306f\u7684\u5916\u308c\u300d", "\u30c1\u30e3\u30c3\u30c8FB\u30dc\u30bf\u30f3", "\u540c\u7a2eFB\u84c4\u7a4d", "K2: \u91cd\u8981\u5ea6\u5224\u5b9a"],
    ]
    add_table(slide, path_b, left=Inches(0.42), top=Inches(y_b + 0.23), width=Inches(12.46),
              col_widths=[3.0, 2.5, 2.0, 3.0],
              header_font_pt=7, cell_font_pt=6, row_height=0.22)

    add_note(slide, "\u30d1\u30b9A: \u69cb\u9020\u5316\u9805\u76ee\u5165\u529b\u2192\u5373Data\u5316\u3002\u30d1\u30b9B: \u540c\u4e00\u30d1\u30bf\u30fc\u30f33\u56de\u84c4\u7a4d\u2192\u30d1\u30bf\u30fc\u30f3\u8a00\u8a9e\u5316\u2192Knowledge Master\u306b\u30eb\u30fc\u30eb\u8ffd\u52a0\u2192\u5168RM\u306eAI\u51fa\u529b\u54c1\u8cea\u5411\u4e0a",
             Inches(0.42), Inches(y_b + 0.23 + 0.22 * 9 + 0.05), Inches(12.46))

    # ── Slide 4: \u2463\u5c55\u958b 4\u5c64\u30e1\u30e2\u30ea\u968e\u5c64 ──
    slide = slides[4]
    set_title(slide, "\u2463\u5c55\u958b: 4\u5c64\u30e1\u30e2\u30ea\u968e\u5c64 + \u6607\u683c\u30eb\u30fc\u30c8",
              "\u500b\u4eba\u2192\u55b6\u696d\u5e97\u2192\u696d\u754c\u2192\u5168\u793e\u306e\u30eb\u30fc\u30c8\u3067\u30ca\u30ec\u30c3\u30b8\u3092\u5171\u6709\u30fb\u6607\u683c\u3059\u308b")

    # Nested boxes
    add_box(slide, "Level 4: \u5168\u793e\uff08Global\uff09 \u2014 \u5168\u793e\u5171\u901a\u306e\u30eb\u30fc\u30eb\u3002\u300c\u5168\u696d\u754c\u3067\u88dc\u52a9\u91d1\u306f\u523a\u3055\u308b\u300d\u7b49",
            Inches(0.42), Inches(1.70), Inches(7.5), Inches(3.50),
            bg_color=RGBColor(0xE8, 0xF0, 0xE8), text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)

    add_box(slide, "Level 3: \u696d\u754c\uff08Industry\uff09 \u2014 \u300c\u88fd\u9020\u696d\u306f\u8a2d\u5099\u6295\u8cc7\u30b5\u30a4\u30af\u30eb5-7\u5e74\u300d",
            Inches(0.72), Inches(2.30), Inches(6.5), Inches(2.65),
            bg_color=RGBColor(0xD0, 0xE4, 0xD0), text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)

    add_box(slide, "Level 2: \u55b6\u696d\u5e97\uff08Branch\uff09 \u2014 \u300c\u3053\u306e\u5730\u57df\u306f\u5efa\u8a2d\u696d\u304c\u591a\u3044\u300d",
            Inches(1.02), Inches(2.85), Inches(5.5), Inches(1.80),
            bg_color=RGBColor(0xC2, 0xE7, 0xD9), text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)

    add_box(slide, "Level 1: \u500b\u4eba\uff08RM\uff09 \u2014 \u300cA\u793e\u306f\u6570\u5b57\u3088\u308a\u4eba\u60c5\u300d",
            Inches(1.32), Inches(3.35), Inches(4.5), Inches(0.80),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=7, bold=True, align=PP_ALIGN.LEFT)

    # Customer memory
    add_box(slide, "\u2605 \u9867\u5ba2\u30e1\u30e2\u30ea\uff08\u6a2a\u65ad\u7684\u95a2\u5fc3\u4e8b\uff09\n\u9867\u5ba2\u5358\u4f4d\u306e\u60c5\u5831\u306f\u5168\u30ec\u30d9\u30eb\u3092\u6a2a\u65ad",
            Inches(0.42), Inches(5.35), Inches(7.5), Inches(0.40),
            bg_color=ACCENT_BLUE, text_color=WHITE, font_size=7, bold=True, align=PP_ALIGN.LEFT)

    # Promotion route table
    add_box(slide, "\u6607\u683c\u30eb\u30fc\u30c8\u4f8b\uff08\u8a73\u7d30\uff09",
            Inches(8.20), Inches(1.70), Inches(4.6), Inches(0.20),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=7, bold=True)

    route_data = [
        ["\u30ec\u30d9\u30eb", "\u4f8b", "\u6607\u683c\u6761\u4ef6"],
        ["\u500b\u4ebaRM", "\u300cA\u793e\u306f\u6570\u5b57\u3067\u8aac\u5f97\u3059\u3079\u304d\u300d", "RM\u500b\u4eba\u306e\u7d4c\u9a13\u3068\u3057\u3066\u84c4\u7a4d"],
        ["\u2193 \u55b6\u696d\u5e97", "\u300c\u88fd\u9020\u696d\u306e\u9867\u5ba2\u306f\u88dc\u52a9\u91d1\u306b\u53cd\u5fdc\u300d", "\u540c\u3058\u55b6\u696d\u5e97\u306e\u8907\u6570RM\u304c\u540c\u50be\u5411FB"],
        ["\u2193 \u696d\u754c", "\u300c\u88fd\u9020\u696d\u306f\u8a2d\u5099\u6295\u8cc7\u30b5\u30a4\u30af\u30eb\u3067\n\u63d0\u6848\u30bf\u30a4\u30df\u30f3\u30b0\u304c\u6c7a\u307e\u308b\u300d", "\u8907\u6570\u55b6\u696d\u5e97\u3067\u540c\u50be\u5411"],
        ["\u2193 \u5168\u793e", "\u5168\u696d\u754c\u5171\u901a\u30eb\u30fc\u30eb\u5316", "\u4ed6\u696d\u754c\u3067\u3082\u540c\u69cb\u9020\u3092\u78ba\u8a8d"],
    ]
    add_table(slide, route_data, left=Inches(8.20), top=Inches(1.93), width=Inches(4.6),
              col_widths=[0.80, 2.00, 1.80],
              header_font_pt=6.5, cell_font_pt=6, row_height=0.32)

    # ── Slide 5: \u2464\u9032\u5316 AI-Human Gap ──
    slide = slides[5]
    set_title(slide, "\u2464\u9032\u5316: AI-Human Gap = \u6700\u9ad8\u306e\u6559\u5e2b\u30c7\u30fc\u30bf",
              "RM\u304cAI\u51fa\u529b\u3092\u4fee\u6b63\u3059\u308b\u305f\u3073\u3001\u305d\u308c\u304c\u6700\u3082\u4fa1\u5024\u3042\u308b\u5b66\u7fd2\u30c7\u30fc\u30bf\u306b\u306a\u308b")

    # Gap diagram
    add_box(slide, "AI\u51fa\u529b",
            Inches(0.42), Inches(1.70), Inches(2.5), Inches(0.35),
            bg_color=ACCENT_BLUE, text_color=WHITE, font_size=10, bold=True)
    add_box(slide, "\u2190 Gap \u2192",
            Inches(3.10), Inches(1.70), Inches(1.8), Inches(0.35),
            bg_color=ACCENT_ORANGE, text_color=WHITE, font_size=10, bold=True)
    add_box(slide, "RM\u306e\u4fee\u6b63",
            Inches(5.10), Inches(1.70), Inches(2.5), Inches(0.35),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=10, bold=True)

    # Gap meaning box
    add_box(slide, "\u3053\u306e\u300c\u5dee\u5206\u300d\u304c\u6700\u9ad8\u306e\u6559\u5e2b\u30c7\u30fc\u30bf:\n\u30fbAI\u304c\u4f55\u3092\u9593\u9055\u3048\u305f\u304b\n\u30fb\u6b63\u89e3\u306f\u4f55\u3060\u3063\u305f\u304b\n\u30fb\u306a\u305c\u9593\u9055\u3048\u305f\u304b\uff08\u6697\u9ed9\u306e\u5224\u65ad\u57fa\u6e96\uff09",
            Inches(7.80), Inches(1.70), Inches(5.0), Inches(0.90),
            bg_color=BG_LIGHT, text_color=CHARCOAL, font_size=8, align=PP_ALIGN.LEFT)

    # Gap quality comparison
    gap_table = [
        ["\u6559\u5e2b\u30c7\u30fc\u30bf\u306e\u7a2e\u985e", "\u54c1\u8cea", "\u8ffd\u52a0\u30b3\u30b9\u30c8", "\u91cf", "\u7279\u5fb4"],
        ["\u6559\u79d1\u66f8\u7684\u306a\u6b63\u89e3\u30c7\u30fc\u30bf", "\u25cb", "\u9ad8\uff08\u4f5c\u6210\u5de5\u6570\uff09", "\u5c11", "\u4e8b\u524d\u306b\u4f5c\u308a\u8fbc\u307f\u304c\u5fc5\u8981"],
        ["\u904e\u53bb\u306e\u4e8b\u4f8b\u30c7\u30fc\u30bf", "\u25b3", "\u4e2d\uff08\u6574\u7406\u5de5\u6570\uff09", "\u4e2d", "\u69cb\u9020\u5316\u306e\u624b\u9593\u304c\u304b\u304b\u308b"],
        ["AI-Human Gap", "\u2605\u2605\u2605", "\u30bc\u30ed\uff08\u65e5\u5e38\u696d\u52d9\u306e\u526f\u7523\u7269\uff09", "\u5927\uff08\u4f7f\u3046\u307b\u3069\u6e9c\u307e\u308b\uff09", "\u6700\u3082\u81ea\u7136\u3067\u6301\u7d9a\u7684"],
    ]
    add_table(slide, gap_table, left=Inches(0.42), top=Inches(2.80), width=Inches(12.46),
              col_widths=[2.5, 0.8, 2.5, 2.5, 2.5],
              header_font_pt=7, cell_font_pt=6.5, row_height=0.26)

    # Positive spiral
    add_box(slide, "\u6b63\u306e\u30b9\u30d1\u30a4\u30e9\u30eb",
            Inches(0.42), Inches(4.00), Inches(12.46), Inches(0.20),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=7, bold=True)

    spiral_steps = [
        ("AI\u3092\u4f7f\u3046", GREEN_MINT),
        ("Gap\u767a\u751f", RGBColor(0xFD, 0xE8, 0xC8)),
        ("\u4fee\u6b63\u3059\u308b", GREEN_PALE),
        ("Gap\u304c\u6559\u5e2b\u30c7\u30fc\u30bf\u306b", ACCENT_ORANGE),
        ("AI\u304c\u8ce2\u304f\u306a\u308b", GREEN_MINT),
        ("\u3082\u3063\u3068\u4f7f\u3044\u305f\u304f\u306a\u308b", GREEN_PALE),
        ("\u3055\u3089\u306bGap\u84c4\u7a4d \u2192 \u3055\u3089\u306b\u8ce2\u304f\u306a\u308b...", RGBColor(0xFD, 0xE8, 0xC8)),
    ]
    x_sp = 0.42
    sp_w = 12.46 / len(spiral_steps) - 0.05
    for i, (label, color) in enumerate(spiral_steps):
        text_c = WHITE if color == ACCENT_ORANGE else CHARCOAL
        add_box(slide, label,
                Inches(x_sp), Inches(4.25), Inches(sp_w), Inches(0.35),
                bg_color=color, text_color=text_c, font_size=6.5, bold=True)
        if i < len(spiral_steps) - 1:
            add_arrow_right(slide, Inches(x_sp + sp_w), Inches(4.28), Inches(0.15), Inches(0.25))
        x_sp += sp_w + 0.05

    # ── Slide 6: \u62e1\u5f35\u53ef\u80fd\u6027 ──
    slide = slides[6]
    set_title(slide, "\u62e1\u5f35\u53ef\u80fd\u6027: \u81ea\u52d5\u9032\u5316 + \u4eba\u7684\u4ecb\u5165\u306b\u3088\u308b\u9032\u5316",
              "\u8ffd\u52a0\u958b\u767a\u306a\u3057\u3067\u65e5\u5e38\u7684\u306b\u9032\u5316\u3059\u308b\u4ed5\u7d44\u307f\u3068\u3001\u5b9a\u671f\u7684\u306a\u30ec\u30d3\u30e5\u30fc\u3067\u6607\u683c\u3059\u308b\u4ed5\u7d44\u307f")

    # Auto evolution
    add_box(slide, "\u81ea\u52d5\u9032\u5316\uff08\u8ffd\u52a0\u958b\u767a\u306a\u3057\uff09",
            Inches(0.42), Inches(1.65), Inches(6.2), Inches(0.22),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    auto_evo = [
        ["\u4ed5\u7d44\u307f", "\u4f55\u304c\u8d77\u304d\u308b\u304b", "\u30bf\u30a4\u30df\u30f3\u30b0", "\u95a2\u9023\u30b9\u30c6\u30fc\u30b8"],
        ["Decision Trace\u84c4\u7a4d", "\u627f\u8a8d\u30d1\u30bf\u30fc\u30f3 \u2192 \u30eb\u30fc\u30eb\u81ea\u52d5\u62bd\u51fa", "\u65e5\u5e38\u7684\uff08RM\u64cd\u4f5c\u306e\u305f\u3073\uff09", "\u2461\u84c4\u7a4d\u2192\u2462\u6607\u683c"],
        ["\u30c1\u30e3\u30c3\u30c8FB\u84c4\u7a4d", "\u7684\u5916\u308c\u691c\u51fa \u2192 \u30d7\u30ed\u30f3\u30d7\u30c8\u6539\u5584\u5019\u88dc", "\u65e5\u5e38\u7684\uff08\u30c1\u30e3\u30c3\u30c8\u5229\u7528\u306e\u305f\u3073\uff09", "\u2464\u9032\u5316"],
        ["\u6210\u7d04/\u5931\u6ce8\u84c4\u7a4d", "\u4e8b\u4f8b\u30d1\u30bf\u30fc\u30f3 \u2192 K5\u81ea\u52d5\u66f4\u65b0", "\u5546\u8ac7\u7d50\u679c\u78ba\u5b9a\u6642", "\u2461\u84c4\u7a4d\u2192\u2462\u6607\u683c"],
    ]
    add_table(slide, auto_evo, left=Inches(0.42), top=Inches(1.90), width=Inches(6.2),
              col_widths=[1.5, 2.0, 1.5, 1.0],
              header_font_pt=6.5, cell_font_pt=6, row_height=0.30)

    # Human evolution
    add_box(slide, "\u4eba\u7684\u4ecb\u5165\u306b\u3088\u308b\u9032\u5316\uff08\u5b9a\u671f\u30ec\u30d3\u30e5\u30fc\uff09",
            Inches(6.90), Inches(1.65), Inches(5.9), Inches(0.22),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    human_evo = [
        ["\u4ed5\u7d44\u307f", "\u4f55\u304c\u8d77\u304d\u308b\u304b", "\u30bf\u30a4\u30df\u30f3\u30b0", "\u95a2\u9023\u30b9\u30c6\u30fc\u30b8"],
        ["\u30a8\u30ad\u30b9\u30d1\u30fc\u30c8\u30d2\u30a2\u30ea\u30f3\u30b0", "\u30d9\u30c6\u30e9\u30f3\u306e\u7d4c\u9a13\u5247 \u2192\nK1-K8\u521d\u671f\u69cb\u7bc9", "PoC\uff5eMVP\u521d\u671f", "\u2460\u30ad\u30e3\u30d7\u30c1\u30e3"],
        ["PJ\u30c1\u30fc\u30e0\u30ec\u30d3\u30e5\u30fc", "\u84c4\u7a4d\u30d1\u30bf\u30fc\u30f3\u306e\u54c1\u8cea\u78ba\u8a8d\n\u2192 K\u6607\u683c", "\u56db\u534a\u671f\u3054\u3068", "\u2462\u6607\u683c"],
        ["\u7d44\u7e54\u6a2a\u65ad\u5206\u6790", "\u55b6\u696d\u5e97\u9593\u306e\u30d1\u30bf\u30fc\u30f3\u6bd4\u8f03\n\u2192 \u5168\u793e\u30eb\u30fc\u30eb", "\u534a\u671f\u3054\u3068", "\u2463\u5c55\u958b"],
    ]
    add_table(slide, human_evo, left=Inches(6.90), top=Inches(1.90), width=Inches(5.9),
              col_widths=[1.3, 2.0, 1.0, 1.0],
              header_font_pt=6.5, cell_font_pt=6, row_height=0.35)

    # Integration note
    add_note(slide, "\u81ea\u52d5\u9032\u5316\u306f\u65e5\u5e38\u696d\u52d9\u306e\u526f\u7523\u7269\u3068\u3057\u3066\u9032\u884c\u3002\u4eba\u7684\u4ecb\u5165\u306f\u54c1\u8cea\u62c5\u4fdd\u3068\u7d44\u7e54\u77e5\u306e\u69cb\u9020\u5316\u306b\u5fc5\u8981\u3002\u4e21\u8005\u304c\u88dc\u5b8c\u7684\u306b\u50cd\u304f\u3053\u3068\u3067\u3001\u6301\u7d9a\u7684\u306a\u9032\u5316\u304c\u5b9f\u73fe\u3059\u308b\u3002",
             Inches(0.42), Inches(3.60), Inches(12.46))

    # ── Slide 7: \u30b5\u30de\u30ea\u30fc ──
    slide = slides[7]
    set_title(slide, "\u30b5\u30de\u30ea\u30fc: \u6697\u9ed9\u77e5\u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb\u306e\u8a2d\u8a08\u539f\u5247",
              "\u5168\u4f53\u3092\u8cab\u304f\u8a2d\u8a08\u539f\u5247\u3068\u5404\u30b9\u30c6\u30fc\u30b8\u306e\u5bfe\u5fdc\u95a2\u4fc2")

    summary_table = [
        ["\u8a2d\u8a08\u539f\u5247", "\u8aac\u660e", "\u5bfe\u5fdc\u30b9\u30c6\u30fc\u30b8", "\u5b9f\u88c5\u65b9\u6cd5"],
        ["\u8ffd\u52a0\u64cd\u4f5c\u30bc\u30ed", "RM\u306e\u65e5\u5e38\u696d\u52d9\u306e\u4e2d\u3067\u6697\u9ed9\u77e5\u3092\n\u6355\u7372\u3002\u5c02\u7528\u753b\u9762\u3092\u4f5c\u3089\u306a\u3044", "\u2460\u30ad\u30e3\u30d7\u30c1\u30e3", "\u9762\u8ac7\u8a18\u9332\u306b\u69cb\u9020\u5316\u9805\u76ee\u57cb\u8fbc\u307f\n\u627f\u8a8d\u30fbFB\u30dcUI\u306b\u7d44\u8fbc"],
        ["2\u30d1\u30b9\u6607\u683c", "\u5373\u6642Data\u5316\u3068\u30d1\u30bf\u30fc\u30f3\u84c4\u7a4d\u306b\u3088\u308b\nKnowledge\u5316\u306e2\u30eb\u30fc\u30c8", "\u2461\u84c4\u7a4d\u2192\u2462\u6607\u683c", "\u30d1\u30b9A: \u69cb\u9020\u5316\u9805\u76ee\u2192\u5373Data\n\u30d1\u30b9B: 3\u56de\u84c4\u7a4d\u2192Knowledge"],
        ["\u591a\u5c64\u30e1\u30e2\u30ea", "\u500b\u4eba\u2192\u55b6\u696d\u5e97\u2192\u696d\u754c\u2192\u5168\u793e\u306e\n\u968e\u5c64\u3067\u30ca\u30ec\u30c3\u30b8\u3092\u5171\u6709\u30fb\u6607\u683c", "\u2463\u5c55\u958b", "4\u5c64 + \u9867\u5ba2\u30e1\u30e2\u30ea\uff08\u6a2a\u65ad\uff09\n\u8907\u6570RM\u540c\u50be\u5411\u3067\u6607\u683c"],
        ["AI-Human Gap\u6d3b\u7528", "RM\u306e\u4fee\u6b63 = \u6700\u9ad8\u306e\u6559\u5e2b\u30c7\u30fc\u30bf\u3002\n\u4f7f\u3046\u307b\u3069AI\u304c\u8ce2\u304f\u306a\u308b\u6b63\u306e\u30b9\u30d1\u30a4\u30e9\u30eb", "\u2464\u9032\u5316", "Gap\u3092Decision Trace\u306b\u8a18\u9332\n\u30d1\u30bf\u30fc\u30f3\u62bd\u51fa\u2192K\u30eb\u30fc\u30eb\u5316"],
        ["\u81ea\u52d5+\u4eba\u7684\u9032\u5316", "\u65e5\u5e38\u696d\u52d9\u306e\u526f\u7523\u7269\u3067\u81ea\u52d5\u9032\u5316\u3057\u3001\n\u5b9a\u671f\u30ec\u30d3\u30e5\u30fc\u3067\u54c1\u8cea\u62c5\u4fdd", "\u62e1\u5f35", "\u81ea\u52d5: \u65e5\u5e38\u7684\u84c4\u7a4d+\u62bd\u51fa\n\u4eba\u7684: \u56db\u534a\u671f/\u534a\u671f\u30ec\u30d3\u30e5\u30fc"],
    ]
    add_table(slide, summary_table, top=Inches(1.70),
              col_widths=[2.0, 3.0, 1.3, 3.5],
              header_font_pt=7, cell_font_pt=6.5, row_height=0.50)

    # Key takeaway
    add_box(slide, "\u8981\u70b9: \u6697\u9ed9\u77e5\u306f\u300c\u8a18\u9332\u3057\u3066\u304f\u3060\u3055\u3044\u300d\u3068\u304a\u9858\u3044\u3059\u308b\u3082\u306e\u3067\u306f\u306a\u304f\u3001\u65e5\u5e38\u696d\u52d9\u306e\u4e2d\u3067\u81ea\u7136\u306b\u6355\u7372\u3055\u308c\u3001\n\u30d1\u30bf\u30fc\u30f3\u691c\u51fa\u3067\u81ea\u52d5\u7684\u306bData/Knowledge\u306b\u6607\u683c\u3057\u3001\u591a\u5c64\u30e1\u30e2\u30ea\u3067\u7d44\u7e54\u5168\u4f53\u306b\u5c55\u958b\u3055\u308c\u3001AI\u304c\u8ce2\u304f\u306a\u308a\u7d9a\u3051\u308b\u3002",
            Inches(0.42), Inches(4.85), Inches(12.46), Inches(0.55),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=8, bold=True, align=PP_ALIGN.LEFT)

    # Save
    output_path = "output/15_tacit_knowledge_lifecycle_detailed.pptx"
    new_prs.save(output_path)
    print(f"Saved detailed version: {output_path} ({len(slides)} slides)")
    return output_path


if __name__ == "__main__":
    build()
