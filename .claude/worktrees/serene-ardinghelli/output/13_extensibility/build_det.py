#!/usr/bin/env python3
"""Build detailed PPTX for 13_extensibility_design — handout version.
More dense, with full YAML examples, all tables, and detailed explanations.
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

TEMPLATE = "project-rikyu-sales-proposals-poc/untracked/original_document/report/提案仮説構築AI_定例会1_20260219_v1.pptx"
WORKDIR = "output/13_extensibility"

# === ACES Color Palette ===
GREEN_DARK = RGBColor(0x00, 0x68, 0x43)
GREEN_MEDIUM = RGBColor(0x19, 0x7A, 0x56)
GREEN_BRIGHT = RGBColor(0x00, 0x9A, 0x62)
GREEN_MINT = RGBColor(0xC2, 0xE7, 0xD9)
GREEN_PALE = RGBColor(0xDF, 0xE6, 0xE0)
CHARCOAL = RGBColor(0x39, 0x39, 0x39)
TEXT_MAIN = RGBColor(0x12, 0x12, 0x12)
TABLE_TEXT = RGBColor(0x45, 0x45, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x52, 0x98, 0xBA)
ACCENT_ORANGE = RGBColor(0xFB, 0xAE, 0x40)
GRAY_LIGHT = RGBColor(0x99, 0x99, 0x99)

TITLE_FONT = "Yu Gothic"
SUB_FONT = "Yu Gothic Medium"
BODY_FONT = "Century Gothic"


def find_shapes(slide):
    shapes_with_y = []
    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            y = s.top / 914400 if s.top else 0
            shapes_with_y.append((y, s))
    shapes_with_y.sort(key=lambda x: x[0])
    title = subtitle = body = None
    for y, s in shapes_with_y:
        if y < 0.7 and title is None: title = s
        elif 0.7 <= y < 1.5 and subtitle is None: subtitle = s
        elif y >= 1.5 and body is None: body = s
    return title, subtitle, body


def prep(slide):
    title, subtitle, body = find_shapes(slide)
    if subtitle:
        subtitle.height = Inches(0.32)
        if subtitle.width < Inches(1):
            subtitle.width = Inches(12.46)
            subtitle.left = Inches(0.42)
    if body:
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = " "
        p.font.size = Pt(1)
        p.font.color.rgb = WHITE
    return title, subtitle, body


def add_table(slide, data, left=Inches(0.42), top=Inches(1.90), width=Inches(12.46),
              col_widths=None, header_font_pt=8, cell_font_pt=7, row_height=0.28):
    rows, cols = len(data), len(data[0])
    height = Inches(row_height * rows)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.name = BODY_FONT
                if row_idx == 0:
                    para.font.size = Pt(header_font_pt)
                    para.font.color.rgb = WHITE
                    para.font.bold = True
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.font.size = Pt(cell_font_pt)
                    para.font.color.rgb = TABLE_TEXT
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_DARK
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_PALE
    return table


def add_box(slide, text, left, top, width, height, bg_color=GREEN_MINT, text_color=GREEN_DARK,
            font_size=9, bold=False, align=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_arrow_down(slide, left, top, width=Inches(0.5), height=Inches(0.20)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "▼"
    p.font.size = Pt(12)
    p.font.color.rgb = GREEN_DARK
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.25)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_note(slide, text, left, top, width, h=Inches(0.30)):
    txBox = slide.shapes.add_textbox(left, top, width, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(6)
    p.font.italic = True
    p.font.color.rgb = GRAY_LIGHT
    p.font.name = BODY_FONT


def build_detailed():
    """Build detailed version: 8 slides — Title + 7 dense content slides."""
    prs = Presentation(TEMPLATE)

    # Layouts
    layout_cover_data = prs.slide_layouts[1]
    layout_content = prs.slide_layouts[10]

    # Build new presentation
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    # Slide plan: Title + 7 content slides
    layouts = [layout_cover_data] + [layout_content] * 7
    for layout in layouts:
        new_prs.slides.add_slide(layout)

    slides = new_prs.slides

    # ── Slide 0: Title ──
    slide = slides[0]
    shapes_with_y = []
    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            y = s.top / 914400 if s.top else 0
            shapes_with_y.append((y, s))
    shapes_with_y.sort(key=lambda x: x[0])
    shape_list = [s for _, s in shapes_with_y]
    if len(shape_list) >= 1:
        tf = shape_list[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "株式会社りそな銀行 御中"
        p.font.size = Pt(32)
        p2 = tf.add_paragraph()
        p2.text = "汎用性・拡張性の設計方針【詳細版】"
        p2.font.size = Pt(28)
    if len(shape_list) >= 2:
        tf = shape_list[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "2026/3/10 | rikyu ソリューションセールス伴奏AI PoC"
        p.font.size = Pt(14)

    # ── Slide 1: LLM first + 5 extension axes ──
    slide = slides[1]
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = "LLM first設計方針 + 5つの拡張軸"
        p.font.name = TITLE_FONT
        p.font.size = Pt(22)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = "人間がルールを書く量を最小化し、コードを変えずにシステムが賢くなる設計"
        p.font.name = SUB_FONT
        p.font.size = Pt(11)

    # Comparison header
    add_box(slide, "従来型", Inches(0.42), Inches(1.80), Inches(3.8), Inches(0.25),
            bg_color=CHARCOAL, text_color=WHITE, font_size=8, bold=True)
    add_box(slide, "LLM first", Inches(4.42), Inches(1.80), Inches(3.8), Inches(0.25),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    add_box(slide, "ルール追加 → コード変更\n新機能追加 → 開発+テスト+デプロイ\n業界追加 → 業界固有のコード",
            Inches(0.42), Inches(2.08), Inches(3.8), Inches(0.80),
            bg_color=RGBColor(0xF5, 0xF5, 0xF5), text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)
    add_box(slide, "ルール追加 → YAML/Markdown追加\n新機能追加 → パイプライン定義追加\n業界追加 → テンプレートYAML追加",
            Inches(4.42), Inches(2.08), Inches(3.8), Inches(0.80),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=7, align=PP_ALIGN.LEFT)

    # 5 axes table
    add_box(slide, "5つの拡張軸", Inches(8.60), Inches(1.80), Inches(4.2), Inches(0.25),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    ext_table = [
        ["軸", "変更方法", "コード変更"],
        ["ナレッジ", "YAML/MD追加・編集", "不要"],
        ["プロンプト", "MDプロンプト編集", "不要"],
        ["ツール", "ツール定義追加", "最小"],
        ["パイプライン", "パイプライン定義追加", "必要"],
        ["エージェント", "エージェント定義+登録", "必要"],
    ]
    add_table(slide, ext_table, left=Inches(8.60), top=Inches(2.08), width=Inches(4.2),
              col_widths=[1.2, 2.0, 1.0], header_font_pt=7, cell_font_pt=6.5, row_height=0.22)

    add_note(slide, "上から順に容易。日常的な改善はナレッジ/プロンプト層で完結。",
             Inches(0.42), Inches(2.95), Inches(7.5))

    # ── Slide 2: ナレッジ拡張 — YAML詳細 ──
    slide = slides[2]
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = "ナレッジ拡張: 業界テンプレート(K1) + 更新ルール(K2)"
        p.font.name = TITLE_FONT
        p.font.size = Pt(20)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = "YAMLファイルを書くだけで、コード変更なしにAIの知識を拡張する"
        p.font.name = SUB_FONT
        p.font.size = Pt(10)

    # K1 YAML
    add_box(slide, "K1: 業界テンプレート", Inches(0.42), Inches(1.70), Inches(6.2), Inches(0.22),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    yaml1 = (
        "# knowledge_master/industry_templates/construction.yaml\n"
        "industry_code: \"D-06\"\n"
        "industry_name: \"建設業\"\n"
        "common_agendas:\n"
        "  - name: \"人材確保・定着\"\n"
        "    description: \"期間工問題、技能者の高齢化、若手採用\"\n"
        "    priority: high\n"
        "    typical_solutions: [\"人材紹介\", \"福利厚生支援\"]\n"
        "  - name: \"DX・BIM推進\"\n"
        "    description: \"設計・施工のデジタル化、BIM導入\"\n"
        "    priority: medium\n"
        "    typical_solutions: [\"IT投資融資\", \"DXコンサル\"]\n"
        "  - name: \"安全管理・コンプライアンス\"\n"
        "    description: \"労災防止、法令順守体制の強化\"\n"
        "    priority: medium\n"
        "segment_patterns:\n"
        "  - \"土木 / 建築\"  - \"国内 / 海外\"  - \"元請 / 下請\""
    )
    add_box(slide, yaml1,
            Inches(0.42), Inches(1.95), Inches(6.2), Inches(3.30),
            bg_color=RGBColor(0x2D, 0x2D, 0x2D), text_color=RGBColor(0xE0, 0xE0, 0xE0),
            font_size=6.5, align=PP_ALIGN.LEFT)

    # K2 YAML
    add_box(slide, "K2: 更新ルール", Inches(6.90), Inches(1.70), Inches(5.9), Inches(0.22),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    yaml2 = (
        "# knowledge_master/update_rules/trigger_patterns.yaml\n"
        "triggers:\n"
        "  - pattern: \"決裁権者が具体的な時期に言及\"\n"
        "    action: \"時期を更新 + 優先度UP\"\n"
        "    confidence: high\n"
        "    source: \"expert_interview_20260312\"\n"
        "  - pattern: \"担当者レベルの伝聞\"\n"
        "    action: \"備考に追記のみ\"\n"
        "    confidence: low"
    )
    add_box(slide, yaml2,
            Inches(6.90), Inches(1.95), Inches(5.9), Inches(1.80),
            bg_color=RGBColor(0x2D, 0x2D, 0x2D), text_color=RGBColor(0xE0, 0xE0, 0xE0),
            font_size=6.5, align=PP_ALIGN.LEFT)

    # Lifecycle
    add_box(slide, "ナレッジのライフサイクル", Inches(6.90), Inches(3.85), Inches(5.9), Inches(0.22),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    steps = ["エキスパートヒアリング", "ナレッジの言語化",
             "YAML/Markdownに構造化", "Knowledge Masterに投入",
             "AIの出力品質が向上", "RMのFBでさらに改善"]
    y_step = 4.12
    for i, label in enumerate(steps):
        add_box(slide, label,
                Inches(6.90), Inches(y_step), Inches(5.9), Inches(0.22),
                bg_color=GREEN_PALE if i < 3 else GREEN_MINT,
                text_color=CHARCOAL, font_size=7)
        y_step += 0.22
        if i < len(steps) - 1:
            add_arrow_down(slide, Inches(9.6), Inches(y_step - 0.02), Inches(0.4), Inches(0.15))
            y_step += 0.12

    add_note(slide, "コード変更不要・デプロイ不要。AIの出力確認のみでナレッジ拡張が完了する。",
             Inches(0.42), Inches(5.40), Inches(12.46))

    # ── Slide 3: ツール拡張 詳細 ──
    slide = slides[3]
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = "ツール拡張: 名前ベースディスパッチ — 詳細設計"
        p.font.name = TITLE_FONT
        p.font.size = Pt(20)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = "agentic harnessパターン: ツール追加がループ構造に影響しない設計"
        p.font.name = SUB_FONT
        p.font.size = Pt(10)

    # Dispatch code box
    code_text = (
        "# ツールのディスパッチマップ（概念）\n"
        "tool_registry = {\n"
        "    \"web_search\": web_search_tool,\n"
        "    \"edinet_api\": edinet_api_tool,\n"
        "    \"tsr_api\": tsr_api_tool,\n"
        "    # 新しいツールはここに1行追加するだけ\n"
        "    \"business_card_api\": business_card_tool,\n"
        "    \"financial_model\": financial_model_tool,\n"
        "}"
    )
    add_box(slide, code_text,
            Inches(0.42), Inches(1.70), Inches(5.5), Inches(2.2),
            bg_color=RGBColor(0x2D, 0x2D, 0x2D), text_color=RGBColor(0xE0, 0xE0, 0xE0),
            font_size=7, align=PP_ALIGN.LEFT)

    # Process flow
    add_box(slide, "LLMが「このツールを使いたい」と判断",
            Inches(0.42), Inches(4.05), Inches(5.5), Inches(0.25),
            bg_color=ACCENT_BLUE, text_color=WHITE, font_size=7, bold=True)
    add_arrow_down(slide, Inches(2.9), Inches(4.30))
    add_box(slide, "ディスパッチマップで実行先を解決",
            Inches(0.42), Inches(4.48), Inches(5.5), Inches(0.25),
            bg_color=GREEN_MINT, text_color=CHARCOAL, font_size=7, bold=True)
    add_arrow_down(slide, Inches(2.9), Inches(4.73))
    add_box(slide, "結果をコンテキストに追加",
            Inches(0.42), Inches(4.90), Inches(5.5), Inches(0.25),
            bg_color=GREEN_PALE, text_color=CHARCOAL, font_size=7, bold=True)

    # Future tools table
    future_tools = [
        ["ツール", "接続先", "追加方法", "影響範囲"],
        ["名刺アプリ連携", "名刺管理サービスAPI", "ツール定義追加", "KPマップのデータソース拡充"],
        ["SFA/CRM連携", "Salesforce等", "ツール定義追加", "面談・商談データの自動取得"],
        ["財務分析ツール", "既存会計モデル", "ツール定義追加", "M5提案の精度向上"],
        ["パワポ出力", "python-pptx等", "ツール定義追加", "出力フォーマットの追加"],
    ]
    add_table(slide, future_tools, left=Inches(6.30), top=Inches(1.70), width=Inches(6.5),
              col_widths=[1.5, 1.8, 1.3, 1.9], header_font_pt=7, cell_font_pt=6.5, row_height=0.26)

    add_box(slide, "新しいツールの追加はループ本体の変更を一切必要としない",
            Inches(6.30), Inches(3.20), Inches(6.5), Inches(0.30),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=7, bold=True)

    # ── Slide 4: パイプライン拡張 詳細 ──
    slide = slides[4]
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = "パイプライン拡張: M1-M10スキル追加型アーキテクチャ"
        p.font.name = TITLE_FONT
        p.font.size = Pt(20)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = "各パイプラインは独立したスキル。既存パイプラインに影響なし"
        p.font.name = SUB_FONT
        p.font.size = Pt(10)

    # Pipeline tree
    add_box(slide, "Pipeline Orchestrator",
            Inches(0.42), Inches(1.70), Inches(4.0), Inches(0.22),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    pipelines = [
        ("M1", "アジェンダ初期生成", True),
        ("M2", "面談→アジェンダ更新", True),
        ("M3", "ニュース検出", True),
        ("M4", "AIチャット", True),
        ("M5", "提案ストーリー", True),
        ("M6", "深掘りペーパー", True),
        ("M7", "担当者用メモ", True),
        ("M8", "アカウントプラン生成", False),
        ("M9", "競合分析レポート", False),
        ("M10", "研修シミュレーション", False),
    ]
    y_pipe = 1.95
    for code, name, existing in pipelines:
        bg = GREEN_MINT if existing else RGBColor(0xEE, 0xF1, 0xF5)
        suffix = "既存" if existing else "将来追加"
        lbl = f"├── {code} {name}  [{suffix}]" if code != "M10" else f"└── {code} {name}  [{suffix}]"
        add_box(slide, lbl,
                Inches(0.42), Inches(y_pipe), Inches(4.0), Inches(0.20),
                bg_color=bg, text_color=CHARCOAL, font_size=6.5, align=PP_ALIGN.LEFT)
        y_pipe += 0.21

    # Requirements table
    req_table = [
        ["要素", "内容", "形式"],
        ["パイプライン定義", "入力→処理ステップ→出力の定義", "Python"],
        ["プロンプトテンプレート", "LLMへの指示", "Markdown"],
        ["Knowledge参照定義", "どのK1-K8を使うか", "YAML"],
        ["API定義", "エンドポイント、パラメータ", "OpenAPI"],
    ]
    add_table(slide, req_table, left=Inches(4.80), top=Inches(1.70), width=Inches(7.9),
              col_widths=[2.0, 3.5, 1.2], header_font_pt=7, cell_font_pt=6.5, row_height=0.24)

    add_box(slide, "ポイント: 新しいスキルを追加しても既存パイプラインに影響しない。M8-M10は将来の拡張候補。",
            Inches(4.80), Inches(3.05), Inches(7.9), Inches(0.30),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=7, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 5: 変更影響の局所化 詳細テーブル ──
    slide = slides[5]
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = "変更影響の局所化 — 全6カテゴリの影響範囲マトリクス"
        p.font.name = TITLE_FONT
        p.font.size = Pt(20)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = "どのレイヤーを変更しても、他のレイヤーに影響が波及しない疎結合設計"
        p.font.name = SUB_FONT
        p.font.size = Pt(10)

    impact_table = [
        ["変更の種類", "影響範囲", "コード変更", "デプロイ", "テスト"],
        ["ナレッジ追加・修正\n(業界テンプレート、ルール等)", "YAMLファイルのみ\n他のコードに影響なし", "不要", "不要", "AI出力確認"],
        ["プロンプト改善\n(出力品質向上)", "Markdownファイルのみ\n他のコードに影響なし", "不要", "不要", "AI出力確認"],
        ["ツール追加\n(名刺アプリ、SFA等)", "ツール定義 + ディスパッチ登録\nループ構造に影響なし", "最小", "必要", "ツール単体テスト"],
        ["パイプライン追加\n(M8, M9...)", "新パイプライン定義\n既存パイプラインに影響なし", "必要", "必要", "統合テスト"],
        ["エージェント追加\n(専門領域の分割)", "エージェント定義 +\nオーケストレーター登録", "必要", "必要", "統合テスト"],
        ["LLMモデル変更\n(GPT-5.2→次世代)", "設定ファイルのみ\nアルゴリズム構造に影響なし", "不要", "設定変更", "回帰テスト"],
    ]
    add_table(slide, impact_table, top=Inches(1.70), col_widths=[2.5, 3.5, 1.0, 1.0, 1.5],
              header_font_pt=7, cell_font_pt=6.5, row_height=0.45)

    # Visual boxes
    y_vis = 5.0
    add_box(slide, "コード変更不要（ナレッジ・プロンプト・LLM変更）",
            Inches(0.42), Inches(y_vis), Inches(4.0), Inches(0.30),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=7, bold=True)
    add_box(slide, "最小コード変更（ツール追加）",
            Inches(4.60), Inches(y_vis), Inches(3.5), Inches(0.30),
            bg_color=ACCENT_ORANGE, text_color=WHITE, font_size=7, bold=True)
    add_box(slide, "コード変更必要（パイプライン・エージェント）",
            Inches(8.30), Inches(y_vis), Inches(4.5), Inches(0.30),
            bg_color=ACCENT_BLUE, text_color=WHITE, font_size=7, bold=True)

    # ── Slide 6: 確定的×判断的 詳細 ──
    slide = slides[6]
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = "確定的(ルール) × 判断的(LLM)の使い分け — 詳細設計"
        p.font.name = TITLE_FONT
        p.font.size = Pt(20)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = "「絶対に守らせたい」はルールエンジンに、「できるだけ従ってほしい」はプロンプトに"
        p.font.name = SUB_FONT
        p.font.size = Pt(10)

    # Detailed table
    rule_table = [
        ["レイヤー", "性質", "rikyuでの対応", "例", "変更方法"],
        ["確定的", "必ずこう動く\n例外なし", "ルールエンジン(YAML)", "信頼度マトリクス\n更新ルールの確定パターン", "YAMLファイル編集\nコード変更不要"],
        ["判断的", "できるだけ\nこう動く", "LLM + プロンプト", "文脈解釈、カスタマイズ\n新規パターン検出", "プロンプト編集\nコード変更不要"],
        ["キャプチャ", "RM操作を\n必ず記録", "Decision Trace", "承認/修正/棄却\n+理由の記録", "自動記録\n設定変更のみ"],
    ]
    add_table(slide, rule_table, top=Inches(1.70), col_widths=[1.3, 1.8, 2.2, 2.8, 2.2],
              header_font_pt=7, cell_font_pt=6.5, row_height=0.50)

    # Three comparison boxes
    box_w = Inches(4.0)
    y_box = 3.80
    add_box(slide, "確定的（ルールエンジン）\n\n「絶対に守らせたい」ルール\n・信頼度マトリクス: 発言者×具体性→信頼度\n・更新ルールの確定パターン: 決裁権者言及→優先度UP\n・入力バリデーション: 必須項目チェック\n\n特徴: 例外なし。予測可能。テスト容易。",
            Inches(0.42), Inches(y_box), box_w, Inches(2.5),
            bg_color=RGBColor(0xD6, 0xEA, 0xF8), text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)

    add_box(slide, "判断的（LLM + プロンプト）\n\n「できるだけ従ってほしい」指針\n・文脈に応じた発言解釈\n・新規パターンの自動検出\n・カスタマイズされた提案生成\n\n特徴: 柔軟。新しい状況に対応可能。",
            Inches(4.62), Inches(y_box), box_w, Inches(2.5),
            bg_color=GREEN_MINT, text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)

    add_box(slide, "キャプチャ（Decision Trace）\n\n「RM操作を必ず記録」する仕組み\n・承認/修正/棄却 + その理由\n・フィードバックループへの入力\n・ルール改善のためのデータ蓄積\n\n特徴: 自動記録。学習基盤。監査対応。",
            Inches(8.82), Inches(y_box), Inches(3.9), Inches(2.5),
            bg_color=GREEN_PALE, text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)

    # ── Slide 7: 設計原則サマリー ──
    slide = slides[7]
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = "設計原則サマリー — 汎用性・拡張性を支える4つの柱"
        p.font.name = TITLE_FONT
        p.font.size = Pt(20)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = "rikyuの拡張性を保証する設計原則と、それぞれの実装パターン"
        p.font.name = SUB_FONT
        p.font.size = Pt(10)

    summary_table = [
        ["設計原則", "説明", "rikyuでの実装", "効果"],
        ["LLM first", "人間がルールを書く量を\n最小化する", "ナレッジ=YAML/MD\nプロンプト=MD", "コード変更なしで\n品質向上"],
        ["名前ベースディスパッチ", "ツール追加がループ\n構造に影響しない", "tool_registry + \nagentic harness", "ツール追加が\n1行で完了"],
        ["スキル追加型", "パイプラインが独立した\nスキルとして並列実行", "M1-M7 + 将来M8-M10\nPipeline Orchestrator", "既存パイプラインに\n影響なし"],
        ["確定的×判断的の分離", "ルールとLLM判断を\n明確に分離", "ルールエンジン(YAML) +\nLLM + Decision Trace", "信頼性と柔軟性の\n両立"],
    ]
    add_table(slide, summary_table, top=Inches(1.70), col_widths=[2.3, 2.8, 3.0, 2.2],
              header_font_pt=8, cell_font_pt=7, row_height=0.55)

    # Key takeaway
    add_box(slide, "要点: rikyuは「コードを変えずに賢くなる」システム。\n"
                   "日常的な改善はナレッジ/プロンプト層で完結し、コード変更・デプロイなしで品質が向上する。\n"
                   "新しいツール・パイプラインの追加も、既存のシステムに影響を与えない疎結合設計。",
            Inches(0.42), Inches(4.2), Inches(12.46), Inches(0.70),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=9, bold=True, align=PP_ALIGN.LEFT)

    # Extensibility roadmap
    roadmap = [
        ["フェーズ", "期間", "拡張内容", "コード変更"],
        ["MVP", "〜M3", "K1-K8ナレッジ構築、M1-M7パイプライン", "初期開発"],
        ["成長期", "M4-M6", "業界テンプレート追加、プロンプト改善", "不要（YAML/MD）"],
        ["拡張期", "M7-M9", "ツール追加（名刺、SFA等）", "最小（定義追加）"],
        ["発展期", "M10-M12", "M8-M10パイプライン追加", "必要（新規開発）"],
    ]
    add_table(slide, roadmap, left=Inches(0.42), top=Inches(5.10), width=Inches(12.46),
              col_widths=[1.5, 1.2, 5.0, 2.5], header_font_pt=7, cell_font_pt=6.5, row_height=0.24)

    # Save
    output_path = "output/13_extensibility_design_detailed.pptx"
    new_prs.save(output_path)
    print(f"Saved detailed version: {output_path}")
    return output_path


if __name__ == "__main__":
    build_detailed()
