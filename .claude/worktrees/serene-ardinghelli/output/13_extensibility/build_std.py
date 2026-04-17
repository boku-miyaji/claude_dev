#!/usr/bin/env python3
"""Build standard PPTX for 13_extensibility_design from template.

Steps:
1. Copy template
2. Build slides (Title, INDEX, Section dividers, Content slides)
3. Apply replacement JSON
4. Add visual content (tables, boxes, diagrams)
"""
import copy
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

TEMPLATE = "project-rikyu-sales-proposals-poc/untracked/original_document/report/提案仮説構築AI_定例会1_20260219_v1.pptx"
WORKDIR = "output/13_extensibility"

# === ACES Color Palette ===
GREEN_DARK = RGBColor(0x00, 0x68, 0x43)
GREEN_MEDIUM = RGBColor(0x19, 0x7A, 0x56)
GREEN_BRIGHT = RGBColor(0x00, 0x9A, 0x62)
GREEN_MINT = RGBColor(0xC2, 0xE7, 0xD9)
GREEN_PALE = RGBColor(0xDF, 0xE6, 0xE0)
CHARCOAL = RGBColor(0x39, 0x39, 0x39)
TEXT_MAIN = RGBColor(0x12, 0x12, 0x12)
TABLE_TEXT = RGBColor(0x45, 0x45, 0x45)
BORDER_GRAY = RGBColor(0xC7, 0xC7, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x52, 0x98, 0xBA)
ACCENT_RED = RGBColor(0xDB, 0x5F, 0x5F)
ACCENT_ORANGE = RGBColor(0xFB, 0xAE, 0x40)
GRAY_LIGHT = RGBColor(0x99, 0x99, 0x99)

TITLE_FONT = "Yu Gothic"
SUB_FONT = "Yu Gothic Medium"
BODY_FONT = "Century Gothic"


# ── Helper functions ──────────────────────────────────

def find_shapes(slide):
    """Find title, subtitle, body shapes by Y-coordinate."""
    shapes_with_y = []
    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            y = s.top / 914400 if s.top else 0
            shapes_with_y.append((y, s))
    shapes_with_y.sort(key=lambda x: x[0])
    title = subtitle = body = None
    for y, s in shapes_with_y:
        if y < 0.7 and title is None: title = s
        elif 0.7 <= y < 1.5 and subtitle is None: subtitle = s
        elif y >= 1.5 and body is None: body = s
    return title, subtitle, body


def prep(slide):
    """Suppress layout bleed-through by clearing body placeholder."""
    title, subtitle, body = find_shapes(slide)
    if subtitle:
        subtitle.height = Inches(0.32)
        if subtitle.width < Inches(1):
            subtitle.width = Inches(12.46)
            subtitle.left = Inches(0.42)
    if body:
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = " "
        p.font.size = Pt(1)
        p.font.color.rgb = WHITE
    return title, subtitle, body


def add_table(slide, data, left=Inches(0.42), top=Inches(1.90), width=Inches(12.46),
              col_widths=None, header_font_pt=9, cell_font_pt=8, row_height=0.32):
    """Styled table with green header + alternating rows."""
    rows, cols = len(data), len(data[0])
    height = Inches(row_height * rows)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for para in cell.text_frame.paragraphs:
                para.font.name = BODY_FONT
                if row_idx == 0:
                    para.font.size = Pt(header_font_pt)
                    para.font.color.rgb = WHITE
                    para.font.bold = True
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.font.size = Pt(cell_font_pt)
                    para.font.color.rgb = TABLE_TEXT

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_DARK
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_PALE
    return table


def add_box(slide, text, left, top, width, height, bg_color=GREEN_MINT, text_color=GREEN_DARK,
            font_size=9, bold=False, align=PP_ALIGN.CENTER):
    """Add a colored text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_arrow_right(slide, left, top, width=Inches(0.6), height=Inches(0.3)):
    """Add right arrow text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_arrow_down(slide, left, top, width=Inches(0.5), height=Inches(0.3)):
    """Add downward arrow text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "▼"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN_DARK
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_note(slide, text, left, top, width, h=Inches(0.35)):
    """Add an italic note text."""
    txBox = slide.shapes.add_textbox(left, top, width, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(6.5)
    p.font.italic = True
    p.font.color.rgb = GRAY_LIGHT
    p.font.name = BODY_FONT


# ── Step 1: Build slide structure from template ──────

def build_slide_structure():
    """Build PPTX with correct slide structure from template."""
    prs = Presentation(TEMPLATE)

    # Template slide indices:
    # 0 = Title (Cover+data), 1 = INDEX (Index_1column), 6 = Section Cover, 21 = Content
    title_slide = prs.slides[0]
    index_slide = prs.slides[1]
    section_slide = prs.slides[6]
    content_slide = prs.slides[21]

    # We need to build a new presentation with:
    # 0: Title (copy of 0)
    # 1: INDEX (copy of 1)
    # 2: Section - LLM first (from 6)
    # 3: Content - 従来型 vs LLM first (from 21)
    # 4: Content - 5 extension axes (from 21)
    # 5: Section - ナレッジ拡張 (from 6)
    # 6: Content - Industry YAML (from 21)
    # 7: Content - Update rule + lifecycle (from 21)
    # 8: Section - ツール拡張 (from 6)
    # 9: Content - Name-based dispatch (from 21)
    # 10: Section - パイプライン拡張 (from 6)
    # 11: Content - M1-M10 (from 21)
    # 12: Section - 変更影響の局所化 (from 6)
    # 13: Content - Change impact (from 21)
    # 14: Section - 確定的×判断的 (from 6)
    # 15: Content - Rule vs LLM (from 21)

    # Structure: [Title, Index, Sec, C, C, Sec, C, C, Sec, C, Sec, C, Sec, C, Sec, C]
    # Source indices:  [0,     1,   6, 21,21, 6, 21,21, 6, 21, 6, 21, 6, 21, 6, 21]
    slide_plan = [0, 1, 6, 21, 21, 6, 21, 21, 6, 21, 6, 21, 6, 21, 6, 21]

    # Build new prs by keeping only the slides we need
    # Approach: duplicate slides from template layouts
    new_prs = Presentation(TEMPLATE)

    # Delete all existing slides first, then add from layouts
    # Actually, a simpler approach: build from slide layouts
    new_prs2 = Presentation()
    new_prs2.slide_width = prs.slide_width
    new_prs2.slide_height = prs.slide_height

    # Get slide layouts from template
    # Layout mapping from template:
    # 0=Cover, 1=Cover+data, 3=Index_1column, 4=Section Cover, 10=Contents_Title&Message_Left
    layout_cover_data = prs.slide_layouts[1]    # Cover+data
    layout_index = prs.slide_layouts[3]         # Index_1column
    layout_section = prs.slide_layouts[4]       # Section Cover
    layout_content = prs.slide_layouts[10]      # Contents_Title&Message_Left

    layout_map = {
        0: layout_cover_data,
        1: layout_index,
        6: layout_section,
        21: layout_content,
    }

    for src_idx in slide_plan:
        layout = layout_map[src_idx]
        new_prs2.slides.add_slide(layout)

    # Save working file
    working_path = f"{WORKDIR}/std-working.pptx"
    new_prs2.save(working_path)
    print(f"Built slide structure: {working_path} ({len(slide_plan)} slides)")
    return working_path


# ── Step 2: Apply replacements ──────────────────────

def apply_replacements(working_path):
    """Apply text replacements using replace.py pattern."""
    import json

    prs = Presentation(working_path)
    with open(f"{WORKDIR}/std-replacement.json") as f:
        replacements = json.load(f)

    SCRIPTS = Path("/home/node/.claude/plugins/cache/anthropic-agent-skills/document-skills/69c0b1a06741/skills/pptx/scripts")
    sys.path.insert(0, str(SCRIPTS))
    from replace import apply_paragraph_properties

    for slide_key, shapes_data in replacements.items():
        slide_idx = int(slide_key.split("-")[1])
        if slide_idx >= len(prs.slides):
            print(f"  Warning: slide {slide_idx} out of range")
            continue

        slide = prs.slides[slide_idx]

        # Find shapes by Y position
        shapes_with_y = []
        for s in slide.shapes:
            if hasattr(s, 'text_frame'):
                y = s.top / 914400 if s.top else 0
                shapes_with_y.append((y, s))
        shapes_with_y.sort(key=lambda x: x[0])

        shape_list = [s for _, s in shapes_with_y]

        for shape_key, shape_repl in shapes_data.items():
            shape_idx = int(shape_key.split("-")[1])
            if shape_idx >= len(shape_list):
                print(f"  Warning: {slide_key}/{shape_key} out of range")
                continue

            shape = shape_list[shape_idx]
            if "paragraphs" not in shape_repl:
                continue

            tf = shape.text_frame
            tf.clear()

            for i, para_data in enumerate(shape_repl["paragraphs"]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                apply_paragraph_properties(p, para_data)

    replaced_path = f"{WORKDIR}/std-replaced.pptx"
    prs.save(replaced_path)
    print(f"Applied replacements: {replaced_path}")
    return replaced_path


# ── Step 3: Add visual content ──────────────────────

def add_visual_content(replaced_path):
    """Add tables, boxes, and diagrams to content slides."""
    prs = Presentation(replaced_path)
    slides = prs.slides

    # ── Slide 3 (idx 3): 従来型 vs LLM first ──
    slide = slides[3]
    title, subtitle, body = prep(slide)

    # Two-column comparison boxes
    box_w = Inches(5.8)
    box_h = Inches(1.6)
    left_x = Inches(0.42)
    right_x = Inches(6.62)
    y_top = Inches(2.0)

    # Header boxes
    add_box(slide, "従来型", left_x, y_top, box_w, Inches(0.35),
            bg_color=CHARCOAL, text_color=WHITE, font_size=11, bold=True)
    add_box(slide, "LLM first", right_x, y_top, box_w, Inches(0.35),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=11, bold=True)

    # Content boxes
    y_content = Inches(2.40)
    add_box(slide, "ルール追加 → コード変更\n新機能追加 → 開発+テスト+デプロイ\n業界追加 → 業界固有のコード\n\n人間がif-else、正規表現、\nワークフロー定義を書いて行動を決める",
            left_x, y_content, box_w, Inches(2.0),
            bg_color=RGBColor(0xF5, 0xF5, 0xF5), text_color=CHARCOAL, font_size=9, align=PP_ALIGN.LEFT)

    add_box(slide, "ルール追加 → YAML/Markdownファイル追加\n新機能追加 → パイプライン定義追加\n業界追加 → テンプレートYAML追加\n\nLLMが自然言語の指示から行動を生成する\n人間はナレッジ(YAML/Markdown)を書くだけ",
            right_x, y_content, box_w, Inches(2.0),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=9, align=PP_ALIGN.LEFT)

    # Key insight
    add_box(slide, "ポイント: 日常的な改善（ナレッジ追加・プロンプト改善）はコード変更なし。開発サイクルから解放される。",
            Inches(0.42), Inches(4.8), Inches(12.46), Inches(0.45),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=9, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 4 (idx 4): 5つの拡張軸 ──
    slide = slides[4]
    title, subtitle, body = prep(slide)

    ext_table = [
        ["拡張軸", "変更の種類", "変更方法", "コード変更", "容易さ"],
        ["ナレッジ", "業界テンプレート追加\nルール改善", "YAML/Markdownファイルの\n追加・編集", "不要", "★★★★★"],
        ["プロンプト", "出力品質の改善\n指示の精緻化", "Markdownプロンプトの編集", "不要", "★★★★☆"],
        ["ツール", "外部API接続\nデータソース追加", "ツール定義の追加\n(名前ベースディスパッチ)", "最小", "★★★☆☆"],
        ["パイプライン", "新しいAI機能の追加\n(M8, M9...)", "パイプライン定義の追加", "必要", "★★☆☆☆"],
        ["エージェント", "専門領域の分割・追加", "エージェント定義 +\nオーケストレーター登録", "必要", "★☆☆☆☆"],
    ]
    add_table(slide, ext_table, top=Inches(1.90), col_widths=[1.5, 2.5, 3.5, 1.2, 1.5],
              header_font_pt=9, cell_font_pt=8, row_height=0.5)

    add_note(slide, "上から順に容易。日常的な改善はナレッジ/プロンプト層で完結。エージェント追加はツール数10超・ドメイン知識が複数領域にまたがる場合の分割手段。",
             Inches(0.42), Inches(5.0), Inches(12.46))

    # ── Slide 6 (idx 6): 業界テンプレートYAML ──
    slide = slides[6]
    title, subtitle, body = prep(slide)

    # YAML example as a box
    yaml_text = (
        "# knowledge_master/industry_templates/construction.yaml\n"
        "industry_code: \"D-06\"\n"
        "industry_name: \"建設業\"\n"
        "common_agendas:\n"
        "  - name: \"人材確保・定着\"\n"
        "    priority: high\n"
        "    typical_solutions: [\"人材紹介\", \"福利厚生支援\"]\n"
        "  - name: \"DX・BIM推進\"\n"
        "    priority: medium\n"
        "    typical_solutions: [\"IT投資融資\", \"DXコンサル\"]\n"
        "segment_patterns:\n"
        "  - \"土木 / 建築\"  - \"元請 / 下請\""
    )
    add_box(slide, yaml_text,
            Inches(0.42), Inches(1.90), Inches(7.5), Inches(3.2),
            bg_color=RGBColor(0x2D, 0x2D, 0x2D), text_color=RGBColor(0xE0, 0xE0, 0xE0),
            font_size=8, align=PP_ALIGN.LEFT)

    # Explanation box on right
    add_box(slide, "追加手順",
            Inches(8.3), Inches(1.90), Inches(4.5), Inches(0.30),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=9, bold=True)

    steps_data = [
        ["#", "ステップ", "所要時間"],
        ["1", "業界の課題をヒアリング", "1-2時間"],
        ["2", "YAMLファイルを作成", "30分"],
        ["3", "Knowledge Masterに配置", "5分"],
        ["4", "AIが自動的に業界知識を活用", "即座"],
    ]
    add_table(slide, steps_data, left=Inches(8.3), top=Inches(2.25), width=Inches(4.5),
              col_widths=[0.4, 2.5, 1.2], header_font_pt=8, cell_font_pt=7, row_height=0.28)

    add_box(slide, "コード変更: 不要\nデプロイ: 不要\nテスト: AIの出力確認のみ",
            Inches(8.3), Inches(3.8), Inches(4.5), Inches(0.75),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=8, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 7 (idx 7): 更新ルール + ライフサイクル ──
    slide = slides[7]
    title, subtitle, body = prep(slide)

    # Update rule YAML box
    yaml2 = (
        "# knowledge_master/update_rules/trigger_patterns.yaml\n"
        "triggers:\n"
        "  - pattern: \"決裁権者が具体的な時期に言及\"\n"
        "    action: \"時期を更新 + 優先度UP\"\n"
        "    confidence: high\n"
        "  - pattern: \"担当者レベルの伝聞\"\n"
        "    action: \"備考に追記のみ\"\n"
        "    confidence: low"
    )
    add_box(slide, yaml2,
            Inches(0.42), Inches(1.90), Inches(6.5), Inches(2.2),
            bg_color=RGBColor(0x2D, 0x2D, 0x2D), text_color=RGBColor(0xE0, 0xE0, 0xE0),
            font_size=8, align=PP_ALIGN.LEFT)

    # Lifecycle diagram
    add_box(slide, "ナレッジのライフサイクル",
            Inches(7.3), Inches(1.90), Inches(5.5), Inches(0.30),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=9, bold=True)

    steps = [
        ("エキスパートヒアリング", GREEN_PALE),
        ("ナレッジの言語化", GREEN_PALE),
        ("YAML/Markdownに構造化", GREEN_MINT),
        ("Knowledge Masterに投入", GREEN_MINT),
        ("AIの出力品質が向上", RGBColor(0xC2, 0xE7, 0xD9)),
        ("RMのFBでさらに改善", RGBColor(0xC2, 0xE7, 0xD9)),
    ]
    y_step = 2.30
    for i, (label, color) in enumerate(steps):
        add_box(slide, label,
                Inches(7.3), Inches(y_step), Inches(5.5), Inches(0.30),
                bg_color=color, text_color=CHARCOAL, font_size=8)
        y_step += 0.30
        if i < len(steps) - 1:
            add_arrow_down(slide, Inches(9.8), Inches(y_step - 0.02), Inches(0.5), Inches(0.20))
            y_step += 0.15

    # ── Slide 9 (idx 9): ツール拡張 — 名前ベースディスパッチ ──
    slide = slides[9]
    title, subtitle, body = prep(slide)

    # Dispatch diagram
    add_box(slide, "tool_registry (ディスパッチマップ)",
            Inches(0.42), Inches(1.90), Inches(5.5), Inches(0.30),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=9, bold=True)

    tools = [
        ("web_search", "Web検索ツール"),
        ("edinet_api", "EDINET APIツール"),
        ("tsr_api", "TSR APIツール"),
        ("business_card_api", "名刺アプリ連携 (将来)"),
        ("financial_model", "財務分析ツール (将来)"),
    ]
    y_tool = 2.30
    for name, desc in tools:
        is_future = "将来" in desc
        bg = RGBColor(0xEE, 0xF1, 0xF5) if is_future else GREEN_MINT
        border_color = ACCENT_BLUE if is_future else GREEN_DARK
        add_box(slide, f'"{name}": {desc}',
                Inches(0.42), Inches(y_tool), Inches(5.5), Inches(0.28),
                bg_color=bg, text_color=CHARCOAL, font_size=8, align=PP_ALIGN.LEFT)
        y_tool += 0.32

    add_box(slide, "新しいツールは1行追加するだけ\nループ本体の変更は一切不要",
            Inches(0.42), Inches(y_tool + 0.10), Inches(5.5), Inches(0.50),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=8, bold=True, align=PP_ALIGN.LEFT)

    # Future tools table
    future_tools = [
        ["ツール", "接続先", "追加方法", "影響範囲"],
        ["名刺アプリ連携", "名刺管理サービスAPI", "ツール定義追加", "KPマップのデータソース拡充"],
        ["SFA/CRM連携", "Salesforce等", "ツール定義追加", "面談・商談データの自動取得"],
        ["財務分析ツール", "既存会計モデル", "ツール定義追加", "M5提案の精度向上"],
        ["パワポ出力", "python-pptx等", "ツール定義追加", "出力フォーマットの追加"],
    ]
    add_table(slide, future_tools, left=Inches(6.3), top=Inches(1.90), width=Inches(6.5),
              col_widths=[1.6, 1.8, 1.3, 1.8], header_font_pt=8, cell_font_pt=7, row_height=0.30)

    # ── Slide 11 (idx 11): パイプライン拡張 M1-M10 ──
    slide = slides[11]
    title, subtitle, body = prep(slide)

    # Pipeline tree
    add_box(slide, "Pipeline Orchestrator",
            Inches(0.42), Inches(1.90), Inches(3.5), Inches(0.30),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=9, bold=True)

    pipelines = [
        ("M1", "アジェンダ初期生成", True),
        ("M2", "面談→アジェンダ更新", True),
        ("M3", "ニュース検出", True),
        ("M4", "AIチャット", True),
        ("M5", "提案ストーリー", True),
        ("M6", "深掘りペーパー", True),
        ("M7", "担当者用メモ", True),
        ("M8", "アカウントプラン生成", False),
        ("M9", "競合分析レポート", False),
        ("M10", "研修シミュレーション", False),
    ]
    y_pipe = 2.30
    for code, name, existing in pipelines:
        bg = GREEN_MINT if existing else RGBColor(0xEE, 0xF1, 0xF5)
        label = f"├── {code} {name}" if code != "M10" else f"└── {code} {name}"
        suffix = "← 既存" if existing else "← 将来追加"
        add_box(slide, f"{label}  {suffix}",
                Inches(0.42), Inches(y_pipe), Inches(3.8), Inches(0.22),
                bg_color=bg, text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)
        y_pipe += 0.23

    # Requirements table
    req_table = [
        ["要素", "内容", "形式"],
        ["パイプライン定義", "入力→処理ステップ→出力の定義", "Python"],
        ["プロンプトテンプレート", "LLMへの指示", "Markdown"],
        ["Knowledge参照定義", "どのK1-K8を使うか", "YAML"],
        ["API定義", "エンドポイント、パラメータ", "OpenAPI"],
    ]
    add_table(slide, req_table, left=Inches(4.7), top=Inches(1.90), width=Inches(8.0),
              col_widths=[2.2, 3.5, 1.2], header_font_pt=8, cell_font_pt=7, row_height=0.28)

    add_box(slide, "ポイント: 新しいスキルを追加しても既存パイプラインに影響しない。独立したスキルとして並列実行可能。",
            Inches(4.7), Inches(3.5), Inches(8.0), Inches(0.45),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=8, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 13 (idx 13): 変更影響の局所化 ──
    slide = slides[13]
    title, subtitle, body = prep(slide)

    # Impact boxes - showing isolation
    changes = [
        ("ナレッジ追加・修正", "YAMLファイルのみ\n他のコードに影響なし", GREEN_MINT, "不要"),
        ("プロンプト改善", "Markdownファイルのみ\n他のコードに影響なし", GREEN_MINT, "不要"),
        ("ツール追加", "ツール定義 + ディスパッチ登録\nループ構造に影響なし", GREEN_PALE, "最小"),
        ("パイプライン追加", "新パイプライン定義\n既存パイプラインに影響なし", RGBColor(0xEE, 0xF1, 0xF5), "必要"),
        ("エージェント追加", "エージェント定義 +\nオーケストレーター登録", RGBColor(0xEE, 0xF1, 0xF5), "必要"),
        ("LLMモデル変更", "設定ファイルのみ\nアルゴリズム構造に影響なし", GREEN_PALE, "不要"),
    ]

    y_change = 1.90
    for label, impact, bg, code_change in changes:
        # Left: change type
        add_box(slide, label,
                Inches(0.42), Inches(y_change), Inches(2.5), Inches(0.55),
                bg_color=bg, text_color=GREEN_DARK, font_size=9, bold=True)
        # Arrow
        add_arrow_right(slide, Inches(3.05), Inches(y_change + 0.10))
        # Right: impact
        add_box(slide, impact,
                Inches(3.75), Inches(y_change), Inches(6.5), Inches(0.55),
                bg_color=WHITE, text_color=CHARCOAL, font_size=8, align=PP_ALIGN.LEFT)
        # Code change badge
        badge_bg = GREEN_DARK if code_change == "不要" else (ACCENT_ORANGE if code_change == "最小" else ACCENT_BLUE)
        add_box(slide, f"コード変更: {code_change}",
                Inches(10.5), Inches(y_change + 0.10), Inches(2.3), Inches(0.30),
                bg_color=badge_bg, text_color=WHITE, font_size=7, bold=True)
        y_change += 0.62

    add_box(slide, "設計原則: 各レイヤーが疎結合であるため、変更の影響が他のレイヤーに波及しない",
            Inches(0.42), Inches(y_change + 0.15), Inches(12.46), Inches(0.40),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=9, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 15 (idx 15): 確定的 × 判断的 ──
    slide = slides[15]
    title, subtitle, body = prep(slide)

    rule_table = [
        ["レイヤー", "性質", "rikyuでの対応", "例"],
        ["確定的", "必ずこう動く。例外なし", "ルールエンジン(YAML)", "信頼度マトリクス\n更新ルールの確定パターン"],
        ["判断的", "できるだけこう動く", "LLM + プロンプト", "文脈解釈、カスタマイズ\n新規パターン検出"],
        ["キャプチャ", "RM操作を必ず記録", "Decision Trace", "承認/修正/棄却+理由の記録"],
    ]
    add_table(slide, rule_table, top=Inches(1.90), col_widths=[1.5, 2.5, 2.8, 4.5],
              header_font_pt=9, cell_font_pt=8, row_height=0.50)

    # Key insight
    add_box(slide, "「絶対に守らせたい」はルールエンジンに、「できるだけ従ってほしい」はプロンプトに。\nこの区別がシステムの信頼性と柔軟性を両立させる。",
            Inches(0.42), Inches(4.0), Inches(12.46), Inches(0.55),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=10, bold=True, align=PP_ALIGN.LEFT)

    # Visual comparison boxes
    add_box(slide, "確定的\n(ルールエンジン)\n\n・信頼度マトリクス\n・更新ルールの確定パターン\n・入力バリデーション",
            Inches(0.42), Inches(4.8), Inches(4.0), Inches(2.0),
            bg_color=RGBColor(0xD6, 0xEA, 0xF8), text_color=CHARCOAL, font_size=8, align=PP_ALIGN.LEFT)

    add_box(slide, "判断的\n(LLM + プロンプト)\n\n・文脈に応じた解釈\n・新規パターンの検出\n・カスタマイズされた提案",
            Inches(4.72), Inches(4.8), Inches(4.0), Inches(2.0),
            bg_color=GREEN_MINT, text_color=CHARCOAL, font_size=8, align=PP_ALIGN.LEFT)

    add_box(slide, "キャプチャ\n(Decision Trace)\n\n・RM操作の記録\n・承認/修正/棄却+理由\n・フィードバックループ",
            Inches(9.02), Inches(4.8), Inches(3.8), Inches(2.0),
            bg_color=GREEN_PALE, text_color=CHARCOAL, font_size=8, align=PP_ALIGN.LEFT)

    # Save
    output_path = "output/13_extensibility_design.pptx"
    prs.save(output_path)
    print(f"Saved standard version: {output_path}")
    return output_path


# ── Main ──────────────────────────────────────────────

if __name__ == "__main__":
    working_path = build_slide_structure()
    replaced_path = apply_replacements(working_path)
    add_visual_content(replaced_path)
