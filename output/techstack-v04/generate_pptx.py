#!/usr/bin/env python3
"""
Generate tech_stack_adoption_policy_v04.pptx (normal) and _detailed.pptx
from template slides.
"""
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = "/workspace/aces-rikyu-sales-proposals-poc/untracked/original_document/report/提案仮説構築AI_定例会1_20260219_v1.pptx"
OUT_NORMAL = "/workspace/output/techstack-v04/tech_stack_adoption_policy_v04.pptx"
OUT_DETAILED = "/workspace/output/techstack-v04/tech_stack_adoption_policy_v04_detailed.pptx"

# ACES colors
HEADER_BG = RGBColor(0x00, 0x68, 0x43)
HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
DATA_TEXT = RGBColor(0x45, 0x45, 0x45)
ALT_ROW_BG = RGBColor(0xDF, 0xE6, 0xE0)
BORDER_COLOR = RGBColor(0xC7, 0xC7, 0xC7)
ACCENT_GREEN = RGBColor(0x00, 0x68, 0x43)

# Template slide indices
TITLE_SLIDE = 0    # Cover slide
INDEX_SLIDE = 1    # Index slide
SECTION_SLIDE = 6  # Section divider
CONTENT_SLIDE = 21 # Content with title + message + body

# ------ Slide 21 placeholder names ------
# shape-1 = title (pos 386366, 400248)
# shape-2 = subtitle/message (pos 386366, 831135)
# shape-6 = body (pos 386366, 1634624)


def duplicate_slide(prs, slide_index):
    """Duplicate template slide by copying XML and relationships."""
    template_slide = prs.slides[slide_index]
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Remove default placeholders from new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy shapes from template
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    return new_slide


def clear_and_set_text(shape, text, font_name="Yu Gothic", font_size=None, bold=None, color=None, alignment=None):
    """Clear shape text and set new text with formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # Remove all paragraphs except the first one via XML
    txBody = tf._txBody
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    paras = txBody.findall('.//a:p', nsmap)
    # Keep only the first paragraph, remove the rest
    for p in paras[1:]:
        txBody.remove(p)
    # Clear the first paragraph's runs
    first_p = paras[0]
    for r in first_p.findall('.//a:r', nsmap):
        first_p.remove(r)
    # Also remove any break elements
    for br in first_p.findall('.//a:br', nsmap):
        first_p.remove(br)

    # Set text
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        run.text = line
        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = font_size
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color
        if alignment:
            para.alignment = alignment


def set_body_bullets(shape, items, font_name="Yu Gothic Medium", font_size=Pt(10), color=None):
    """Set body text as bullet points."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # Remove all paragraphs except the first one via XML
    txBody = tf._txBody
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    paras = txBody.findall('.//a:p', nsmap)
    for p in paras[1:]:
        txBody.remove(p)
    first_p = paras[0]
    for r in first_p.findall('.//a:r', nsmap):
        first_p.remove(r)
    for br in first_p.findall('.//a:br', nsmap):
        first_p.remove(br)

    for i, item in enumerate(items):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        # Support indentation via tuple (level, text)
        if isinstance(item, tuple):
            level, text = item
            para.level = level
        else:
            text = item
            para.level = 0

        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = font_size
        if color:
            run.font.color.rgb = color
        else:
            run.font.color.rgb = DATA_TEXT


def add_table(slide, rows_data, col_widths, left=Emu(386366), top=Emu(1790000),
              header_bg=HEADER_BG, header_text=HEADER_TEXT, font_size=Pt(8)):
    """Add a table to slide."""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0]) if rows_data else 0
    width = sum(col_widths)
    height = Emu(250000) * num_rows

    # Limit table height
    max_height = Emu(4800000)
    if height > max_height:
        height = max_height

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, Emu(width), height)
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w)

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = str(cell_text)
            run.font.size = font_size
            run.font.name = "Yu Gothic"

            if r == 0:
                # Header
                run.font.color.rgb = header_text
                run.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                run.font.color.rgb = DATA_TEXT
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ALT_ROW_BG
                else:
                    cell.fill.background()

    return table_shape


def get_shape_by_index(slide, idx):
    """Get shape by index."""
    shapes = list(slide.shapes)
    if idx < len(shapes):
        return shapes[idx]
    return None


def make_title_slide(prs):
    """Slide 0: Title slide"""
    slide = duplicate_slide(prs, TITLE_SLIDE)
    s0 = get_shape_by_index(slide, 0)
    s1 = get_shape_by_index(slide, 1)
    clear_and_set_text(s0, "SOMPOケア株式会社 御中\n技術スタック採用方針 v0.4", font_size=Pt(24), bold=True)
    clear_and_set_text(s1, "2026年4月 | ACES Inc.", font_size=Pt(12))
    return slide


def make_index_slide(prs, items):
    """Slide 1: INDEX slide"""
    slide = duplicate_slide(prs, INDEX_SLIDE)
    # Index uses shape-1 for content
    s1 = get_shape_by_index(slide, 1)
    clear_and_set_text(s1, "\n".join(items), font_size=Pt(12))
    return slide


def make_section_slide(prs, title):
    """Section divider slide"""
    slide = duplicate_slide(prs, SECTION_SLIDE)
    s1 = get_shape_by_index(slide, 1)
    clear_and_set_text(s1, title, font_size=Pt(22), bold=True)
    return slide


def make_content_slide(prs, title, message, body_items=None, table_data=None, col_widths=None):
    """Content slide with title, message, and body (bullets or table)."""
    slide = duplicate_slide(prs, CONTENT_SLIDE)
    shapes = list(slide.shapes)

    # shape-1 = title (index 1 in template)
    # shape-2 = message/subtitle (index 2)
    # shape-6 = body (index 6)
    # Find placeholders by name pattern
    title_shape = None
    message_shape = None
    body_shape = None

    for s in shapes:
        if not s.has_text_frame:
            continue
        name = s.name
        if 'プレースホルダー 2' in name:
            title_shape = s
        elif 'プレースホルダー 3' in name:
            message_shape = s
        elif 'プレースホルダー 4' in name:
            body_shape = s

    if title_shape:
        clear_and_set_text(title_shape, title, font_name="Yu Gothic", font_size=Pt(14), bold=True)
    if message_shape:
        clear_and_set_text(message_shape, message, font_name="Yu Gothic Medium", font_size=Pt(10),
                          color=ACCENT_GREEN, bold=True)

    if table_data and col_widths:
        # Clear body text and add table
        if body_shape:
            clear_and_set_text(body_shape, "", font_size=Pt(1))
        add_table(slide, table_data, col_widths, top=Emu(1790000))
    elif body_items:
        if body_shape:
            set_body_bullets(body_shape, body_items)

    return slide


def remove_template_slides(prs, count):
    """Remove the first N template slides."""
    for _ in range(count):
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            # Try alternate attribute
            rId_elem = prs.slides._sldIdLst[0]
            for attr_name in rId_elem.attrib:
                if 'id' in attr_name.lower() and attr_name != 'id':
                    rId = rId_elem.attrib[attr_name]
                    break
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def build_normal_version():
    """Build the normal (~22 slides) version."""
    prs = Presentation(TEMPLATE_PATH)
    original_count = len(prs.slides)

    # --- Slide 0: Title ---
    make_title_slide(prs)

    # --- Slide 1: INDEX ---
    make_index_slide(prs, [
        "Part 0: SOMPOケアが目指す姿",
        "Part 1: Foundry移行の設計判断",
        "Part 2: 構造課題と対応方針",
        "Part 3: 技術スタック採用4原則",
        "Part 4: RFI評価と推奨",
    ])

    # --- Slide 2: Section "Part 0" ---
    make_section_slide(prs, "Part 0: SOMPOケアが目指す姿")

    # --- Slide 3: 3つのゴール ---
    make_content_slide(prs,
        title="SOMPOケアが目指す姿 -- 3つのゴール",
        message='本PJTは「Foundryの引っ越し」ではない。3つのゴールを同時に達成する設計が必要',
        body_items=[
            "G1 確実な移行: Foundryからの脱却と現行業務の継続",
            (1, "CMC/モニター/ナースコール分析/施設360の既存BI継続"),
            (1, "1,000+施設の日常モニタリングが中断しないこと"),
            "",
            "G2 自律運用: 自律的に運用できる基盤への転換",
            (1, "「ベンダーしか知らない」状態からの脱却"),
            (1, "特別なSEを抱えなくても分析が回せる運用モデル"),
            (1, "推奨運用体制: 2名程度(データエンジニア中心)"),
            "",
            "G3 AI拡張: AI活用基盤としての拡張性確保",
            (1, "DDA17プロジェクト全ての土台となる設計"),
            (1, "介護記録テキスト分析、カメラ映像AI解析等の将来ユースケース"),
        ]
    )

    # --- Slide 4: DDA戦略 ---
    make_content_slide(prs,
        title="DDA戦略と新基盤の位置づけ",
        message="新分析環境はDDA17プロジェクト全ての土台。今選ぶ基盤が将来のAI活用の天井になりうる",
        table_data=[
            ["対象", "目指す価値"],
            ["ご利用者さま", "自立支援、カスタムメイドケア、ACP"],
            ["社員", "ムリ・ムダ・ムラ軽減、働きがい向上"],
            ["未来社会", "需給Gap解消、地域社会連携"],
            ["", ""],
            ["DDA 5段階取組み", "内容"],
            ["1", "アセスメント音声自動フォーマット化"],
            ["2", "ケアプラン自動化・標準化"],
            ["3", "介護S/W統一(NDS CareBase)"],
            ["4", "カメラ・センサーの最大活用"],
            ["5", "介護リアルデータのカスタムメイドケア活用"],
        ],
        col_widths=[3000000, 8000000]
    )

    # --- Slide 5: 一気通貫 ---
    make_content_slide(prs,
        title="ゴール -> 4原則 -> RFP評価軸の一気通貫",
        message="As-Is課題 -> To-Be -> 4原則 -> RFP評価軸が論理的に繋がっている",
        body_items=[
            "As-Is課題:",
            (1, "OverStack/高コスト -> G1: 確実な移行 -> 原則3 ガバナンス -> 統合カタログ"),
            (1, "属人化/BBox化 -> G2: 自律運用 -> 原則4 ナレッジ蓄積 -> セルフサービス性"),
            (1, "AI環境なし -> G3: AI拡張 -> 原則1 AI接続 -> SQL API"),
            (1, "非構造未対応 -> G3: AI拡張 -> 原則2 学習蓄積 -> MLflow統合"),
            "",
            "この構造により、RFPの評価軸一つ一つが",
            "「なぜ必要か」をAs-Is課題まで遡って説明できる。",
        ]
    )

    # --- Slide 6: Section "Foundry移行" ---
    make_section_slide(prs, "Foundry移行の設計判断")

    # --- Slide 7: 4レイヤー分解 ---
    make_content_slide(prs,
        title="基盤の4レイヤー分解",
        message="L0-L1はHD共通化、L2-L3はケア独自設計。本資料はL2-L3の設計方針を定める",
        table_data=[
            ["レイヤー", "内容", "HD基盤に乗るか", "判断根拠"],
            ["L0: インフラ", "AWS, VPC, IdP", "乗る", "規模の経済"],
            ["L1: ストレージ", "Databricks/Snowflake", "条件付き", "マルチテナント設計が前提"],
            ["L2: データアーキ", "スキーマ, セマンティック", "乗れない", "ドメイン固有"],
            ["L3: アプリ", "BI, AI/ML環境", "乗れない", "ケアのユースケース直結"],
        ],
        col_widths=[2000000, 2800000, 2200000, 4000000]
    )

    # --- Slide 8: Foundryとの差分 ---
    make_content_slide(prs,
        title="Foundryとの差分整理",
        message="Foundryの課題を削ぎ落とし、意図的に変えた設計。同じものを作り直すのではない",
        table_data=[
            ["区分", "項目", "内容"],
            ["引き継ぐ", "Bronze Layer", "生データを新基盤に全量移行"],
            ["引き継ぐ", "オントロジー", "セマンティックレイヤー/dbtとして再定義"],
            ["引き継ぐ", "KPI定義", "セマンティックモデルとしてコード化"],
            ["引き継ぐ", "パイプラインロジック", "デコレーター除去のみでほぼ動作(219本)"],
            ["変える", "OverStack", "必要機能に絞ったコスト最適構成"],
            ["変える", "ベンダーロックイン", "Delta Lake/Parquet + OSS中心構成"],
            ["変える", "セルフBI不可", "SQL中心 + 直感的BIツール"],
            ["変える", "AI実験環境なし", "MLflow + Model Registry + Git連携"],
            ["変える", "独自言語体系", "SQL/Python標準 + dbt"],
        ],
        col_widths=[1500000, 2500000, 7000000]
    )

    # --- Slide 9: ガバナンス代替 ---
    make_content_slide(prs,
        title="ガバナンス機能の代替マッピング",
        message="Foundryの5機能中3機能は完全代替可能。残り2機能はCI/CDとOSSで補完",
        table_data=[
            ["Foundry機能", "新基盤での代替", "適合度", "備考"],
            ["Markings(タグ自動伝播)", "Tags(手動/API)", "△", "自動伝播なし、CI/CDで補完"],
            ["RLS(行レベルセキュリティ)", "行フィルタ+列マスキング", "◎", "Unity Catalog標準機能"],
            ["自動リネージ", "Unity Catalog Lineage", "○", "PySpark GA済"],
            ["監査ログ", "System Tables", "◎", "SQLで柔軟に分析可"],
            ["PII検出", "Presidio(OSS)/AWS Macie", "△", "別途導入必要"],
        ],
        col_widths=[2800000, 3000000, 1000000, 4200000]
    )

    # --- Slide 10: Section "構造課題" ---
    make_section_slide(prs, "構造課題と対応方針")

    # --- Slide 11: 課題x4原則マトリクス ---
    make_content_slide(prs,
        title="課題 x 4原則 対応マトリクス",
        message="As-Is課題がどの原則・設計判断で打ち取られるかを明示",
        table_data=[
            ["As-Is課題", "対応原則", "解決策", "RFP評価軸"],
            ["セルフBI不可", "4.ナレッジ蓄積", "SQL中心+セルフサービスBI", "学習曲線"],
            ["属人化/BBox化", "4.ナレッジ蓄積", "パイプラインGUI可視化+KPIコード管理", "運用ダッシュボード"],
            ["OverStack/高コスト", "4.ナレッジ蓄積", "必要機能に絞った従量課金構成", "コスト透明性"],
            ["SSOT不在", "1.AI接続", "セマンティックレイヤーでKPI一元管理", "セマンティックモデル"],
            ["AI実験環境なし", "2.学習蓄積", "MLflow+ModelRegistry+Git連携", "実験管理統合度"],
            ["ガバナンス都度構築", "3.ガバナンス", "統合カタログでアクセス制御一元管理", "統合カタログ"],
        ],
        col_widths=[2500000, 2000000, 3500000, 3000000]
    )

    # --- Slide 12: 非構造化データ ---
    make_content_slide(prs,
        title="非構造化データの取り扱い方針",
        message="基盤に入れるべきデータと外部に任せるデータの判断基準を定義",
        body_items=[
            "基盤に入れるべきもの:",
            (1, "SOMPOケア固有のドメインデータ(約28,800室 x 290施設)"),
            (1, "複数ソースを横断して分析するために結合が必要なデータ"),
            (1, "時系列で蓄積し、AIの学習データとして使うもの"),
            (1, "外部ソリューションが分析した「結果」(構造化されたインサイト)"),
            "",
            "基盤に入れなくてよいもの:",
            (1, "外部ベンダーの分析環境で完結する生データ(カメラ映像等)"),
            (1, "個人のメモ・下書き等、品質管理されていないデータ"),
            "",
            "外部ソリューション連携: カメラ映像 -> 外部ベンダーで分析",
            (1, "-> 構造化された結果(転倒イベント, ADLスコア等)のみ基盤に取り込み"),
            (1, "-> 他のデータと結合して施設横断の分析・予測に活用"),
        ]
    )

    # --- Slide 13: Section "4原則" ---
    make_section_slide(prs, "技術スタック採用4原則")

    # --- Slide 14: 4原則全体像 ---
    make_content_slide(prs,
        title="技術スタック採用4原則の全体像",
        message="個別製品の比較ではなく、製品に依存しない4つの設計原則でRFPの評価軸を定める",
        body_items=[
            "原則1: AIと対話できるインターフェースを持つ基盤",
            (1, "SQL API / MCP / セマンティックモデル / マルチモデル対応"),
            "",
            "原則2: 組織の学習が蓄積する基盤",
            (1, "実験管理(MLflow) / コード管理(Git) / 評価基盤 / エージェント管理"),
            "",
            "原則3: ガバナンスが拡張可能な基盤",
            (1, "権限継承 / 監査ログ / データ分類 / SSOT"),
            "",
            "原則4: ナレッジが社内に残る基盤",
            (1, "段階的自走 / 運用透明性 / 変更安全性 / ナレッジ移転"),
            "",
            "なぜこの4つか: 技術は半年で世代交代する。だが、4原則が守るものは陳腐化しない。",
            (1, "データの蓄積、データ定義、ガバナンスの設計思想、組織のナレッジ蓄積"),
        ]
    )

    # --- Slide 15: 7つの評価観点 ---
    make_content_slide(prs,
        title="7つの評価観点(SOMPOケア具体化版)",
        message="1-2は今回選定、3-6は余地確認、7は将来検討",
        table_data=[
            ["観点", "分類", "SOMPOケア固有の確認事項"],
            ["1.連携性", "今回選定", "ほのぼの/Axist/眠りSCAN接続、CareBase刷新対応"],
            ["2.安全性", "今回選定", "1,000施設の行RLS、要配慮個人情報分類、監査ログ"],
            ["3.進化性", "余地確認", "Delta Lake/Parquet、マルチモデル、非構造対応余地"],
            ["4.実効性", "余地確認", "眠りSCAN 1.92TB移行、219パイプライン移植"],
            ["5.運用性", "余地確認", "SEなしで運用可能か、2名体制、GUI完結度"],
            ["6.開発速度", "余地確認", "PoC環境構築2週間以内、MLflow/Git連携"],
            ["7.規模", "将来検討", "AIエージェント統一管理、部門別コスト可視化"],
        ],
        col_widths=[1800000, 1800000, 7400000]
    )

    # --- Slide 16: 運用モデル ---
    make_content_slide(prs,
        title="運用モデル -- 自走のロードマップ",
        message="特別なSEを抱えなくても分析が回せる。Phase 1-4で段階的に内製化範囲を拡大",
        body_items=[
            "Phase 1(構築~半年): ベンダー主導、社内は学習",
            (1, "Unity Catalog操作、パイプライン監視の習得"),
            (1, "「なぜこう設計したか」のドキュメント = 最重要成果物"),
            "",
            "Phase 2(半年~1年): 日常運用を社内に移管",
            (1, "KPI定義変更、権限管理、パイプライン監視を社内で実施"),
            "",
            "Phase 3(1年~): 社内が主導、ベンダーは大規模改修のみ",
            (1, "新規ユースケースの要件定義も社内主導"),
            "",
            "Phase 4(2年~): LLMを活用した内製化の拡大",
            (1, "LLMツール(Cursor/Claude Code等)でパイプライン修正・AI PoC"),
            (1, "ドメイン知識を持つ社内人材がAI開発の主体に"),
        ]
    )

    # --- Slide 17: RFPスコープ ---
    make_content_slide(prs,
        title="RFPスコープの明確化",
        message="必須/加点/将来検討の3段階で、ベンダーに求める回答レベルを明示",
        table_data=[
            ["観点", "分類", "必須", "加点", "将来検討"],
            ["1.連携性", "今回選定", "ソース接続, CB刷新", "--", "HD基盤, IoT"],
            ["2.安全性", "今回選定", "行/列RLS, 監査ログ", "--", "時間軸, AI経由"],
            ["3.進化性", "余地確認", "オープンフォーマット", "セマンティック構築", "マルチモデル"],
            ["4.実効性", "余地確認", "BI再現, データ移行", "本番相当PoC", "--"],
            ["5.運用性", "余地確認", "スキル要件, 自走RM", "ナレッジ移転", "--"],
            ["6.開発速度", "余地確認", "--", "PoC, MLflow, Git", "--"],
            ["7.規模", "将来検討", "--", "--", "エージェント管理"],
        ],
        col_widths=[1500000, 1500000, 3000000, 2500000, 2500000]
    )

    # --- Slide 18: Section "RFI評価" ---
    make_section_slide(prs, "RFI評価と推奨")

    # --- Slide 19: RFI 5社評価 ---
    make_content_slide(prs,
        title="重み付きRFI 5社評価",
        message="Foundry移行実現性と運用自走性の重みが最大。Databricksを推奨するがベンダー=/=ツール",
        table_data=[
            ["評価観点", "重み", "富士通(DB)", "デロイト(DB)", "日立(SF)", "PwC(AWS)", "Salesforce"],
            ["Foundry移行実現性", "★5", "◎", "○", "△", "△", "x"],
            ["運用の自走性", "★5", "○", "○", "◎", "△", "○"],
            ["統合カタログ", "★4", "◎", "◎", "○", "△", "△"],
            ["コスト構造", "★4", "○", "○", "○", "○", "△"],
            ["セルフサービスBI", "★3", "○", "○", "◎", "△", "○"],
            ["SQL API/セマンティック", "★3", "◎", "◎", "○", "△", "x"],
            ["MLflow/実験管理", "★2", "◎", "○", "△", "△", "x"],
            ["非構造データ対応", "★2", "○", "○", "○", "○", "x"],
            ["マルチモデル", "★1", "○", "○", "○", "○", "△"],
        ],
        col_widths=[2200000, 800000, 1400000, 1600000, 1400000, 1400000, 1400000]
    )

    # --- Slide 20: 推奨技術スタック ---
    make_content_slide(prs,
        title="推奨技術スタック総括",
        message="Databricks + dbt + Tableau。PySpark移植コスト最小、2名程度で運用可能",
        body_items=[
            "推奨構成: Databricks + dbt Core + Tableau",
            (1, "PySpark互換: デコレーター除去のみでほぼ移行可能(G1直結)"),
            (1, "Unity Catalog: 1,000施設の行RLS+要配慮個人情報を一元管理"),
            (1, "Delta Lake標準: オープンフォーマット、将来の基盤変更時にデータ持出し可能"),
            (1, "AI-BI Genie: セマンティックモデルベースの自然言語KPI照会"),
            (1, "MLflow統合: 追加ライセンス不要でAI実験管理が利用可能"),
            (1, "従量課金: DBU課金でOverStack回避"),
            "",
            "リスク:",
            (1, "DBU単価の不透明性(値引き交渉が必要)"),
            (1, "Markings自動伝播の欠如(CI/CDで補完)"),
            (1, "セマンティックレイヤー初期構築の組織的合意形成コスト"),
            "",
            "総合評価: 富士通(A) / デロイト(A-) / 日立(B+) / PwC(B) / Salesforce(C)",
        ]
    )

    # --- Slide 21: Next Steps ---
    make_content_slide(prs,
        title="Next Steps",
        message="やるべきことは明確。設計判断レベルの追加コストで対応可能",
        table_data=[
            ["#", "ステップ", "タイミング"],
            ["1", "4原則のRFP反映", "RFP策定時"],
            ["2", "HDとのレイヤー合意(L0-L1共通化範囲)", "HD基盤検討と並行"],
            ["3", "ソースシステムの繋ぎ込み調査", "要件定義フェーズ"],
            ["4", "Foundryオントロジー棚卸し", "RFP策定時"],
            ["5", "PoC候補の具体化(自然言語KPI照会)", "RFP策定と並行"],
            ["6", "マスタデータ名寄せ方針", "移行設計フェーズ"],
            ["7", "介護記録テキストデータの実態調査(340GB)", "PoC設計前"],
            ["8", "横断人材の確保方針", "Phase 1開始前"],
            ["9", "運用自走ロードマップの策定", "RFP策定時"],
        ],
        col_widths=[500000, 6500000, 4000000]
    )

    # Remove original template slides
    remove_template_slides(prs, original_count)
    prs.save(OUT_NORMAL)
    print(f"Saved normal version: {OUT_NORMAL}")


def build_detailed_version():
    """Build the detailed (~14 slides) version - merged slides, no section dividers."""
    prs = Presentation(TEMPLATE_PATH)
    original_count = len(prs.slides)

    # --- Slide 0: Title ---
    make_title_slide(prs)

    # --- Slide 1: INDEX ---
    make_index_slide(prs, [
        "SOMPOケアが目指す姿 + DDA戦略",
        "ゴール -> 4原則 -> RFP評価軸",
        "4レイヤー分解 + Foundryとの差分",
        "ガバナンス機能の代替マッピング",
        "課題 x 4原則 対応マトリクス",
        "非構造化データの取り扱い方針",
        "技術スタック採用4原則 + 7つの評価観点",
        "運用モデル + RFPスコープ",
        "RFI 5社評価 + 推奨技術スタック",
        "Next Steps",
    ])

    # --- Slide 2: ゴール + DDA (merged 3+4) ---
    make_content_slide(prs,
        title="SOMPOケアが目指す姿 + DDA戦略",
        message='本PJTは「Foundryの引っ越し」ではない。3つのゴールを同時に達成し、DDA17プロジェクトの土台を構築する',
        table_data=[
            ["ゴール", "概要", "DDA戦略との関連"],
            ["G1 確実な移行", "Foundry脱却+現行BI継続(CMC/モニター/NC/施設360)", "全ての前提条件"],
            ["G2 自律運用", "ベンダー依存脱却、2名程度で運用可能", "三方良しの社員向け価値実現"],
            ["G3 AI拡張", "DDA17PJの土台、非構造データ対応", "カスタムメイドケア/予兆検知の基盤"],
            ["", "", ""],
            ["DDA取組み", "内容", "基盤への要件"],
            ["1.音声フォーマット化", "アセスメント音声の自動化", "音声取込みレーン(将来)"],
            ["2.ケアプラン自動化", "標準化+AI最適化", "セマンティックレイヤー必須"],
            ["3.介護S/W統一", "NDS CareBase移行(2028)", "スキーマ変更耐性"],
            ["4.カメラ・センサー", "最大活用", "構造化結果の取込み"],
            ["5.リアルデータ活用", "重度化予測等", "時点スナップショット保持"],
        ],
        col_widths=[2500000, 4500000, 4000000]
    )

    # --- Slide 3: 一気通貫 ---
    make_content_slide(prs,
        title="As-Is -> To-Be -> 4原則 -> RFP評価軸",
        message="全ての設計判断が一気通貫で繋がっている",
        body_items=[
            "As-Is: OverStack/高コスト -> G1 -> 原則3 ガバナンス -> 統合カタログ",
            "As-Is: 属人化/BBox化 -> G2 -> 原則4 ナレッジ蓄積 -> セルフサービス性",
            "As-Is: セルフBI不可 -> G2 -> 原則4 ナレッジ蓄積 -> 学習曲線",
            "As-Is: AI環境なし -> G3 -> 原則1 AI接続 -> SQL API",
            "As-Is: 非構造未対応 -> G3 -> 原則2 学習蓄積 -> MLflow統合",
            "As-Is: ナレッジ消失 -> G3 -> 原則2 学習蓄積 -> セマンティックモデル",
            "",
            "RFPの評価軸一つ一つが「なぜ必要か」をAs-Is課題まで遡って説明できる。",
        ]
    )

    # --- Slide 4: レイヤー+差分 (merged 7+8) ---
    make_content_slide(prs,
        title="4レイヤー分解 + Foundryとの差分",
        message="L0-L1はHD共通化、L2-L3はケア独自。Foundryの課題を削ぎ落とし意図的に変えた設計",
        table_data=[
            ["レイヤー", "HD基盤", "Foundryから引き継ぐ", "意図的に変える"],
            ["L0:インフラ", "乗る", "AWS/VPC/IdP共通", "--"],
            ["L1:ストレージ", "条件付き", "Spark互換(219本PySpark)", "オープンフォーマット(Delta)"],
            ["L2:データアーキ", "乗れない", "オントロジー->セマンティック", "独自言語->SQL/dbt標準"],
            ["L3:アプリ", "乗れない", "BI画面要件、KPI定義", "セルフBI化+MLflow統合"],
        ],
        col_widths=[1800000, 1500000, 3700000, 4000000]
    )

    # --- Slide 5: ガバナンス代替 ---
    make_content_slide(prs,
        title="ガバナンス機能の代替マッピング",
        message="Foundryの5機能中3機能は完全代替可能。残り2機能はCI/CDとOSSで補完",
        table_data=[
            ["Foundry機能", "新基盤での代替", "適合度", "備考"],
            ["Markings(タグ自動伝播)", "Tags(手動/API)", "△", "自動伝播なし、CI/CDで補完"],
            ["RLS(行レベルセキュリティ)", "行フィルタ+列マスキング", "◎", "Unity Catalog標準機能"],
            ["自動リネージ", "Unity Catalog Lineage", "○", "PySpark GA済"],
            ["監査ログ", "System Tables", "◎", "SQLで柔軟に分析可"],
            ["PII検出", "Presidio(OSS)/AWS Macie", "△", "別途導入必要"],
        ],
        col_widths=[2800000, 3000000, 1000000, 4200000]
    )

    # --- Slide 6: 課題x4原則 ---
    make_content_slide(prs,
        title="課題 x 4原則 対応マトリクス",
        message="As-Is課題がどの原則・設計判断で打ち取られるかを明示",
        table_data=[
            ["As-Is課題", "対応原則", "解決策", "RFP評価軸"],
            ["セルフBI不可", "4.ナレッジ蓄積", "SQL中心+セルフサービスBI", "学習曲線"],
            ["属人化/BBox化", "4.ナレッジ蓄積", "パイプラインGUI可視化+KPIコード管理", "運用ダッシュボード"],
            ["OverStack/高コスト", "4.ナレッジ蓄積", "必要機能に絞った従量課金構成", "コスト透明性"],
            ["SSOT不在", "1.AI接続", "セマンティックレイヤーでKPI一元管理", "セマンティックモデル"],
            ["AI実験環境なし", "2.学習蓄積", "MLflow+ModelRegistry+Git連携", "実験管理統合度"],
            ["ガバナンス都度構築", "3.ガバナンス", "統合カタログでアクセス制御一元管理", "統合カタログ"],
        ],
        col_widths=[2500000, 2000000, 3500000, 3000000]
    )

    # --- Slide 7: 非構造化データ ---
    make_content_slide(prs,
        title="非構造化データの取り扱い方針",
        message="基盤に入れるべきデータと外部に任せるデータの判断基準を定義",
        body_items=[
            "基盤に入れるべき: SOMPOケア固有ドメインデータ/横断結合データ/AI学習用時系列",
            "基盤に入れない: 外部完結する生データ(カメラ映像等)/品質未管理データ",
            "",
            "外部ソリューション連携: カメラ映像->外部分析->構造化結果のみ取込み",
            "",
            "Phase 1: 非構造データレーンの「設計」確保(実装は後回し)",
            "Phase 2: 介護記録テキストのPII検出・パターン分析",
            "Phase 3: 外部AI結果取込み、音声データ、マルチモーダル対応",
        ]
    )

    # --- Slide 8: 4原則+7観点 (merged 14+15) ---
    make_content_slide(prs,
        title="技術スタック採用4原則 + 7つの評価観点",
        message="4つの設計原則でRFPの評価軸を定め、7観点をSOMPOケア固有に具体化",
        table_data=[
            ["原則/観点", "内容", "SOMPOケア固有の確認事項"],
            ["原則1: AI接続", "SQL API/MCP/セマンティック", "AI-BI Genie, SQL精度3倍向上"],
            ["原則2: 学習蓄積", "MLflow/Git/評価基盤", "ベンダー変更後もナレッジ継続"],
            ["原則3: ガバナンス", "権限継承/監査/SSOT", "1,000施設RLS+要配慮個人情報"],
            ["原則4: ナレッジ蓄積", "段階的自走/運用透明性", "2名体制, SEなし運用"],
            ["", "", ""],
            ["1-2: 連携性+安全性", "今回選定", "ソース接続,RLS,監査ログ"],
            ["3-6: 進化/実効/運用/速度", "余地確認", "Delta Lake,219本移植,GUI完結"],
            ["7: 規模", "将来検討", "AIエージェント管理"],
        ],
        col_widths=[2800000, 3200000, 5000000]
    )

    # --- Slide 9: 運用+スコープ (merged 16+17) ---
    make_content_slide(prs,
        title="運用モデル + RFPスコープ",
        message="Phase 1-4で段階的に内製化。必須/加点/将来検討の3段階でRFPを明確化",
        body_items=[
            "運用ロードマップ:",
            (1, "Phase 1(~半年): ベンダー主導+社内学習"),
            (1, "Phase 2(~1年): 日常運用を社内移管(KPI変更/権限管理/監視)"),
            (1, "Phase 3(1年~): 社内主導, ベンダーは大規模改修のみ"),
            (1, "Phase 4(2年~): LLMツール活用で内製化拡大"),
            "",
            "RFPスコープ:",
            (1, "必須: ソース接続, CB刷新対応, 行/列RLS, 監査ログ, BI再現, 自走RM"),
            (1, "加点: セマンティック構築支援, 本番相当PoC, MLflow, Git, ナレッジ移転"),
            (1, "将来: HD基盤連携, 時間軸制御, AI経由制御, マルチモデル, エージェント管理"),
        ]
    )

    # --- Slide 10: RFI+推奨 (merged 19+20) ---
    make_content_slide(prs,
        title="RFI 5社評価 + 推奨技術スタック",
        message="Databricks + dbt + Tableauを推奨。PySpark移植コスト最小、2名程度で運用可能",
        table_data=[
            ["評価観点", "重み", "富士通", "デロイト", "日立", "PwC", "SF"],
            ["Foundry移行", "★5", "◎", "○", "△", "△", "x"],
            ["運用自走性", "★5", "○", "○", "◎", "△", "○"],
            ["統合カタログ", "★4", "◎", "◎", "○", "△", "△"],
            ["コスト構造", "★4", "○", "○", "○", "○", "△"],
            ["セルフBI", "★3", "○", "○", "◎", "△", "○"],
            ["SQL/セマンティック", "★3", "◎", "◎", "○", "△", "x"],
            ["MLflow", "★2", "◎", "○", "△", "△", "x"],
            ["総合", "--", "A", "A-", "B+", "B", "C"],
        ],
        col_widths=[2200000, 800000, 1200000, 1400000, 1200000, 1200000, 1200000]
    )

    # --- Slide 11: Next Steps ---
    make_content_slide(prs,
        title="Next Steps",
        message="やるべきことは明確。設計判断レベルの追加コストで対応可能",
        table_data=[
            ["#", "ステップ", "タイミング"],
            ["1", "4原則のRFP反映", "RFP策定時"],
            ["2", "HDとのレイヤー合意(L0-L1共通化範囲)", "HD基盤検討と並行"],
            ["3", "ソースシステムの繋ぎ込み調査", "要件定義フェーズ"],
            ["4", "Foundryオントロジー棚卸し", "RFP策定時"],
            ["5", "PoC候補の具体化(自然言語KPI照会)", "RFP策定と並行"],
            ["6", "マスタデータ名寄せ方針", "移行設計フェーズ"],
            ["7", "介護記録テキストデータの実態調査(340GB)", "PoC設計前"],
            ["8", "横断人材の確保方針", "Phase 1開始前"],
            ["9", "運用自走ロードマップの策定", "RFP策定時"],
        ],
        col_widths=[500000, 6500000, 4000000]
    )

    # Remove original template slides
    remove_template_slides(prs, original_count)
    prs.save(OUT_DETAILED)
    print(f"Saved detailed version: {OUT_DETAILED}")


if __name__ == "__main__":
    build_normal_version()
    build_detailed_version()
