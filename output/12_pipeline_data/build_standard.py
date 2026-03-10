#!/usr/bin/env python3
"""Build standard PPTX for doc 12: AIパイプライン × データナレッジ関係図.

~14 slides:
  0: Title
  1: INDEX
  2: Section — 全体マップ
  3: M1-M7 × K1-K8 consumption matrix
  4: Section — パイプライン別
  5: M1 input/knowledge/output
  6: M2 input/knowledge/output
  7: M4 AIチャット input/knowledge/output
  8: M5/M6/M7 combined
  9: Section — Data→Knowledge変換
 10: Data→K1-K8 conversion map
 11: Section — フィードバックループ
 12: Decision Trace → positive spiral
 13: Section — 即時回答設計
 14: L1/L2/L3 placement table
"""
import copy
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

TEMPLATE = "/workspace/project-rikyu-sales-proposals-poc/untracked/original_document/report/提案仮説構築AI_定例会1_20260219_v1.pptx"
OUTFILE = "/workspace/output/12_pipeline_data/12_pipeline_data_relation.pptx"

# ── ACES Colors ────────────────────────────────────────
C = {
    "dk_green":   RGBColor(0x00, 0x68, 0x43),
    "mint":       RGBColor(0xC2, 0xE7, 0xD9),
    "pale":       RGBColor(0xDF, 0xE6, 0xE0),
    "charcoal":   RGBColor(0x39, 0x39, 0x39),
    "tbl_text":   RGBColor(0x45, 0x45, 0x45),
    "accent":     RGBColor(0x52, 0x98, 0xBA),
    "white":      RGBColor(0xFF, 0xFF, 0xFF),
    "black":      RGBColor(0x00, 0x00, 0x00),
    "gray":       RGBColor(0x6E, 0x6F, 0x73),
    "lt_gray":    RGBColor(0xF2, 0xF2, 0xF2),
    "orange":     RGBColor(0xEC, 0x6C, 0x1F),
    "green":      RGBColor(0x00, 0x9A, 0x62),
    "high_bg":    RGBColor(0xE8, 0xF5, 0xE9),
    "mid_bar":    RGBColor(0x8B, 0x9D, 0xAF),
    "mid_bg":     RGBColor(0xEE, 0xF1, 0xF5),
    "low_bar":    RGBColor(0xB0, 0xB0, 0xB0),
    "low_bg":     RGBColor(0xF5, 0xF5, 0xF5),
    "star_gold":  RGBColor(0xD4, 0xA0, 0x17),
}
FL = "Century Gothic"
FJ = "Yu Gothic Medium"
FT = "Yu Gothic"

# ── Layout constants ───────────────────────────────────
X_FULL = 0.25
W_FULL = 9.50
Y0 = 1.55
Y_MAX = 7.10


# ╔══════════════════════════════════════════════════════╗
# ║  HELPERS                                             ║
# ╚══════════════════════════════════════════════════════╝

def _ea(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", FJ)
    rPr.append(ea)


def _run(para, text, sz=9, bold=False, color=None, font=None):
    r = para.add_run()
    r.text = text
    r.font.name = font or FL
    r.font.size = Pt(sz)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    _ea(r)
    return r


def _rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def _txbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def section_hdr(slide, text, x, y, w, sz=11):
    tf = _txbox(slide, x, y, w, 0.24)
    _run(tf.paragraphs[0], text, sz=sz, bold=True, color=C["dk_green"])
    _rect(slide, x, y + 0.24, w, 0.015, C["green"])
    return y + 0.30


def add_table(slide, data, x=0.25, y=1.10, w=9.50, col_widths=None,
              hdr_sz=8, cell_sz=7.5, rh=0.26):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(
        rows, cols, Inches(x), Inches(y), Inches(w), Inches(rh * rows)
    )
    tbl = ts.table
    tblPr = tbl._tbl.tblPr
    for attr in ("bandRow", "bandCol", "firstRow", "lastRow"):
        tblPr.attrib.pop(attr, None)
    for child in list(tblPr):
        if child.tag.endswith("}tblStyle"):
            tblPr.remove(child)
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for ri, rd in enumerate(data):
        tbl.rows[ri].height = Inches(rh)
        for ci, val in enumerate(rd):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            _run(
                p,
                str(val),
                sz=hdr_sz if ri == 0 else cell_sz,
                bold=(ri == 0),
                color=C["white"] if ri == 0 else C["tbl_text"],
            )
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["dk_green"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["white"]
    return tbl


def add_box(slide, x, y, w, h, fill, text, sz=9, bold=False, text_color=None, align="center"):
    """Add a colored box with centered text."""
    s = _rect(slide, x, y, w, h, fill)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    _run(p, text, sz=sz, bold=bold, color=text_color or C["charcoal"])
    s.text_frame.paragraphs[0].space_before = Pt(2)
    return s


def add_arrow(slide, x, y, w=0.3, h=0.2, direction="right"):
    """Add a small arrow shape."""
    if direction == "right":
        shape_type = MSO_SHAPE.RIGHT_ARROW
    else:
        shape_type = MSO_SHAPE.DOWN_ARROW
    s = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = C["dk_green"]
    s.line.fill.background()
    return s


def bullets(slide, items, x, y, w, sz=9):
    total_lines = sum(
        1 + str(item[0] if isinstance(item, tuple) else item).count("\n")
        for item in items
    )
    h = total_lines * sz * 1.5 / 72 + 0.15
    tf = _txbox(slide, x, y, w, max(h, 0.3))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        if isinstance(item, tuple):
            _run(p, item[0], sz=sz, color=C["charcoal"])
            if item[1]:
                _run(p, f"  {item[1]}", sz=max(sz - 1.5, 6), color=C["gray"])
        else:
            _run(p, str(item), sz=sz, color=C["charcoal"])
    return y + h


# ╔══════════════════════════════════════════════════════╗
# ║  SLIDE INFRASTRUCTURE                                ║
# ╚══════════════════════════════════════════════════════╝

def find_shapes(slide):
    """Return shapes sorted by Y-coordinate."""
    return sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))


def prep(slide):
    """Suppress layout bleed-through by clearing inherited text."""
    for sh in slide.shapes:
        if hasattr(sh, "text_frame"):
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.text = ""


def remove_idx147(slide):
    """OOXML cleanup: remove elements with idx=147."""
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for elem in slide._element.iter():
        if elem.get("idx") == "147":
            parent = elem.getparent()
            if parent is not None:
                parent.remove(elem)


def duplicate_slide(prs, src_idx):
    """Duplicate a slide from the template and return it."""
    template_slide = prs.slides[src_idx]
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    # Copy shapes from template
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)
    return new_slide


def make_content_slide(prs, title_text, subtitle_text=""):
    """Create a content slide from template slide 21 layout."""
    layout = prs.slide_layouts[25]  # Contents_Title&Message_Left
    slide = prs.slides.add_slide(layout)
    prep(slide)
    remove_idx147(slide)
    shapes = find_shapes(slide)
    # Set title in the first text shape found near top
    for sh in shapes:
        if hasattr(sh, "text_frame") and (sh.top or 0) < Inches(1.0):
            tf = sh.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            if tf.paragraphs:
                _run(tf.paragraphs[0], title_text, sz=14, bold=True, color=C["charcoal"], font=FT)
            break
    return slide


def make_section_slide(prs, title_text):
    """Create a section divider slide."""
    layout = prs.slide_layouts[4]  # Section Cover
    slide = prs.slides.add_slide(layout)
    prep(slide)
    remove_idx147(slide)
    shapes = find_shapes(slide)
    for sh in shapes:
        if hasattr(sh, "text_frame") and (sh.height or 0) > Inches(0.5):
            tf = sh.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            _run(tf.paragraphs[0], title_text, sz=24, bold=True, color=C["white"], font=FT)
            break
    return slide


def make_title_slide(prs, title_text, date_text):
    """Create a title slide."""
    layout = prs.slide_layouts[1]  # Cover+data
    slide = prs.slides.add_slide(layout)
    prep(slide)
    remove_idx147(slide)
    shapes = find_shapes(slide)
    title_set = False
    date_set = False
    for sh in shapes:
        if hasattr(sh, "text_frame"):
            tf = sh.text_frame
            if not title_set and (sh.top or 0) < Inches(4.0) and (sh.height or 0) > Inches(0.3):
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.text = ""
                _run(tf.paragraphs[0], title_text, sz=18, bold=True, color=C["charcoal"], font=FT)
                title_set = True
            elif not date_set and (sh.top or 0) >= Inches(4.0):
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.text = ""
                _run(tf.paragraphs[0], date_text, sz=10, color=C["gray"])
                date_set = True
    return slide


def make_index_slide(prs, items):
    """Create an INDEX slide."""
    layout = prs.slide_layouts[3]  # Index_1column
    slide = prs.slides.add_slide(layout)
    prep(slide)
    remove_idx147(slide)
    shapes = find_shapes(slide)
    # Set title
    for sh in shapes:
        if hasattr(sh, "text_frame"):
            tf = sh.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            _run(tf.paragraphs[0], "INDEX", sz=16, bold=True, color=C["charcoal"], font=FT)
            break
    # Add index items as text box
    y = 1.40
    for i, item in enumerate(items):
        tf = _txbox(slide, 0.50, y, 9.00, 0.30)
        p = tf.paragraphs[0]
        _run(p, f"{i + 1}. ", sz=11, bold=True, color=C["dk_green"])
        _run(p, item, sz=11, color=C["charcoal"])
        y += 0.36
    return slide


# ╔══════════════════════════════════════════════════════╗
# ║  SLIDE CONTENT BUILDERS                              ║
# ╚══════════════════════════════════════════════════════╝

def build_consumption_matrix(slide):
    """Slide 3: M1-M7 × K1-K8 consumption matrix."""
    section_hdr(slide, "パイプライン × Knowledge Master 消費マトリクス", X_FULL, Y0, W_FULL)

    data = [
        ["", "K1\n業界テンプ\nレート", "K2\n更新\nルール", "K3\n信頼度\nマトリクス",
         "K4\n商品マッチ\nング", "K5\n事例\nパターン", "K6\n商談進行\nパターン",
         "K7\n顧客KM", "K8\nタイミング\n知識"],
        ["M1 初期生成", "★★★", "─", "─", "─", "─", "─", "★", "─"],
        ["M2 面談更新", "★", "★★★", "★★★", "─", "─", "─", "★★", "─"],
        ["M3 ニュース", "★★", "★", "─", "─", "─", "─", "★", "★★"],
        ["M4 チャット", "★★", "★", "★", "★★★", "★★", "★★★", "★★", "★★"],
        ["M5 提案", "★★", "─", "─", "★★★", "★★★", "★★", "★", "★★"],
        ["M6 深掘り", "★★", "─", "─", "★", "★★★", "─", "★", "★"],
        ["M7 メモ", "★", "★", "★", "─", "─", "★★★", "★★", "★★★"],
    ]

    tbl = add_table(
        slide, data,
        x=0.15, y=Y0 + 0.35, w=9.70,
        col_widths=[1.10, 1.10, 0.95, 1.05, 1.10, 0.95, 1.10, 0.90, 1.05],
        hdr_sz=7, cell_sz=8, rh=0.42,
    )

    # Color-code cells with stars
    rows, cols = len(data), len(data[0])
    for ri in range(1, rows):
        for ci in range(1, cols):
            cell = tbl.cell(ri, ci)
            val = data[ri][ci]
            if "★★★" in val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["high_bg"]
            elif "★★" in val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["mid_bg"]
            elif "★" in val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["low_bg"]

    # Legend
    tf = _txbox(slide, X_FULL, 5.70, W_FULL, 0.30)
    p = tf.paragraphs[0]
    _run(p, "★★★", sz=9, bold=True, color=C["dk_green"])
    _run(p, "=主要消費  ", sz=8, color=C["gray"])
    _run(p, "★★", sz=9, bold=True, color=C["mid_bar"])
    _run(p, "=参照  ", sz=8, color=C["gray"])
    _run(p, "★", sz=9, bold=True, color=C["low_bar"])
    _run(p, "=補助参照  ", sz=8, color=C["gray"])
    _run(p, "─", sz=9, color=C["gray"])
    _run(p, "=不使用", sz=8, color=C["gray"])


def build_m1_pipeline(slide):
    """Slide 5: M1 アジェンダ初期生成 pipeline."""
    section_hdr(slide, "M1 アジェンダ初期生成: 入力Data × Knowledge × 出力", X_FULL, Y0, W_FULL)

    # Visual diagram: Input boxes → Knowledge → Output
    bx, by = 0.25, Y0 + 0.35
    # Input boxes
    inputs = ["企業基本情報\n(TSR/TDB)", "財務情報\n(EDINET)", "業界ニュース", "企業HP/IR"]
    for i, inp in enumerate(inputs):
        add_box(slide, bx, by + i * 0.55, 1.80, 0.45, C["pale"], inp, sz=8, text_color=C["charcoal"])

    # Arrow
    add_arrow(slide, 2.30, by + 0.80, w=0.40, h=0.25)

    # Knowledge boxes
    add_box(slide, 2.95, by + 0.30, 2.20, 0.50, C["mint"], "K1 業界テンプレート", sz=9, bold=True, text_color=C["dk_green"])
    add_box(slide, 2.95, by + 0.95, 2.20, 0.50, C["mint"], "K7 顧客KM (既存時)", sz=9, bold=True, text_color=C["dk_green"])

    # Arrow
    add_arrow(slide, 5.40, by + 0.80, w=0.40, h=0.25)

    # Output boxes
    add_box(slide, 6.05, by + 0.30, 3.50, 0.50, C["high_bg"], "経営アジェンダシート\n(項目×優先度×確信度×時期)", sz=8, text_color=C["dk_green"])
    add_box(slide, 6.05, by + 0.95, 3.50, 0.50, C["high_bg"], "新規アジェンダ候補リスト", sz=8, text_color=C["dk_green"])

    # Detail table
    data = [
        ["入力Data", "対象", "取得方法", "消費Knowledge", "処理方式"],
        ["企業基本情報", "顧客(法人)", "自動(API)", "K1: 業種コード→テンプレート選択", "ルールエンジン(即時)"],
        ["財務情報(5年)", "顧客(法人)", "自動(EDINET)", "K1: 財務→カスタマイズ", "LLM"],
        ["業界ニュース", "業界/業種", "自動(API)", "—", "LLM(補完)"],
        ["企業HP/IR", "顧客(法人)", "自動(スクレイピング)", "—", "LLM(補完)"],
    ]
    add_table(
        slide, data,
        x=0.15, y=4.20, w=9.70,
        col_widths=[1.30, 1.00, 1.50, 3.10, 2.80],
        hdr_sz=7.5, cell_sz=7, rh=0.28,
    )

    # Note
    tf = _txbox(slide, X_FULL, 5.70, W_FULL, 0.35)
    p = tf.paragraphs[0]
    _run(p, "即時回答ポイント: ", sz=8, bold=True, color=C["dk_green"])
    _run(p, "K1業界テンプレートの読み込みは<1秒。テンプレートベースの初期アジェンダは即座に生成し、LLMカスタマイズは非同期で追加。", sz=8, color=C["gray"])


def build_m2_pipeline(slide):
    """Slide 6: M2 面談→アジェンダ更新提案."""
    section_hdr(slide, "M2 面談→アジェンダ更新提案: 入力Data × Knowledge × 出力", X_FULL, Y0, W_FULL)

    # Visual diagram
    bx, by = 0.25, Y0 + 0.35
    inputs = ["面談記録", "既存アジェンダ", "KP情報\n(役職/決裁権)"]
    for i, inp in enumerate(inputs):
        add_box(slide, bx, by + i * 0.55, 1.80, 0.45, C["pale"], inp, sz=8, text_color=C["charcoal"])

    add_arrow(slide, 2.30, by + 0.55, w=0.40, h=0.25)

    # Knowledge
    kbs = ["K2 更新ルール", "K3 信頼度マトリクス", "K7 顧客KM"]
    for i, kb in enumerate(kbs):
        add_box(slide, 2.95, by + i * 0.55, 2.20, 0.45, C["mint"], kb, sz=9, bold=True, text_color=C["dk_green"])

    add_arrow(slide, 5.40, by + 0.55, w=0.40, h=0.25)

    # Output
    add_box(slide, 6.05, by + 0.10, 3.50, 0.55, C["high_bg"], "更新提案\n(変更箇所+根拠+信頼度)", sz=8, text_color=C["dk_green"])
    add_box(slide, 6.05, by + 0.80, 3.50, 0.55, C["high_bg"], "Decision Trace記録", sz=8, text_color=C["dk_green"])

    # Table
    data = [
        ["入力Data", "消費Knowledge", "処理方式", "応答速度"],
        ["面談記録(自由テキスト)", "K3: 発言→信頼度判定", "ルールエンジン(マトリクス照合)", "<1秒"],
        ["面談記録(構造化済み)", "K2: 発言→更新ルール適用", "ルールエンジン + LLM", "3-10秒"],
        ["KP情報", "K3: 発言者の役職→信頼度", "ルールエンジン(即時)", "<1秒"],
    ]
    add_table(
        slide, data,
        x=0.15, y=4.10, w=9.70,
        col_widths=[2.00, 2.60, 2.80, 1.00],
        hdr_sz=7.5, cell_sz=7, rh=0.30,
    )

    tf = _txbox(slide, X_FULL, 5.40, W_FULL, 0.35)
    p = tf.paragraphs[0]
    _run(p, "即時回答ポイント: ", sz=8, bold=True, color=C["dk_green"])
    _run(p, "信頼度判定はマトリクス照合で即時。更新ルールの確定パターンもルールエンジンで即時。LLMは曖昧な文脈解釈のみ。", sz=8, color=C["gray"])


def build_m4_pipeline(slide):
    """Slide 7: M4 AIチャット pipeline."""
    section_hdr(slide, "M4 AIチャット（折衝アシスタント）: 入力Data × Knowledge × 出力", X_FULL, Y0, W_FULL)

    # Visual diagram
    bx, by = 0.25, Y0 + 0.35
    inputs = ["RMの質問", "顧客情報+\nアジェンダ", "面談履歴", "KPマップ"]
    for i, inp in enumerate(inputs):
        add_box(slide, bx, by + i * 0.48, 1.70, 0.40, C["pale"], inp, sz=7.5, text_color=C["charcoal"])

    add_arrow(slide, 2.20, by + 0.55, w=0.35, h=0.20)

    add_box(slide, 2.80, by + 0.30, 2.40, 0.80, C["mint"], "K1-K8全Knowledge\n(質問種別に応じて\n動的に選択)", sz=8, bold=True, text_color=C["dk_green"])

    add_arrow(slide, 5.45, by + 0.55, w=0.35, h=0.20)

    add_box(slide, 6.05, by + 0.10, 3.50, 0.50, C["high_bg"], "チャット応答\n(ストリーミング)", sz=8, text_color=C["dk_green"])
    add_box(slide, 6.05, by + 0.75, 3.50, 0.50, C["high_bg"], "FB記録\n(暗黙知キャプチャ)", sz=8, text_color=C["dk_green"])

    # Question patterns table
    data = [
        ["質問パターン", "消費Knowledge", "処理方式", "応答速度"],
        ["「何を提案すべき？」", "K4商品マッチング + K1業界", "Knowledge参照 → LLM生成", "初回<0.5秒"],
        ["「どう進めるべき？」", "K6商談進行 + K8タイミング", "Knowledge参照 → LLM生成", "初回<0.5秒"],
        ["「この業界の動き？」", "K1業界 + ニュースData", "Knowledge参照 + Web補完", "2-5秒"],
        ["「この出力おかしくない？」", "K2更新ルール + K3信頼度", "ルール検証 → LLM説明", "1-3秒"],
    ]
    add_table(
        slide, data,
        x=0.15, y=4.30, w=9.70,
        col_widths=[2.20, 2.60, 2.50, 1.10],
        hdr_sz=7.5, cell_sz=7, rh=0.28,
    )

    tf = _txbox(slide, X_FULL, 5.80, W_FULL, 0.30)
    p = tf.paragraphs[0]
    _run(p, "即時回答ポイント: ", sz=8, bold=True, color=C["dk_green"])
    _run(p, "Knowledge Masterを事前にコンテキストに注入済み。RMの質問に対して即座にストリーミング開始。", sz=8, color=C["gray"])


def build_m567_combined(slide):
    """Slide 8: M5/M6/M7 combined table."""
    section_hdr(slide, "M5 提案ストーリー / M6 深掘りペーパー / M7 メモ", X_FULL, Y0, W_FULL)

    data = [
        ["パイプライン", "主要入力Data", "主要Knowledge", "即時対応", "LLM依存部分"],
        ["M5 提案\nストーリー", "アジェンダ\n+ 財務 + KP", "K4商品 + K5事例\n+ K6商談", "商品候補リスト\n(K4即時)", "ストーリー構成\n・表現"],
        ["M6 深掘り\nペーパー", "テーマ\n+ 業界Data", "K1業界\n+ K5事例", "業界パターン\n(K1即時)", "分析・考察の\n生成"],
        ["M7 メモ", "アジェンダ\n+ 面談履歴", "K6商談進行\n+ K8タイミング", "次アクション候補\n(K6即時)", "メモ文面の\n生成"],
    ]

    add_table(
        slide, data,
        x=0.15, y=Y0 + 0.35, w=9.70,
        col_widths=[1.30, 1.80, 2.20, 2.00, 2.00],
        hdr_sz=8, cell_sz=7.5, rh=0.60,
    )

    # Explanation bullets
    y = Y0 + 0.35 + 0.60 * 4 + 0.20
    bullets(slide, [
        "M5: 商品マッチング(K4)はルールエンジンで即時候補リスト生成 → ストーリー化はLLMで非同期",
        "M6: 業界テンプレート(K1)から深掘りの軸を即時提示 → 分析はLLMで生成",
        "M7: 商談進行パターン(K6)から次アクションを即時提案 → メモの文面はLLM",
    ], X_FULL + 0.10, y, W_FULL - 0.20, sz=9)


def build_data_knowledge_map(slide):
    """Slide 10: Data → Knowledge Master conversion map."""
    section_hdr(slide, "Data層 → Knowledge Master 変換マップ", X_FULL, Y0, W_FULL)

    # Visual layout: Data boxes on left → arrows → Knowledge boxes on right
    left_x = 0.20
    left_w = 3.80
    right_x = 5.60
    right_w = 4.20
    arrow_x = 4.20

    conversions = [
        ("業界/業種:\n市場規模・競合・規制", "K1 業界テンプレート\n(業種コード→共通課題パターン)"),
        ("渉外担当者:\n面談記録 + 承認ログ", "K2 更新ルール\n(発言パターン→更新アクション)"),
        ("顧客法人:\n役員・KP + 面談記録", "K3 信頼度マトリクス\n(発言者×具体性→信頼度)"),
        ("自社りそな:\n商品マスタ + 適合条件", "K4 商品マッチング\n(課題タグ→商品候補リスト)"),
        ("渉外担当者:\n提案履歴 + 成約/失注", "K5 事例パターン\n(類似事例→成功/失敗パターン)"),
        ("渉外担当者:\n面談記録 + 活動実績", "K6 商談進行パターン\n(状況→推奨アクション)"),
        ("全Data × 顧客単位で集約", "K7 顧客Knowledge Master\n(顧客ごとの統合コンテキスト)"),
        ("経済/世界情勢 + 業界/業種", "K8 タイミング知識\n(経済環境→提案タイミング)"),
    ]

    by = Y0 + 0.35
    rh = 0.52
    for i, (data_text, knowledge_text) in enumerate(conversions):
        y = by + i * (rh + 0.06)
        # Data box (left)
        add_box(slide, left_x, y, left_w, rh, C["pale"], data_text, sz=7, text_color=C["charcoal"])
        # Arrow
        add_arrow(slide, arrow_x, y + rh / 2 - 0.10, w=1.10, h=0.20)
        # Knowledge box (right)
        add_box(slide, right_x, y, right_w, rh, C["mint"], knowledge_text, sz=7, bold=True, text_color=C["dk_green"])


def build_feedback_loop(slide):
    """Slide 12: Feedback loop diagram."""
    section_hdr(slide, "暗黙知のフィードバックループ", X_FULL, Y0, W_FULL)

    # Build a vertical flow diagram
    cx = 2.00  # center column x
    cw = 6.00  # center column width
    by = Y0 + 0.40

    # Step 1: RM操作
    add_box(slide, cx, by, cw, 0.45, C["pale"], "RMの操作（承認 / 修正 / 棄却）", sz=10, bold=True, text_color=C["charcoal"])
    by += 0.50
    add_arrow(slide, cx + cw / 2 - 0.15, by, w=0.30, h=0.25, direction="down")
    by += 0.30

    # Step 2: Decision Trace
    add_box(slide, cx, by, cw, 0.45, C["mint"], "Decision Trace記録", sz=10, bold=True, text_color=C["dk_green"])
    by += 0.55

    # Two paths
    # Path A
    path_a_x = 0.50
    path_a_w = 4.20
    add_box(slide, path_a_x, by, path_a_w, 0.40, C["high_bg"], "パスA: Data表に追加（即時）", sz=9, bold=True, text_color=C["dk_green"])

    # Path B
    path_b_x = 5.30
    path_b_w = 4.50
    add_box(slide, path_b_x, by, path_b_w, 0.40, C["mid_bg"], "パスB: 3回同一パターン → Knowledge Master更新（昇格）", sz=8, bold=True, text_color=C["accent"])

    by += 0.55
    add_arrow(slide, cx + cw / 2 - 0.15, by, w=0.30, h=0.25, direction="down")
    by += 0.30

    # Step 3: AI精度向上
    add_box(slide, cx, by, cw, 0.40, C["dk_green"], "AI出力精度が向上", sz=10, bold=True, text_color=C["white"])
    by += 0.50
    add_arrow(slide, cx + cw / 2 - 0.15, by, w=0.30, h=0.25, direction="down")
    by += 0.30

    # Step 4: Positive spiral
    add_box(slide, cx, by, cw, 0.40, C["dk_green"], "RMがもっと使う → さらにData蓄積 → 正のスパイラル", sz=10, bold=True, text_color=C["white"])

    # Note about agent harness
    by += 0.55
    tf = _txbox(slide, X_FULL, by, W_FULL, 0.40)
    p = tf.paragraphs[0]
    _run(p, "エージェントハーネスとの対応: ", sz=8, bold=True, color=C["dk_green"])
    _run(p, "Decision Traceは、agentic harnessにおけるPostToolUse Hookと同じ位置づけ。RM操作（=ツール実行）の後に自動で発火し、暗黙知を確定的にキャプチャする。", sz=8, color=C["gray"])


def build_l1l2l3_table(slide):
    """Slide 14: L1/L2/L3 placement table."""
    section_hdr(slide, "即時回答の設計: 何をL1に置くか", X_FULL, Y0, W_FULL)

    data = [
        ["情報", "階層", "アクセス速度", "理由"],
        ["K1 業界テンプレート", "L1", "<1秒", "M1/M4で頻繁に参照。業種コードで即引き"],
        ["K2 更新ルール", "L1", "<1秒", "M2で毎回使用。YAML→即判定"],
        ["K3 信頼度マトリクス", "L1", "<1秒", "M2で毎回使用。テーブル照合のみ"],
        ["K4 商品マッチング", "L1", "<1秒", "M4/M5で頻繁に参照。課題タグで即引き"],
        ["K6 商談進行パターン", "L1", "<1秒", "M4/M7で頻繁に参照"],
        ["K8 タイミング知識", "L1", "<1秒", "M3/M4/M7で参照。カレンダー＋ルール"],
        ["顧客個別の財務Data", "L2", "1-5秒", "顧客アクセス時にロード"],
        ["面談履歴・提案履歴", "L2", "1-5秒", "セッション開始時にロード"],
        ["最新ニュース・HP", "L3", "5-15秒", "必要時のみWeb取得"],
    ]

    tbl = add_table(
        slide, data,
        x=0.15, y=Y0 + 0.35, w=9.70,
        col_widths=[2.20, 0.60, 1.10, 5.40],
        hdr_sz=8, cell_sz=7.5, rh=0.32,
    )

    # Color code by level
    for ri in range(1, len(data)):
        level = data[ri][1]
        cell_level = tbl.cell(ri, 1)
        if level == "L1":
            cell_level.fill.solid()
            cell_level.fill.fore_color.rgb = C["high_bg"]
        elif level == "L2":
            cell_level.fill.solid()
            cell_level.fill.fore_color.rgb = C["mid_bg"]
        elif level == "L3":
            cell_level.fill.solid()
            cell_level.fill.fore_color.rgb = C["low_bg"]

    # Legend
    y = Y0 + 0.35 + 0.32 * len(data) + 0.15
    tf = _txbox(slide, X_FULL, y, W_FULL, 0.30)
    p = tf.paragraphs[0]
    _run(p, "L1", sz=9, bold=True, color=C["dk_green"])
    _run(p, "=メモリ常駐(即時)  ", sz=8, color=C["gray"])
    _run(p, "L2", sz=9, bold=True, color=C["accent"])
    _run(p, "=セッション時ロード  ", sz=8, color=C["gray"])
    _run(p, "L3", sz=9, bold=True, color=C["low_bar"])
    _run(p, "=必要時のみ取得", sz=8, color=C["gray"])


# ╔══════════════════════════════════════════════════════╗
# ║  MAIN                                                ║
# ╚══════════════════════════════════════════════════════╝

def main():
    prs = Presentation(TEMPLATE)

    # Remove all existing slides (build from scratch using layouts)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Slide 0: Title
    make_title_slide(prs, "AIパイプライン × データナレッジ関係図\nrikyu ソリューションセールス伴奏AI PoC", "2026-03-10  |  POC最終成果物")

    # Slide 1: INDEX
    make_index_slide(prs, [
        "全体マップ: パイプライン × Knowledge Master",
        "パイプライン別: 入力Data × Knowledge × 出力",
        "Data→Knowledge変換マップ",
        "フィードバックループ",
        "即時回答設計: L1/L2/L3配置",
    ])

    # Slide 2: Section — 全体マップ
    make_section_slide(prs, "1. 全体マップ")

    # Slide 3: Consumption matrix
    s3 = make_content_slide(prs, "パイプライン × Knowledge Master 消費マトリクス")
    build_consumption_matrix(s3)

    # Slide 4: Section — パイプライン別
    make_section_slide(prs, "2. パイプライン別")

    # Slide 5: M1
    s5 = make_content_slide(prs, "M1 アジェンダ初期生成")
    build_m1_pipeline(s5)

    # Slide 6: M2
    s6 = make_content_slide(prs, "M2 面談→アジェンダ更新提案")
    build_m2_pipeline(s6)

    # Slide 7: M4
    s7 = make_content_slide(prs, "M4 AIチャット（折衝アシスタント）")
    build_m4_pipeline(s7)

    # Slide 8: M5/M6/M7
    s8 = make_content_slide(prs, "M5/M6/M7 提案・深掘り・メモ")
    build_m567_combined(s8)

    # Slide 9: Section — Data→Knowledge変換
    make_section_slide(prs, "3. Data→Knowledge変換")

    # Slide 10: Conversion map
    s10 = make_content_slide(prs, "Data層 → Knowledge Master 変換マップ")
    build_data_knowledge_map(s10)

    # Slide 11: Section — フィードバックループ
    make_section_slide(prs, "4. フィードバックループ")

    # Slide 12: Feedback loop
    s12 = make_content_slide(prs, "暗黙知のフィードバックループ")
    build_feedback_loop(s12)

    # Slide 13: Section — 即時回答設計
    make_section_slide(prs, "5. 即時回答設計")

    # Slide 14: L1/L2/L3 table
    s14 = make_content_slide(prs, "即時回答の設計: 何をL1に置くか")
    build_l1l2l3_table(s14)

    prs.save(OUTFILE)
    print(f"Saved standard PPTX ({len(prs.slides)} slides) to {OUTFILE}")


if __name__ == "__main__":
    main()
