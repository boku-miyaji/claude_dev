#!/usr/bin/env python3
"""Build detailed PPTX for doc 12: AIパイプライン × データナレッジ関係図.

~8 slides (dense, handout-style):
  0: Title
  1: INDEX
  2: 全体消費マトリクス（詳細注釈付き）
  3: M1/M2 パイプライン詳細
  4: M4 パイプライン詳細 + 質問パターン全量
  5: M5/M6/M7 + Data→Knowledge変換マップ
  6: フィードバックループ詳細 + Decision Traceフロー
  7: L1/L2/L3配置 + 設計根拠
"""
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

TEMPLATE = "/workspace/project-rikyu-sales-proposals-poc/untracked/original_document/report/提案仮説構築AI_定例会1_20260219_v1.pptx"
OUTFILE = "/workspace/output/12_pipeline_data/12_pipeline_data_relation_detailed.pptx"

# ── ACES Colors ────────────────────────────────────────
C = {
    "dk_green":   RGBColor(0x00, 0x68, 0x43),
    "mint":       RGBColor(0xC2, 0xE7, 0xD9),
    "pale":       RGBColor(0xDF, 0xE6, 0xE0),
    "charcoal":   RGBColor(0x39, 0x39, 0x39),
    "tbl_text":   RGBColor(0x45, 0x45, 0x45),
    "accent":     RGBColor(0x52, 0x98, 0xBA),
    "white":      RGBColor(0xFF, 0xFF, 0xFF),
    "black":      RGBColor(0x00, 0x00, 0x00),
    "gray":       RGBColor(0x6E, 0x6F, 0x73),
    "lt_gray":    RGBColor(0xF2, 0xF2, 0xF2),
    "orange":     RGBColor(0xEC, 0x6C, 0x1F),
    "green":      RGBColor(0x00, 0x9A, 0x62),
    "high_bg":    RGBColor(0xE8, 0xF5, 0xE9),
    "mid_bar":    RGBColor(0x8B, 0x9D, 0xAF),
    "mid_bg":     RGBColor(0xEE, 0xF1, 0xF5),
    "low_bar":    RGBColor(0xB0, 0xB0, 0xB0),
    "low_bg":     RGBColor(0xF5, 0xF5, 0xF5),
}
FL = "Century Gothic"
FJ = "Yu Gothic Medium"
FT = "Yu Gothic"

X_FULL = 0.17
W_FULL = 9.67
Y0 = 1.55


# ╔══════════════════════════════════════════════════════╗
# ║  HELPERS                                             ║
# ╚══════════════════════════════════════════════════════╝

def _ea(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", FJ)
    rPr.append(ea)


def _run(para, text, sz=9, bold=False, color=None, font=None):
    r = para.add_run()
    r.text = text
    r.font.name = font or FL
    r.font.size = Pt(sz)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    _ea(r)
    return r


def _rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def _txbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def section_hdr(slide, text, x, y, w, sz=9):
    tf = _txbox(slide, x, y, w, 0.22)
    _run(tf.paragraphs[0], text, sz=sz, bold=True, color=C["dk_green"])
    _rect(slide, x, y + 0.22, w, 0.015, C["green"])
    return y + 0.28


def add_table(slide, data, x=0.17, y=1.10, w=9.67, col_widths=None,
              hdr_sz=7, cell_sz=6.5, rh=0.22):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(
        rows, cols, Inches(x), Inches(y), Inches(w), Inches(rh * rows)
    )
    tbl = ts.table
    tblPr = tbl._tbl.tblPr
    for attr in ("bandRow", "bandCol", "firstRow", "lastRow"):
        tblPr.attrib.pop(attr, None)
    for child in list(tblPr):
        if child.tag.endswith("}tblStyle"):
            tblPr.remove(child)
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for ri, rd in enumerate(data):
        tbl.rows[ri].height = Inches(rh)
        for ci, val in enumerate(rd):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            _run(
                p, str(val),
                sz=hdr_sz if ri == 0 else cell_sz,
                bold=(ri == 0),
                color=C["white"] if ri == 0 else C["tbl_text"],
            )
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["dk_green"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["white"]
    return tbl


def add_box(slide, x, y, w, h, fill, text, sz=7, bold=False, text_color=None, align="center"):
    s = _rect(slide, x, y, w, h, fill)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    _run(p, text, sz=sz, bold=bold, color=text_color or C["charcoal"])
    return s


def add_arrow(slide, x, y, w=0.3, h=0.2, direction="right"):
    shape_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    s = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = C["dk_green"]
    s.line.fill.background()
    return s


def bullets(slide, items, x, y, w, sz=7):
    total_lines = sum(
        1 + str(item[0] if isinstance(item, tuple) else item).count("\n")
        for item in items
    )
    h = total_lines * sz * 1.5 / 72 + 0.12
    tf = _txbox(slide, x, y, w, max(h, 0.3))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        if isinstance(item, tuple):
            _run(p, item[0], sz=sz, color=C["charcoal"])
            if item[1]:
                _run(p, f"  {item[1]}", sz=max(sz - 1.5, 5.5), color=C["gray"])
        else:
            _run(p, str(item), sz=sz, color=C["charcoal"])
    return y + h


def find_shapes(slide):
    return sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))


def prep(slide):
    for sh in slide.shapes:
        if hasattr(sh, "text_frame"):
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.text = ""


def remove_idx147(slide):
    for elem in slide._element.iter():
        if elem.get("idx") == "147":
            parent = elem.getparent()
            if parent is not None:
                parent.remove(elem)


def make_content_slide(prs, title_text):
    layout = prs.slide_layouts[25]  # Contents_Title&Message_Left
    slide = prs.slides.add_slide(layout)
    prep(slide)
    remove_idx147(slide)
    shapes = find_shapes(slide)
    for sh in shapes:
        if hasattr(sh, "text_frame") and (sh.top or 0) < Inches(1.0):
            tf = sh.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            if tf.paragraphs:
                _run(tf.paragraphs[0], title_text, sz=13, bold=True, color=C["charcoal"], font=FT)
            break
    return slide


def make_title_slide(prs, title_text, date_text):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    prep(slide)
    remove_idx147(slide)
    shapes = find_shapes(slide)
    title_set = False
    date_set = False
    for sh in shapes:
        if hasattr(sh, "text_frame"):
            tf = sh.text_frame
            if not title_set and (sh.top or 0) < Inches(4.0) and (sh.height or 0) > Inches(0.3):
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.text = ""
                _run(tf.paragraphs[0], title_text, sz=16, bold=True, color=C["charcoal"], font=FT)
                title_set = True
            elif not date_set and (sh.top or 0) >= Inches(4.0):
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.text = ""
                _run(tf.paragraphs[0], date_text, sz=10, color=C["gray"])
                date_set = True
    return slide


def make_index_slide(prs, items):
    layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(layout)
    prep(slide)
    remove_idx147(slide)
    shapes = find_shapes(slide)
    for sh in shapes:
        if hasattr(sh, "text_frame"):
            tf = sh.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            _run(tf.paragraphs[0], "INDEX — 詳細版", sz=14, bold=True, color=C["charcoal"], font=FT)
            break
    y = 1.40
    for i, item in enumerate(items):
        tf = _txbox(slide, 0.50, y, 9.00, 0.28)
        p = tf.paragraphs[0]
        _run(p, f"{i + 1}. ", sz=10, bold=True, color=C["dk_green"])
        _run(p, item, sz=10, color=C["charcoal"])
        y += 0.33
    return slide


# ╔══════════════════════════════════════════════════════╗
# ║  SLIDE CONTENT BUILDERS                              ║
# ╚══════════════════════════════════════════════════════╝

def build_slide_2(slide):
    """全体消費マトリクス + 注釈."""
    y = section_hdr(slide, "パイプライン × Knowledge Master 消費マトリクス（詳細）", X_FULL, Y0, W_FULL)

    data = [
        ["", "K1\n業界テンプ\nレート", "K2\n更新\nルール", "K3\n信頼度\nマトリクス",
         "K4\n商品\nマッチング", "K5\n事例\nパターン", "K6\n商談\n進行", "K7\n顧客KM", "K8\nタイミング"],
        ["M1 初期生成", "★★★", "─", "─", "─", "─", "─", "★", "─"],
        ["M2 面談更新", "★", "★★★", "★★★", "─", "─", "─", "★★", "─"],
        ["M3 ニュース", "★★", "★", "─", "─", "─", "─", "★", "★★"],
        ["M4 チャット", "★★", "★", "★", "★★★", "★★", "★★★", "★★", "★★"],
        ["M5 提案", "★★", "─", "─", "★★★", "★★★", "★★", "★", "★★"],
        ["M6 深掘り", "★★", "─", "─", "★", "★★★", "─", "★", "★"],
        ["M7 メモ", "★", "★", "★", "─", "─", "★★★", "★★", "★★★"],
    ]

    tbl = add_table(
        slide, data,
        x=X_FULL, y=y, w=W_FULL,
        col_widths=[1.00, 1.05, 0.90, 1.05, 1.00, 0.90, 0.90, 0.80, 0.90],
        hdr_sz=6, cell_sz=7, rh=0.36,
    )

    # Color cells
    for ri in range(1, len(data)):
        for ci in range(1, len(data[0])):
            cell = tbl.cell(ri, ci)
            val = data[ri][ci]
            if "★★★" in val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["high_bg"]
            elif "★★" in val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["mid_bg"]
            elif "★" in val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["low_bg"]

    y2 = y + 0.36 * len(data) + 0.08
    y2 = section_hdr(slide, "消費パターンの解説", X_FULL, y2, W_FULL)
    bullets(slide, [
        "M4(AIチャット)が最もKnowledge消費が多い — 質問種別に応じてK1-K8を動的に選択",
        "M1(初期生成)はK1のみに依存 — テンプレートベースで即時生成可能",
        "K5(事例パターン)はM4/M5/M6で共有消費 — 類似事例マッチングが鍵",
        "K7(顧客KM)は全パイプラインで補助参照 — 顧客コンテキストの統合ハブ",
        "★★★=主要消費(そのパイプラインの核)  ★★=参照(品質向上)  ★=補助参照(あれば参照)  ─=不使用",
    ], X_FULL + 0.05, y2, W_FULL - 0.10, sz=6.5)


def build_slide_3(slide):
    """M1/M2 パイプライン詳細."""
    y = section_hdr(slide, "M1 アジェンダ初期生成 — 入力/Knowledge/出力 詳細", X_FULL, Y0, W_FULL)

    data_m1 = [
        ["入力Data", "対象", "取得方法", "消費Knowledge", "処理方式"],
        ["企業基本情報", "顧客(法人)", "自動(API)", "K1: 業種コード→テンプレート選択", "ルールエンジン(即時)"],
        ["財務情報(5年)", "顧客(法人)", "自動(EDINET)", "K1: 財務→カスタマイズ", "LLM"],
        ["業界ニュース", "業界/業種", "自動(API)", "—", "LLM(補完)"],
        ["企業HP/IR", "顧客(法人)", "自動(スクレイピング)", "—", "LLM(補完)"],
    ]
    add_table(
        slide, data_m1,
        x=X_FULL, y=y, w=W_FULL,
        col_widths=[1.20, 0.90, 1.40, 3.07, 2.70],
        hdr_sz=6.5, cell_sz=6, rh=0.22,
    )

    # M1 note
    y2 = y + 0.22 * len(data_m1) + 0.05
    tf = _txbox(slide, X_FULL, y2, W_FULL, 0.18)
    p = tf.paragraphs[0]
    _run(p, "M1即時回答: ", sz=6.5, bold=True, color=C["dk_green"])
    _run(p, "K1業界テンプレートの読み込みは<1秒。テンプレートベースの初期アジェンダは即座に生成し、LLMカスタマイズは非同期で追加。", sz=6, color=C["gray"])

    y3 = y2 + 0.22
    y3 = section_hdr(slide, "M2 面談→アジェンダ更新提案 — 入力/Knowledge/出力 詳細", X_FULL, y3, W_FULL)

    data_m2 = [
        ["入力Data", "消費Knowledge", "処理方式", "応答速度", "備考"],
        ["面談記録(自由テキスト)", "K3: 発言→信頼度判定", "ルールエンジン(マトリクス照合)", "<1秒", "発言者役職×具体性で即判定"],
        ["面談記録(構造化済み)", "K2: 発言→更新ルール適用", "ルールエンジン + LLM", "3-10秒", "確定パターンはルールで即時"],
        ["KP情報", "K3: 発言者役職→信頼度", "ルールエンジン(即時)", "<1秒", "役職テーブル照合のみ"],
    ]
    add_table(
        slide, data_m2,
        x=X_FULL, y=y3, w=W_FULL,
        col_widths=[1.60, 2.10, 2.30, 0.65, 2.62],
        hdr_sz=6.5, cell_sz=6, rh=0.22,
    )

    y4 = y3 + 0.22 * len(data_m2) + 0.05
    tf = _txbox(slide, X_FULL, y4, W_FULL, 0.18)
    p = tf.paragraphs[0]
    _run(p, "M2即時回答: ", sz=6.5, bold=True, color=C["dk_green"])
    _run(p, "信頼度判定はマトリクス照合で即時。更新ルールの確定パターンもルールエンジンで即時。LLMは曖昧な文脈解釈のみ。", sz=6, color=C["gray"])

    y5 = y4 + 0.22
    y5 = section_hdr(slide, "M1/M2 出力比較", X_FULL, y5, W_FULL)
    data_out = [
        ["パイプライン", "主要出力", "即時生成部分", "LLM生成部分", "Decision Trace"],
        ["M1 初期生成", "経営アジェンダシート\n+ 新規候補リスト", "テンプレート選択\n+ 業種別共通課題", "財務データ解釈\n+ IR情報統合", "—(初回生成)"],
        ["M2 面談更新", "更新提案\n+ Decision Trace", "信頼度判定\n+ 確定パターン更新", "曖昧文脈の解釈\n+ 更新根拠生成", "RMの承認/修正/棄却を\n自動記録"],
    ]
    add_table(
        slide, data_out,
        x=X_FULL, y=y5, w=W_FULL,
        col_widths=[1.00, 1.80, 1.80, 1.80, 2.87],
        hdr_sz=6.5, cell_sz=6, rh=0.36,
    )


def build_slide_4(slide):
    """M4 AIチャット詳細 + 全質問パターン."""
    y = section_hdr(slide, "M4 AIチャット（折衝アシスタント）— 全質問パターン詳細", X_FULL, Y0, W_FULL)

    data = [
        ["質問パターン", "消費Knowledge", "処理方式", "応答速度", "Knowledge選択ロジック"],
        ["「何を提案すべき？」", "K4商品マッチング\n+ K1業界", "Knowledge参照\n→ LLM生成", "初回<0.5秒", "課題タグ→K4で商品候補\n→K1で業界文脈付与"],
        ["「どう進めるべき？」", "K6商談進行\n+ K8タイミング", "Knowledge参照\n→ LLM生成", "初回<0.5秒", "商談ステージ→K6で推奨\n→K8でタイミング最適化"],
        ["「この業界の動き？」", "K1業界\n+ ニュースData", "Knowledge参照\n+ Web補完", "2-5秒", "業種コード→K1テンプレート\n+ 最新ニュースをL3取得"],
        ["「この出力おかしくない？」", "K2更新ルール\n+ K3信頼度", "ルール検証\n→ LLM説明", "1-3秒", "出力のルール適合検証\n→信頼度再計算→説明生成"],
        ["「類似事例は？」", "K5事例パターン\n+ K7顧客KM", "パターン検索\n→ LLM要約", "1-2秒", "業種×課題→K5マッチ\n→K7で顧客コンテキスト付与"],
        ["「競合の状況は？」", "K1業界\n+ K5事例", "Knowledge参照\n→ LLM分析", "2-5秒", "業界テンプレートの競合セクション\n+ 類似事例の競合動向"],
    ]

    add_table(
        slide, data,
        x=X_FULL, y=y, w=W_FULL,
        col_widths=[1.50, 1.40, 1.10, 0.70, 2.80],
        hdr_sz=6.5, cell_sz=5.5, rh=0.42,
    )

    y2 = y + 0.42 * len(data) + 0.08
    y2 = section_hdr(slide, "M4 入力Data詳細", X_FULL, y2, W_FULL)
    data2 = [
        ["入力Data", "取得元", "コンテキスト注入方式", "更新頻度"],
        ["RMの質問", "チャットUI", "プロンプト直接", "リアルタイム"],
        ["顧客情報+アジェンダ", "L2(セッション時ロード)", "システムプロンプトに事前注入", "セッション開始時"],
        ["面談履歴", "L2(セッション時ロード)", "RAGで関連面談を動的選択", "セッション開始時"],
        ["KPマップ", "L2(セッション時ロード)", "構造化データとして注入", "セッション開始時"],
    ]
    add_table(
        slide, data2,
        x=X_FULL, y=y2, w=W_FULL,
        col_widths=[1.60, 2.10, 3.00, 1.30],
        hdr_sz=6.5, cell_sz=6, rh=0.22,
    )

    y3 = y2 + 0.22 * len(data2) + 0.05
    tf = _txbox(slide, X_FULL, y3, W_FULL, 0.20)
    p = tf.paragraphs[0]
    _run(p, "設計方針: ", sz=6.5, bold=True, color=C["dk_green"])
    _run(p, "Knowledge Masterを事前にコンテキストに注入済み。RMの質問に対して即座にストリーミング開始。Web検索は必要時のみ(L3)。FB記録(暗黙知キャプチャ)は全質問で自動発火。", sz=6, color=C["gray"])


def build_slide_5(slide):
    """M5/M6/M7 + Data→Knowledge変換マップ."""
    y = section_hdr(slide, "M5/M6/M7 パイプライン比較", X_FULL, Y0, W_FULL)

    data = [
        ["パイプライン", "主要入力Data", "主要Knowledge", "即時対応", "LLM依存部分", "出力形式"],
        ["M5 提案\nストーリー", "アジェンダ+財務+KP", "K4商品+K5事例\n+K6商談", "商品候補リスト\n(K4即時)", "ストーリー構成\n・表現", "提案書ドラフト\n(Word/PPTX)"],
        ["M6 深掘り\nペーパー", "テーマ+業界Data", "K1業界+K5事例", "業界パターン\n(K1即時)", "分析・考察の\n生成", "深掘りレポート\n(Word)"],
        ["M7 メモ", "アジェンダ+面談履歴", "K6商談進行\n+K8タイミング", "次アクション候補\n(K6即時)", "メモ文面の\n生成", "面談メモ\n(テキスト)"],
    ]
    add_table(
        slide, data,
        x=X_FULL, y=y, w=W_FULL,
        col_widths=[0.95, 1.40, 1.40, 1.30, 1.20, 1.25],
        hdr_sz=6.5, cell_sz=6, rh=0.40,
    )

    y2 = y + 0.40 * len(data) + 0.10
    y2 = section_hdr(slide, "Data層 → Knowledge Master 変換マップ", X_FULL, y2, W_FULL)

    conv_data = [
        ["Data表（7対象・34項目）", "→", "Knowledge Master（K1-K8）", "変換ロジック"],
        ["業界/業種: 市場規模・競合・規制", "→", "K1 業界テンプレート", "業種コード→共通課題パターン生成"],
        ["渉外担当者: 面談記録+承認ログ", "→", "K2 更新ルール", "発言パターン→更新アクションマッピング"],
        ["顧客法人: 役員・KP+面談記録", "→", "K3 信頼度マトリクス", "発言者×具体性→信頼度スコア算出"],
        ["自社りそな: 商品マスタ+適合条件", "→", "K4 商品マッチング", "課題タグ→商品候補リスト生成"],
        ["渉外担当者: 提案履歴+成約/失注", "→", "K5 事例パターン", "類似事例→成功/失敗パターン抽出"],
        ["渉外担当者: 面談記録+活動実績", "→", "K6 商談進行パターン", "状況→推奨アクション導出"],
        ["全Data × 顧客単位で集約", "→", "K7 顧客Knowledge Master", "顧客ごとの統合コンテキスト構築"],
        ["経済/世界情勢 + 業界/業種", "→", "K8 タイミング知識", "経済環境→提案タイミング判定"],
    ]
    add_table(
        slide, conv_data,
        x=X_FULL, y=y2, w=W_FULL,
        col_widths=[2.60, 0.30, 2.10, 4.27],
        hdr_sz=6.5, cell_sz=6, rh=0.22,
    )


def build_slide_6(slide):
    """フィードバックループ詳細 + Decision Traceフロー."""
    y = section_hdr(slide, "暗黙知フィードバックループ — 詳細フロー", X_FULL, Y0, W_FULL)

    # Flow table
    data = [
        ["ステップ", "アクション", "データ種別", "保存先", "次ステップへのトリガー"],
        ["1. RM操作", "承認 / 修正 / 棄却", "操作ログ", "Decision Trace DB", "自動発火(PostToolUse Hook)"],
        ["2. Trace記録", "操作+コンテキストを記録", "構造化ログ", "Decision Trace DB", "パターン検出ジョブ(非同期)"],
        ["3a. パスA: 即時", "Data表に追加", "確定データ", "Data表(7対象)", "次回パイプライン実行時に反映"],
        ["3b. パスB: 昇格", "3回同一パターン検出", "パターン候補", "Knowledge候補キュー", "管理者レビュー or 自動昇格"],
        ["4. KM更新", "Knowledge Master更新", "Knowledge Master", "K1-K8のいずれか", "全パイプラインに反映"],
        ["5. 精度向上", "AI出力の品質が上昇", "—", "—", "RMの利用頻度増→Data蓄積加速"],
    ]
    add_table(
        slide, data,
        x=X_FULL, y=y, w=W_FULL,
        col_widths=[0.80, 1.70, 1.00, 1.50, 2.50],
        hdr_sz=6.5, cell_sz=6, rh=0.26,
    )

    y2 = y + 0.26 * len(data) + 0.10
    y2 = section_hdr(slide, "正のスパイラルの定量イメージ", X_FULL, y2, W_FULL)

    spiral_data = [
        ["フェーズ", "期間", "Data蓄積量", "Knowledge更新", "AI精度イメージ", "RM利用頻度"],
        ["PoC開始", "0-1ヶ月", "初期Data投入", "K1-K4初期構築", "ベースライン", "試用(週1-2回)"],
        ["初期運用", "1-3ヶ月", "+100-300件", "K5/K6にパターン蓄積開始", "+10-20%", "定着(週3-5回)"],
        ["安定運用", "3-6ヶ月", "+500-1000件", "K2/K3の精度向上。K8も充実", "+30-50%", "日常利用(毎日)"],
        ["成熟期", "6-12ヶ月", "+2000件以上", "全KMが自律的に進化", "+50-70%", "必須ツール化"],
    ]
    add_table(
        slide, spiral_data,
        x=X_FULL, y=y2, w=W_FULL,
        col_widths=[0.70, 0.70, 1.10, 2.20, 1.10, 1.50],
        hdr_sz=6.5, cell_sz=6, rh=0.24,
    )

    y3 = y2 + 0.24 * len(spiral_data) + 0.05
    tf = _txbox(slide, X_FULL, y3, W_FULL, 0.25)
    p = tf.paragraphs[0]
    _run(p, "エージェントハーネス対応: ", sz=6.5, bold=True, color=C["dk_green"])
    _run(p, "Decision Traceは、agentic harnessのPostToolUse Hookと同じ位置づけ。RM操作(=ツール実行)の後に自動発火し、暗黙知を確定的にキャプチャする。パスA(即時)はin-context memory、パスB(昇格)はlong-term memoryに相当。", sz=6, color=C["gray"])


def build_slide_7(slide):
    """L1/L2/L3配置 + 設計根拠."""
    y = section_hdr(slide, "即時回答設計: L1/L2/L3配置と設計根拠", X_FULL, Y0, W_FULL)

    data = [
        ["情報", "階層", "速度", "理由", "参照パイプライン", "更新頻度"],
        ["K1 業界テンプレート", "L1", "<1秒", "M1/M4で頻繁に参照。業種コードで即引き", "M1,M3,M4,M5,M6", "月次(業界変動時)"],
        ["K2 更新ルール", "L1", "<1秒", "M2で毎回使用。YAML→即判定", "M2,M4,M7", "週次(パターン蓄積)"],
        ["K3 信頼度マトリクス", "L1", "<1秒", "M2で毎回使用。テーブル照合のみ", "M2,M4", "月次(役職変動時)"],
        ["K4 商品マッチング", "L1", "<1秒", "M4/M5で頻繁に参照。課題タグで即引き", "M4,M5", "月次(商品改廃時)"],
        ["K6 商談進行パターン", "L1", "<1秒", "M4/M7で頻繁に参照", "M4,M5,M7", "週次(パターン蓄積)"],
        ["K8 タイミング知識", "L1", "<1秒", "M3/M4/M7で参照。カレンダー＋ルール", "M3,M4,M5,M7", "日次(経済イベント)"],
        ["顧客個別の財務Data", "L2", "1-5秒", "顧客アクセス時にロード", "M1,M2,M5", "セッション開始時"],
        ["面談履歴・提案履歴", "L2", "1-5秒", "セッション開始時にロード", "M2,M4,M5,M7", "セッション開始時"],
        ["最新ニュース・HP", "L3", "5-15秒", "必要時のみWeb取得", "M3,M4", "オンデマンド"],
    ]

    tbl = add_table(
        slide, data,
        x=X_FULL, y=y, w=W_FULL,
        col_widths=[1.50, 0.40, 0.55, 2.60, 1.20, 1.25],
        hdr_sz=6.5, cell_sz=6, rh=0.24,
    )

    # Color code by level
    for ri in range(1, len(data)):
        level = data[ri][1]
        cell_level = tbl.cell(ri, 1)
        if level == "L1":
            cell_level.fill.solid()
            cell_level.fill.fore_color.rgb = C["high_bg"]
        elif level == "L2":
            cell_level.fill.solid()
            cell_level.fill.fore_color.rgb = C["mid_bg"]
        elif level == "L3":
            cell_level.fill.solid()
            cell_level.fill.fore_color.rgb = C["low_bg"]

    y2 = y + 0.24 * len(data) + 0.08
    y2 = section_hdr(slide, "L1配置の設計根拠", X_FULL, y2, W_FULL)
    bullets(slide, [
        "L1配置基準: (1)アクセス頻度が高い (2)データサイズが小さい(メモリに収まる) (3)参照パイプラインが3つ以上",
        "K5(事例パターン)がL1でない理由: データ量が大きい(事例DB全体)→L2でRAG検索",
        "K7(顧客KM)がL1でない理由: 顧客単位で分散→セッション時に対象顧客分のみL2ロード",
        "L3(Web取得)は応答速度に影響 → M4チャットでは「取得中...」表示+ストリーミングで体感速度を維持",
    ], X_FULL + 0.05, y2, W_FULL - 0.10, sz=6)


# ╔══════════════════════════════════════════════════════╗
# ║  MAIN                                                ║
# ╚══════════════════════════════════════════════════════╝

def main():
    prs = Presentation(TEMPLATE)

    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Slide 0: Title
    make_title_slide(prs, "AIパイプライン × データナレッジ関係図\n【詳細版】rikyu PoC", "2026-03-10  |  POC最終成果物 — 詳細ハンドアウト")

    # Slide 1: INDEX
    make_index_slide(prs, [
        "全体消費マトリクス（詳細注釈付き）",
        "M1/M2 パイプライン詳細",
        "M4 AIチャット — 全質問パターン",
        "M5/M6/M7 + Data→Knowledge変換マップ",
        "フィードバックループ詳細 + 正のスパイラル定量イメージ",
        "L1/L2/L3配置 + 設計根拠",
    ])

    # Slide 2: Consumption matrix detailed
    s2 = make_content_slide(prs, "パイプライン × Knowledge Master 消費マトリクス")
    build_slide_2(s2)

    # Slide 3: M1/M2 detailed
    s3 = make_content_slide(prs, "M1/M2 パイプライン詳細")
    build_slide_3(s3)

    # Slide 4: M4 detailed
    s4 = make_content_slide(prs, "M4 AIチャット（折衝アシスタント）詳細")
    build_slide_4(s4)

    # Slide 5: M5/M6/M7 + Data→KM map
    s5 = make_content_slide(prs, "M5/M6/M7 + Data→Knowledge変換")
    build_slide_5(s5)

    # Slide 6: Feedback loop detailed
    s6 = make_content_slide(prs, "フィードバックループ詳細")
    build_slide_6(s6)

    # Slide 7: L1/L2/L3 + design rationale
    s7 = make_content_slide(prs, "即時回答設計: L1/L2/L3配置")
    build_slide_7(s7)

    prs.save(OUTFILE)
    print(f"Saved detailed PPTX ({len(prs.slides)} slides) to {OUTFILE}")


if __name__ == "__main__":
    main()
