#!/usr/bin/env python3
"""Build standard PPTX for 15_tacit_knowledge_lifecycle (~16 slides).

Covers:
- What tacit knowledge is ("state" not "type")
- 5-stage lifecycle
- Capture: 5 touchpoints (C1-C5), zero effort principle
- Accumulate -> Promote: Path A (5 items), Path B (8 items)
- Deploy: 4-layer memory hierarchy
- Evolve: AI-Human Gap concept
- Extension: auto + human evolution
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TEMPLATE = "project-rikyu-sales-proposals-poc/untracked/original_document/report/\u63d0\u6848\u4eee\u8aac\u69cb\u7bc9AI_\u5b9a\u4f8b\u4f1a1_20260219_v1.pptx"
WORKDIR = "output/15_tacit_lifecycle"

# === ACES Color Palette ===
GREEN_DARK = RGBColor(0x00, 0x68, 0x43)
GREEN_MEDIUM = RGBColor(0x19, 0x7A, 0x56)
GREEN_BRIGHT = RGBColor(0x00, 0x9A, 0x62)
GREEN_MINT = RGBColor(0xC2, 0xE7, 0xD9)
GREEN_PALE = RGBColor(0xDF, 0xE6, 0xE0)
CHARCOAL = RGBColor(0x39, 0x39, 0x39)
TEXT_MAIN = RGBColor(0x12, 0x12, 0x12)
TABLE_TEXT = RGBColor(0x45, 0x45, 0x45)
BORDER_GRAY = RGBColor(0xC7, 0xC7, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x52, 0x98, 0xBA)
ACCENT_RED = RGBColor(0xDB, 0x5F, 0x5F)
ACCENT_ORANGE = RGBColor(0xFB, 0xAE, 0x40)
ACCENT_GOLD = RGBColor(0xC0, 0xB7, 0x6A)
GRAY_LIGHT = RGBColor(0x99, 0x99, 0x99)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

TITLE_FONT = "Yu Gothic"
SUB_FONT = "Yu Gothic Medium"
BODY_FONT = "Century Gothic"


# ── Helper functions ──────────────────────────────────

def find_shapes(slide):
    """Find title, subtitle, body shapes by Y-coordinate."""
    shapes_with_y = []
    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            y = s.top / 914400 if s.top else 0
            shapes_with_y.append((y, s))
    shapes_with_y.sort(key=lambda x: x[0])
    title = subtitle = body = None
    for y, s in shapes_with_y:
        if y < 0.7 and title is None: title = s
        elif 0.7 <= y < 1.5 and subtitle is None: subtitle = s
        elif y >= 1.5 and body is None: body = s
    return title, subtitle, body


def prep(slide):
    """Suppress layout bleed-through by clearing body placeholder."""
    title, subtitle, body = find_shapes(slide)
    if subtitle:
        subtitle.height = Inches(0.32)
        if subtitle.width < Inches(1):
            subtitle.width = Inches(12.46)
            subtitle.left = Inches(0.42)
    if body:
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = " "
        p.font.size = Pt(1)
        p.font.color.rgb = WHITE
    return title, subtitle, body


def add_table(slide, data, left=Inches(0.42), top=Inches(1.90), width=Inches(12.46),
              col_widths=None, header_font_pt=9, cell_font_pt=8, row_height=0.32):
    """Styled table with green header + alternating rows."""
    rows, cols = len(data), len(data[0])
    height = Inches(row_height * rows)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.name = BODY_FONT
                if row_idx == 0:
                    para.font.size = Pt(header_font_pt)
                    para.font.color.rgb = WHITE
                    para.font.bold = True
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.font.size = Pt(cell_font_pt)
                    para.font.color.rgb = TABLE_TEXT
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_DARK
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_PALE
    return table


def add_box(slide, text, left, top, width, height, bg_color=GREEN_MINT, text_color=GREEN_DARK,
            font_size=9, bold=False, align=PP_ALIGN.CENTER):
    """Add a colored text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_arrow_down(slide, left, top, width=Inches(0.5), height=Inches(0.3)):
    """Add downward arrow text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u25bc"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN_DARK
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_arrow_right(slide, left, top, width=Inches(0.6), height=Inches(0.3)):
    """Add right arrow text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u2192"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_note(slide, text, left, top, width, h=Inches(0.35)):
    """Add an italic note text."""
    txBox = slide.shapes.add_textbox(left, top, width, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(6.5)
    p.font.italic = True
    p.font.color.rgb = GRAY_LIGHT
    p.font.name = BODY_FONT


def set_title(slide, title_text, subtitle_text):
    """Set title and subtitle for a content slide."""
    title, subtitle, body = prep(slide)
    if title:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = title_text
        p.font.name = TITLE_FONT
        p.font.size = Pt(22)
    if subtitle:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        p.text = subtitle_text
        p.font.name = SUB_FONT
        p.font.size = Pt(11)


def set_section_title(slide, section_text):
    """Set title for a section divider slide."""
    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            tf = s.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = section_text
            p.font.name = TITLE_FONT
            p.font.size = Pt(28)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT
            break


# ── Build slide structure ──────────────────────────────

def build():
    """Build complete standard PPTX."""
    prs = Presentation(TEMPLATE)

    # Get slide layouts
    layout_cover_data = prs.slide_layouts[1]   # Cover+data
    layout_index = prs.slide_layouts[3]        # Index_1column
    layout_section = prs.slide_layouts[4]      # Section Cover
    layout_content = prs.slide_layouts[10]     # Contents_Title&Message_Left

    # Build new presentation
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    # Slide plan: 16 slides
    # 0: Title (cover)
    # 1: INDEX
    # 2: Section - \u6697\u9ed9\u77e5\u3068\u306f\u4f55\u304b
    # 3: Content - "\u72b6\u614b"\u306e\u8aac\u660e
    # 4: Section - \u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb\u5168\u4f53\u50cf
    # 5: Content - 5-stage flow
    # 6: Section - \u2460\u30ad\u30e3\u30d7\u30c1\u30e3
    # 7: Content - 5 touchpoints table
    # 8: Content - Zero effort principle
    # 9: Section - \u2461\u84c4\u7a4d\u2192\u2462\u6607\u683c
    # 10: Content - Path A table
    # 11: Content - Path B table
    # 12: Section - \u2463\u5c55\u958b
    # 13: Content - 4-layer memory hierarchy
    # 14: Section - \u2464\u9032\u5316
    # 15: Content - AI-Human Gap + spiral
    # (16 slides total)

    slide_types = [
        layout_cover_data,  # 0
        layout_index,       # 1
        layout_section,     # 2
        layout_content,     # 3
        layout_section,     # 4
        layout_content,     # 5
        layout_section,     # 6
        layout_content,     # 7
        layout_content,     # 8
        layout_section,     # 9
        layout_content,     # 10
        layout_content,     # 11
        layout_section,     # 12
        layout_content,     # 13
        layout_section,     # 14
        layout_content,     # 15
    ]

    for layout in slide_types:
        new_prs.slides.add_slide(layout)

    slides = new_prs.slides

    # ── Slide 0: Title ──
    slide = slides[0]
    shapes_with_y = []
    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            y = s.top / 914400 if s.top else 0
            shapes_with_y.append((y, s))
    shapes_with_y.sort(key=lambda x: x[0])
    shape_list = [s for _, s in shapes_with_y]
    if len(shape_list) >= 1:
        tf = shape_list[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "\u682a\u5f0f\u4f1a\u793e\u308a\u305d\u306a\u9280\u884c \u5fa1\u4e2d"
        p.font.size = Pt(32)
        p2 = tf.add_paragraph()
        p2.text = "\u6697\u9ed9\u77e5\u306e\u30c7\u30fc\u30bf\u5316\u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb"
        p2.font.size = Pt(28)
    if len(shape_list) >= 2:
        tf = shape_list[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "2026/3/10 | rikyu \u30bd\u30ea\u30e5\u30fc\u30b7\u30e7\u30f3\u30bb\u30fc\u30eb\u30b9\u4f34\u8d70AI PoC"
        p.font.size = Pt(14)

    # ── Slide 1: INDEX ──
    slide = slides[1]
    title_s, subtitle_s, body_s = find_shapes(slide)
    if title_s:
        title_s.text_frame.clear()
        p = title_s.text_frame.paragraphs[0]
        p.text = "INDEX"
        p.font.name = TITLE_FONT
        p.font.size = Pt(28)
    if body_s:
        tf = body_s.text_frame
        tf.clear()
        items = [
            "1. \u6697\u9ed9\u77e5\u3068\u306f\u4f55\u304b \u2014 \u300c\u7a2e\u985e\u300d\u3067\u306f\u306a\u304f\u300c\u72b6\u614b\u300d",
            "2. \u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb\u5168\u4f53\u50cf \u2014 5\u6bb5\u968e\u306e\u30d5\u30ed\u30fc",
            "3. \u2460\u30ad\u30e3\u30d7\u30c1\u30e3 \u2014 \u8ffd\u52a0\u64cd\u4f5c\u306a\u3057\u3067\u6355\u7372\u3059\u308b",
            "4. \u2461\u84c4\u7a4d\u2192\u2462\u6607\u683c \u2014 2\u3064\u306e\u30d1\u30b9",
            "5. \u2463\u5c55\u958b \u2014 \u591a\u5c64\u30e1\u30e2\u30ea\u3067\u7d44\u7e54\u306b\u914d\u5e03",
            "6. \u2464\u9032\u5316 \u2014 AI-Human Gap = \u6700\u9ad8\u306e\u6559\u5e2b\u30c7\u30fc\u30bf",
        ]
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.name = BODY_FONT
            p.font.size = Pt(14)
            p.font.color.rgb = CHARCOAL
            p.space_after = Pt(8)

    # ── Slide 2: Section - \u6697\u9ed9\u77e5\u3068\u306f\u4f55\u304b ──
    set_section_title(slides[2], "1. \u6697\u9ed9\u77e5\u3068\u306f\u4f55\u304b")

    # ── Slide 3: "\u72b6\u614b"\u306e\u8aac\u660e ──
    slide = slides[3]
    set_title(slide, "\u6697\u9ed9\u77e5\u306f\u300c\u7a2e\u985e\u300d\u3067\u306f\u306a\u304f\u300c\u72b6\u614b\u300d",
              "\u672a\u8a18\u9332\u30fb\u672a\u8a00\u8a9e\u5316\u3068\u3044\u3046\u300c\u72b6\u614b\u300d\u304b\u3089\u3001Data/Knowledge\u3078\u5909\u63db\u3059\u308b\u3060\u3051")

    # 3 example boxes with before/after
    examples = [
        ("\u300cA\u793e\u306e\u7dcf\u52d9\u90e8\u9577\u3068\u5c02\u52d9\u306f\u4e0d\u4ef2\u300d", "\u672a\u8a18\u9332\u306a\u3089 \u2192 \u6697\u9ed9\u77e5\nKP\u30de\u30c3\u30d7\u306b\u8a18\u9332\u3057\u305f\u3089 \u2192 Data", ACCENT_BLUE),
        ("\u300c\u7269\u6d41\u696d\u306f\u4eba\u4ef6\u8cbb\u304cKey Driver\u300d", "\u30d9\u30c6\u30e9\u30f3\u306e\u982d\u306e\u4e2d\u306a\u3089 \u2192 \u6697\u9ed9\u77e5\n\u30eb\u30fc\u30eb\u5316\u3057\u305f\u3089 \u2192 Knowledge", ACCENT_ORANGE),
        ("\u300c\u3053\u306e\u51fa\u529b\u306f\u7684\u5916\u308c\u3060\u300d\u3068\u3044\u3046\u76f4\u611f", "\u6697\u9ed9\u77e5\nFB\u7406\u7531\u3092\u8a00\u8a9e\u5316\u3057\u305f\u3089 \u2192 Knowledge", ACCENT_GOLD),
    ]
    y_pos = 2.0
    for label, desc, color in examples:
        # Label box
        add_box(slide, label,
                Inches(0.42), Inches(y_pos), Inches(4.5), Inches(0.50),
                bg_color=color, text_color=WHITE, font_size=11, bold=True, align=PP_ALIGN.LEFT)
        # Arrow
        add_arrow_right(slide, Inches(5.10), Inches(y_pos + 0.08))
        # Description
        add_box(slide, desc,
                Inches(5.80), Inches(y_pos), Inches(6.9), Inches(0.50),
                bg_color=BG_LIGHT, text_color=CHARCOAL, font_size=9, align=PP_ALIGN.LEFT)
        y_pos += 0.65

    # Key insight
    add_box(slide, "\u30dd\u30a4\u30f3\u30c8: \u6697\u9ed9\u77e5\u306f\u300c\u5225\u306e\u60c5\u5831\u300d\u3067\u306f\u306a\u304f\u3001\u307e\u3060Data\u306b\u3082Knowledge\u306b\u3082\u306a\u3063\u3066\u3044\u306a\u3044\u300c\u72b6\u614b\u300d\u3002\u8a18\u9332\u30fb\u8a00\u8a9e\u5316\u3059\u308c\u3070\u6d3b\u7528\u53ef\u80fd\u306b\u306a\u308b\u3002",
            Inches(0.42), Inches(4.2), Inches(12.46), Inches(0.45),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=10, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 4: Section - \u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb\u5168\u4f53\u50cf ──
    set_section_title(slides[4], "2. \u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb\u5168\u4f53\u50cf")

    # ── Slide 5: 5-stage flow diagram ──
    slide = slides[5]
    set_title(slide, "\u6697\u9ed9\u77e5\u30e9\u30a4\u30d5\u30b5\u30a4\u30af\u30eb: 5\u6bb5\u968e\u306e\u30d5\u30ed\u30fc",
              "\u30ad\u30e3\u30d7\u30c1\u30e3\u304b\u3089\u9032\u5316\u307e\u3067\u3001\u6b63\u306e\u30b9\u30d1\u30a4\u30e9\u30eb\u3067AI\u304c\u8ce2\u304f\u306a\u308b")

    stages = [
        ("\u2460 \u30ad\u30e3\u30d7\u30c1\u30e3\uff08\u6355\u7372\uff09", "RM\u306e\u65e5\u5e38\u696d\u52d9\u306e\u4e2d\u3067\u3001\u8ffd\u52a0\u64cd\u4f5c\u306a\u3057\u306b\u6697\u9ed9\u77e5\u3092\u8a18\u9332", GREEN_DARK, WHITE),
        ("\u2461 \u84c4\u7a4d\uff08\u30d1\u30bf\u30fc\u30f3\u691c\u51fa\uff09", "\u540c\u3058\u7a2e\u985e\u306e\u64cd\u4f5c\u30fb\u5224\u65ad\u304c\u84c4\u7a4d\u3055\u308c\u308b\u3068\u30d1\u30bf\u30fc\u30f3\u304c\u898b\u3048\u308b", GREEN_MEDIUM, WHITE),
        ("\u2462 \u6607\u683c\uff08Data/Knowledge\u5316\uff09", "\u30d1\u30bf\u30fc\u30f3\u304c\u95be\u5024\u3092\u8d85\u3048\u305f\u3089Data\u307e\u305f\u306fKnowledge\u306b\u5909\u63db", GREEN_BRIGHT, WHITE),
        ("\u2463 \u5c55\u958b\uff08\u591a\u5c64\u30e1\u30e2\u30ea\u3067\u914d\u5e03\uff09", "\u500b\u4eba\u2192\u55b6\u696d\u5e97\u2192\u696d\u754c\u2192\u5168\u793e\u306e\u30eb\u30fc\u30c8\u3067\u5171\u6709", ACCENT_BLUE, WHITE),
        ("\u2464 \u9032\u5316\uff08AI\u51fa\u529b\u7cbe\u5ea6\u5411\u4e0a\uff09", "\u6607\u683c\u3057\u305fKnowledge\u304cAI\u306e\u5224\u65ad\u3092\u6539\u5584\u3059\u308b", ACCENT_GOLD, CHARCOAL),
    ]
    y_stage = 1.90
    box_w = Inches(11.50)
    box_h = Inches(0.55)
    for i, (title_text, desc, bg, tc) in enumerate(stages):
        add_box(slide, f"{title_text}\n{desc}",
                Inches(0.42), Inches(y_stage), box_w, box_h,
                bg_color=bg, text_color=tc, font_size=9, bold=False, align=PP_ALIGN.LEFT)
        y_stage += 0.55
        if i < len(stages) - 1:
            add_arrow_down(slide, Inches(6.0), Inches(y_stage - 0.02), Inches(0.5), Inches(0.22))
            y_stage += 0.20

    # Spiral loop
    add_box(slide, "\u21b3 \u2464\u9032\u5316 \u2192 \u2460\u30ad\u30e3\u30d7\u30c1\u30e3\u306b\u623b\u308b\uff08\u6b63\u306e\u30b9\u30d1\u30a4\u30e9\u30eb\uff09",
            Inches(0.42), Inches(y_stage + 0.10), Inches(11.50), Inches(0.35),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=10, bold=True, align=PP_ALIGN.CENTER)

    # ── Slide 6: Section - \u2460\u30ad\u30e3\u30d7\u30c1\u30e3 ──
    set_section_title(slides[6], "3. \u2460\u30ad\u30e3\u30d7\u30c1\u30e3")

    # ── Slide 7: 5 touchpoints ──
    slide = slides[7]
    set_title(slide, "\u30ad\u30e3\u30d7\u30c1\u30e3\u306e5\u3064\u306e\u30bf\u30c3\u30c1\u30dd\u30a4\u30f3\u30c8",
              "RM\u304c\u666e\u6bb5\u901a\u308a\u306e\u696d\u52d9\u3092\u3059\u308b\u3060\u3051\u3067\u3001\u6697\u9ed9\u77e5\u304c\u81ea\u52d5\u7684\u306b\u8a18\u9332\u3055\u308c\u308b")

    tp_data = [
        ["#", "\u30bf\u30c3\u30c1\u30dd\u30a4\u30f3\u30c8", "RM\u306e\u64cd\u4f5c", "\u30ad\u30e3\u30d7\u30c1\u30e3\u3055\u308c\u308b\u6697\u9ed9\u77e5", "\u8a18\u9332\u5148"],
        ["C1", "\u9762\u8ac7\u8a18\u9332\u5165\u529b", "\u9762\u8ac7\u5f8c\u306b\u8a18\u9332\u3092\u5165\u529b", "\u4ed6\u884c\u6761\u4ef6, KP\u95a2\u4fc2\u6027, \u8a2d\u5099\u72b6\u6cc1,\n\u30aa\u30fc\u30ca\u30fc\u60c5\u5831, \u7af6\u5408\u52d5\u5411", "\u9762\u8ac7DB"],
        ["C2", "AI\u66f4\u65b0\u63d0\u6848\u30ec\u30d3\u30e5\u30fc", "\u627f\u8a8d/\u4fee\u6b63/\u68c4\u5374+\u7406\u7531", "\u300c\u3053\u306e\u767a\u8a00\u306a\u3089\u66f4\u65b0\u3059\u3079\u304d\u300d\u306e\u52d8\u6240,\n\u300c\u3053\u306e\u4fe1\u983c\u5ea6\u306f\u9055\u3046\u300d\u306e\u611f\u899a", "\u627f\u8a8d\u30ed\u30b0DB"],
        ["C3", "AI\u30c1\u30e3\u30c3\u30c8\u5229\u7528", "\u8cea\u554f+FB(\u7684\u5916\u308c\u30dc\u30bf\u30f3)", "\u300c\u3053\u306e\u51fa\u529b\u306f\u7684\u5916\u308c\u300d\u306e\u7406\u7531,\n\u3088\u304f\u805e\u304b\u308c\u308b\u8cea\u554f\u30d1\u30bf\u30fc\u30f3", "\u30c1\u30e3\u30c3\u30c8\u30ed\u30b0DB"],
        ["C4", "\u6210\u7d04/\u5931\u6ce8\u8a18\u9332", "\u7d50\u679c+\u8981\u56e0\u3092\u8a18\u9332", "\u300c\u523a\u3055\u3063\u305f/\u523a\u3055\u3089\u306a\u304b\u3063\u305f\u300d\u306e\u8a18\u61b6,\n\u5546\u54c1\u9078\u5b9a\u7406\u7531", "\u63d0\u6848DB"],
        ["C5", "\u30a8\u30ad\u30b9\u30d1\u30fc\u30c8\u30d2\u30a2\u30ea\u30f3\u30b0", "\u30d2\u30a2\u30ea\u30f3\u30b0\u306b\u56de\u7b54", "\u696d\u754c\u30d1\u30bf\u30fc\u30f3, \u5546\u8ac7\u9032\u884c\u306e\u30b3\u30c4,\nNG\u4e8b\u9805", "Knowledge\nMaster"],
    ]
    add_table(slide, tp_data, top=Inches(1.85),
              col_widths=[0.5, 2.0, 2.0, 4.5, 1.5],
              header_font_pt=8, cell_font_pt=7, row_height=0.55)

    # ── Slide 8: Zero effort principle ──
    slide = slides[8]
    set_title(slide, "\u300c\u8ffd\u52a0\u64cd\u4f5c\u306a\u3057\u300d\u306e\u539f\u5247",
              "\u6697\u9ed9\u77e5\u8a18\u9332\u306e\u305f\u3081\u306e\u5c02\u7528\u753b\u9762\u306f\u4f5c\u3089\u306a\u3044")

    # Bad design
    add_box(slide, "\u2717 \u60aa\u3044\u8a2d\u8a08", Inches(0.42), Inches(2.00), Inches(5.8), Inches(0.35),
            bg_color=ACCENT_RED, text_color=WHITE, font_size=11, bold=True)
    add_box(slide, "\u300c\u6697\u9ed9\u77e5\u3092\u8a18\u9332\u3057\u3066\u304f\u3060\u3055\u3044\u300d\n\u2192 RM\u306e\u8ca0\u8377\u589e \u2192 \u4f7f\u308f\u308c\u306a\u3044",
            Inches(0.42), Inches(2.40), Inches(5.8), Inches(0.80),
            bg_color=RGBColor(0xFD, 0xE8, 0xE8), text_color=CHARCOAL, font_size=10, align=PP_ALIGN.LEFT)

    # Good design
    add_box(slide, "\u2713 \u826f\u3044\u8a2d\u8a08", Inches(6.50), Inches(2.00), Inches(6.3), Inches(0.35),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=11, bold=True)
    add_box(slide, "\u666e\u6bb5\u306e\u64cd\u4f5c\uff08\u9762\u8ac7\u8a18\u9332\u3001\u627f\u8a8d\u3001\u30c1\u30e3\u30c3\u30c8\uff09\u306e\u4e2d\u306b\n\u69cb\u9020\u5316\u9805\u76ee\u3092\u57cb\u3081\u8fbc\u3080\n\u2192 \u81ea\u7136\u306b\u8a18\u9332\u3055\u308c\u308b",
            Inches(6.50), Inches(2.40), Inches(6.3), Inches(0.80),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=10, align=PP_ALIGN.LEFT)

    # Example
    add_box(slide, "\u4f8b: \u9762\u8ac7\u8a18\u9332\u753b\u9762\u306e\u300c\u4ed6\u884c\u60c5\u5831\u300d\u300cKP\u95a2\u4fc2\u6027\u300d\u9805\u76ee\u306f\u3001RM\u304c\u9762\u8ac7\u3067\u805e\u3044\u305f\u60c5\u5831\u3092\u81ea\u7136\u306b\u8a18\u9332\u3059\u308b\u5c0e\u7dda\u3002\n\u308f\u3056\u308f\u3056\u300c\u6697\u9ed9\u77e5\u5165\u529b\u300d\u753b\u9762\u3092\u4f5c\u3089\u306a\u3044\u3002",
            Inches(0.42), Inches(3.60), Inches(12.46), Inches(0.55),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=9, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 9: Section - \u2461\u84c4\u7a4d\u2192\u2462\u6607\u683c ──
    set_section_title(slides[9], "4. \u2461\u84c4\u7a4d \u2192 \u2462\u6607\u683c")

    # ── Slide 10: Path A ──
    slide = slides[10]
    set_title(slide, "\u30d1\u30b9A: \u8a18\u9332\u3059\u308c\u3070\u5373Data\uff085\u9805\u76ee\uff09",
              "\u9762\u8ac7\u8a18\u9332\u306e\u69cb\u9020\u5316\u9805\u76ee\u306b\u5165\u529b \u2192 \u8a18\u9332\u3057\u305f\u6642\u70b9\u3067Data\u5c64\u306b\u5373\u8ffd\u52a0")

    path_a = [
        ["\u6697\u9ed9\u77e5", "\u8a18\u9332\u65b9\u6cd5", "\u6607\u683c\u5148", "\u6d3b\u7528\u30bf\u30a4\u30df\u30f3\u30b0"],
        ["\u9762\u8ac7\u3067\u805e\u3044\u305f\u4ed6\u884c\u6761\u4ef6", "\u300c\u4ed6\u884c\u60c5\u5831\u300d\u9805\u76ee", "Data: \u4ed6\u884c\u63d0\u793a\u6761\u4ef6", "\u6b21\u56de\u306e\u63d0\u6848\u6642"],
        ["KP\u9593\u306e\u4eba\u9593\u95a2\u4fc2\u30fb\u529b\u5b66", "\u300cKP\u95a2\u4fc2\u6027\u300d\u9805\u76ee", "Data: \u5f79\u54e1\u30fbKP", "KP\u30de\u30c3\u30d7\u306b\u5373\u53cd\u6620"],
        ["\u8a2d\u5099\u30fb\u6295\u8cc7\u306e\u73fe\u6cc1", "\u300c\u8a2d\u5099\u30fb\u6295\u8cc7\u300d\u9805\u76ee", "Data: \u4e8b\u696d\u69cb\u9020", "\u30a2\u30b8\u30a7\u30f3\u30c0\u66f4\u65b0\u63d0\u6848\u306b\u53cd\u6620"],
        ["\u7d4c\u55b6\u8005\u306e\u500b\u4eba\u7684\u72b6\u6cc1", "\u300c\u30aa\u30fc\u30ca\u30fc\u60c5\u5831\u300d\u9805\u76ee", "Data: \u7d4c\u55b6\u8005\u57fa\u672c\u60c5\u5831", "\u4e8b\u696d\u627f\u7d99\u5224\u5b9a\u306b\u53cd\u6620"],
        ["\u7af6\u5408\u653b\u52e2\u306e\u60c5\u5831", "\u300c\u7af6\u5408\u52d5\u5411\u300d\u9805\u76ee", "Data: \u7af6\u5408\u653b\u52e2\u60c5\u5831", "\u63d0\u6848\u6226\u7565\u306b\u53cd\u6620"],
    ]
    add_table(slide, path_a, top=Inches(1.85),
              col_widths=[3.0, 2.5, 3.0, 3.0],
              header_font_pt=9, cell_font_pt=8, row_height=0.36)

    # Flow explanation
    add_box(slide, "RM\u306e\u9762\u8ac7\u8a18\u9332\u5165\u529b \u2192 \u69cb\u9020\u5316\u9805\u76ee\u306b\u5165\u529b \u2192 \u8a18\u9332\u3057\u305f\u6642\u70b9\u3067Data\u5c64\u306b\u5373\u8ffd\u52a0 \u2192 \u6b21\u56de\u306eAI\u51e6\u7406\u3067\u5373\u6d3b\u7528",
            Inches(0.42), Inches(4.30), Inches(12.46), Inches(0.38),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=9, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 11: Path B ──
    slide = slides[11]
    set_title(slide, "\u30d1\u30b9B: \u30d1\u30bf\u30fc\u30f3\u84c4\u7a4d \u2192 Knowledge\u5316\uff088\u9805\u76ee\uff09",
              "\u540c\u4e00\u30d1\u30bf\u30fc\u30f3\u304c3\u56de\u84c4\u7a4d \u2192 \u30d1\u30bf\u30fc\u30f3\u3092\u8a00\u8a9e\u5316 \u2192 Knowledge Master\u306b\u8ffd\u52a0")

    path_b = [
        ["\u6697\u9ed9\u77e5", "\u84c4\u7a4d\u65b9\u6cd5", "\u6607\u683c\u6761\u4ef6", "\u6607\u683c\u5148Knowledge"],
        ["\u300c\u3053\u306e\u767a\u8a00\u306a\u3089\u66f4\u65b0\u3059\u3079\u304d\u300d", "\u627f\u8a8d/\u4fee\u6b63/\u68c4\u5374\u30d1\u30bf\u30fc\u30f3", "\u540c\u4e003\u56de", "K2: \u66f4\u65b0\u30eb\u30fc\u30eb"],
        ["\u300c\u3053\u306e\u4fe1\u983c\u5ea6\u306f\u9055\u3046\u300d", "\u4fe1\u983c\u5ea6\u4e0a\u66f8\u304d\u4fee\u6b63", "\u540c\u4e00\u4fee\u6b63\u84c4\u7a4d", "K3: \u4fe1\u983c\u5ea6\u30de\u30c8\u30ea\u30af\u30b9"],
        ["\u300c\u3053\u3046\u805e\u304f\u3068\u524d\u306b\u9032\u3080\u300d", "\u6210\u529f\u30a2\u30b8\u30a7\u30f3\u30c0\u306e\u9762\u8ac7\u5206\u6790", "\u6210\u529f\u30d1\u30bf\u30fc\u30f3\u62bd\u51fa", "K6: \u5546\u8ac7\u9032\u884c\u30d1\u30bf\u30fc\u30f3"],
        ["\u300c\u3053\u306e\u9867\u5ba2\u306b\u306f\u3053\u3046\u5165\u308b\u300d", "\u30c1\u30e3\u30c3\u30c8FB", "PJ\u30c1\u30fc\u30e0\u30ec\u30d3\u30e5\u30fc", "K6: \u9867\u5ba2\u30bf\u30a4\u30d7\u5225"],
        ["\u300c\u3053\u306e\u696d\u754c\u306f\u4eca\u3053\u3046\u3044\u3046\u7a7a\u6c17\u300d", "\u30c1\u30e3\u30c3\u30c8FB", "\u8907\u6570RM\u540c\u50be\u5411", "K1: \u696d\u754c\u5225\u8ab2\u984c"],
        ["\u300c\u3053\u306e\u8ab2\u984c\u306b\u306f\u3053\u306e\u5546\u54c1\u300d", "\u6210\u7d04\u6642\u306e\u9078\u5b9a\u7406\u7531", "\u6210\u529f\u30d1\u30bf\u30fc\u30f3\u84c4\u7a4d", "K4: \u8ab2\u984c\u00d7\u5546\u54c1"],
        ["\u300c\u523a\u3055\u3063\u305f/\u523a\u3055\u3089\u306a\u304b\u3063\u305f\u300d", "\u6210\u7d04/\u5931\u6ce8\u8981\u56e0", "\u4e8b\u4f8b\u69cb\u9020\u5316", "K5: \u6210\u7d04/\u5931\u6ce8\u30d1\u30bf\u30fc\u30f3"],
        ["\u300c\u3053\u306e\u51fa\u529b\u306f\u7684\u5916\u308c\u300d", "\u30c1\u30e3\u30c3\u30c8FB\u30dc\u30bf\u30f3", "\u540c\u7a2eFB\u84c4\u7a4d", "K2: \u91cd\u8981\u5ea6\u5224\u5b9a"],
    ]
    add_table(slide, path_b, top=Inches(1.85),
              col_widths=[3.0, 2.5, 2.0, 3.0],
              header_font_pt=8, cell_font_pt=7, row_height=0.34)

    # Flow explanation
    add_box(slide, "RM\u306e\u627f\u8a8d/\u4fee\u6b63/\u68c4\u5374 \u2192 Decision Trace\u306b\u8a18\u9332 \u2192 \u540c\u4e00\u30d1\u30bf\u30fc\u30f33\u56de\u84c4\u7a4d \u2192 \u30d1\u30bf\u30fc\u30f3\u3092\u8a00\u8a9e\u5316 \u2192 K\u30eb\u30fc\u30eb\u3068\u3057\u3066\u8ffd\u52a0 \u2192 \u5168RM\u306eAI\u51fa\u529b\u54c1\u8cea\u304c\u5411\u4e0a",
            Inches(0.42), Inches(5.10), Inches(12.46), Inches(0.38),
            bg_color=GREEN_PALE, text_color=GREEN_DARK, font_size=8, bold=True, align=PP_ALIGN.LEFT)

    # ── Slide 12: Section - \u2463\u5c55\u958b ──
    set_section_title(slides[12], "5. \u2463\u5c55\u958b")

    # ── Slide 13: 4-layer memory hierarchy ──
    slide = slides[13]
    set_title(slide, "\u591a\u5c64\u30e1\u30e2\u30ea\u968e\u5c64: 4\u5c64 + \u9867\u5ba2\u30e1\u30e2\u30ea",
              "\u500b\u4eba\u2192\u55b6\u696d\u5e97\u2192\u696d\u754c\u2192\u5168\u793e\u306e\u30eb\u30fc\u30c8\u3067\u30ca\u30ec\u30c3\u30b8\u3092\u5171\u6709\u30fb\u6607\u683c")

    # Nested boxes for memory hierarchy (outermost first)
    # L4: Global (full width)
    add_box(slide, "Level 4: \u5168\u793e\uff08Global\uff09\n\u5168\u793e\u5171\u901a\u306e\u30eb\u30fc\u30eb\u3002\u300c\u5168\u696d\u754c\u3067\u88dc\u52a9\u91d1\u306f\u523a\u3055\u308b\u300d\u7b49",
            Inches(0.42), Inches(1.85), Inches(8.0), Inches(3.80),
            bg_color=RGBColor(0xE8, 0xF0, 0xE8), text_color=CHARCOAL, font_size=9, align=PP_ALIGN.LEFT)

    # L3: Industry
    add_box(slide, "Level 3: \u696d\u754c\uff08Industry\uff09\n\u696d\u754c\u3054\u3068\u306e\u30d1\u30bf\u30fc\u30f3\u3002\u300c\u88fd\u9020\u696d\u306f\u8a2d\u5099\u6295\u8cc7\u30b5\u30a4\u30af\u30eb5-7\u5e74\u300d",
            Inches(0.72), Inches(2.60), Inches(7.0), Inches(2.80),
            bg_color=RGBColor(0xD0, 0xE4, 0xD0), text_color=CHARCOAL, font_size=9, align=PP_ALIGN.LEFT)

    # L2: Branch
    add_box(slide, "Level 2: \u55b6\u696d\u5e97\uff08Branch\uff09\n\u5730\u57df\u7279\u6027\u3002\u300c\u3053\u306e\u5730\u57df\u306f\u5efa\u8a2d\u696d\u304c\u591a\u3044\u300d",
            Inches(1.02), Inches(3.30), Inches(6.0), Inches(1.80),
            bg_color=RGBColor(0xC2, 0xE7, 0xD9), text_color=CHARCOAL, font_size=9, align=PP_ALIGN.LEFT)

    # L1: Personal RM
    add_box(slide, "Level 1: \u500b\u4eba\uff08RM\uff09\nRM\u500b\u4eba\u306e\u7d4c\u9a13\u3002\u300cA\u793e\u306f\u6570\u5b57\u3088\u308a\u4eba\u60c5\u300d",
            Inches(1.32), Inches(3.95), Inches(5.0), Inches(0.85),
            bg_color=GREEN_MINT, text_color=GREEN_DARK, font_size=9, bold=True, align=PP_ALIGN.LEFT)

    # Customer memory (side box)
    add_box(slide, "\u2605 \u9867\u5ba2\u30e1\u30e2\u30ea\uff08\u6a2a\u65ad\u7684\u95a2\u5fc3\u4e8b\uff09\n\u9867\u5ba2\u5358\u4f4d\u306e\u60c5\u5831\u306f\u5168\u30ec\u30d9\u30eb\u3092\u6a2a\u65ad",
            Inches(8.80), Inches(1.85), Inches(4.0), Inches(1.20),
            bg_color=ACCENT_BLUE, text_color=WHITE, font_size=9, bold=True, align=PP_ALIGN.LEFT)

    # Promotion route
    add_box(slide, "\u6607\u683c\u30eb\u30fc\u30c8\u4f8b",
            Inches(8.80), Inches(3.25), Inches(4.0), Inches(0.28),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=9, bold=True)

    routes = [
        "\u500b\u4ebaRM\u300cA\u793e\u306f\u6570\u5b57\u3067\u8aac\u5f97\u3059\u3079\u304d\u300d",
        "\u2193 \u540c\u3058\u55b6\u696d\u5e97\u306e\u8907\u6570RM\u304c\u540c\u50be\u5411FB",
        "\u55b6\u696d\u5e97\u30ec\u30d9\u30eb\u306e\u30d1\u30bf\u30fc\u30f3\u306b\u6607\u683c",
        "\u2193 \u8907\u6570\u55b6\u696d\u5e97\u3067\u540c\u50be\u5411",
        "\u696d\u754c\u30ec\u30d9\u30eb\u306e\u30d1\u30bf\u30fc\u30f3\u306b\u6607\u683c",
        "\u2193 \u4ed6\u696d\u754c\u3067\u3082\u540c\u69cb\u9020\u3092\u78ba\u8a8d",
        "\u5168\u793e\u30eb\u30fc\u30eb\u306b\u6607\u683c",
    ]
    y_rt = 3.58
    for rt in routes:
        add_box(slide, rt,
                Inches(8.80), Inches(y_rt), Inches(4.0), Inches(0.22),
                bg_color=GREEN_PALE, text_color=CHARCOAL, font_size=7, align=PP_ALIGN.LEFT)
        y_rt += 0.24

    # ── Slide 14: Section - \u2464\u9032\u5316 ──
    set_section_title(slides[14], "6. \u2464\u9032\u5316")

    # ── Slide 15: AI-Human Gap + spiral + extension ──
    slide = slides[15]
    set_title(slide, "AI-Human Gap = \u6700\u9ad8\u306e\u6559\u5e2b\u30c7\u30fc\u30bf",
              "RM\u304cAI\u51fa\u529b\u3092\u4fee\u6b63\u3059\u308b\u305f\u3073\u3001\u305d\u308c\u304c\u6700\u3082\u4fa1\u5024\u3042\u308b\u5b66\u7fd2\u30c7\u30fc\u30bf\u306b\u306a\u308b")

    # Gap diagram
    add_box(slide, "AI\u51fa\u529b",
            Inches(0.42), Inches(1.90), Inches(2.8), Inches(0.40),
            bg_color=ACCENT_BLUE, text_color=WHITE, font_size=11, bold=True)
    add_box(slide, "\u2190 Gap \u2192\n\u3053\u306e\u300c\u5dee\u5206\u300d\u304c\u6700\u9ad8\u306e\u6559\u5e2b\u30c7\u30fc\u30bf",
            Inches(3.42), Inches(1.90), Inches(2.5), Inches(0.40),
            bg_color=ACCENT_ORANGE, text_color=WHITE, font_size=9, bold=True)
    add_box(slide, "RM\u306e\u4fee\u6b63",
            Inches(6.12), Inches(1.90), Inches(2.8), Inches(0.40),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=11, bold=True)

    # Gap quality comparison table
    gap_table = [
        ["\u6559\u5e2b\u30c7\u30fc\u30bf\u306e\u7a2e\u985e", "\u54c1\u8cea", "\u8ffd\u52a0\u30b3\u30b9\u30c8", "\u91cf"],
        ["\u6559\u79d1\u66f8\u7684\u306a\u6b63\u89e3\u30c7\u30fc\u30bf", "\u25cb", "\u9ad8\uff08\u4f5c\u6210\u5de5\u6570\uff09", "\u5c11"],
        ["\u904e\u53bb\u306e\u4e8b\u4f8b\u30c7\u30fc\u30bf", "\u25b3", "\u4e2d\uff08\u6574\u7406\u5de5\u6570\uff09", "\u4e2d"],
        ["AI-Human Gap", "\u2605\u2605\u2605", "\u30bc\u30ed\uff08\u65e5\u5e38\u696d\u52d9\u306e\u526f\u7523\u7269\uff09", "\u5927\uff08\u4f7f\u3046\u307b\u3069\u6e9c\u307e\u308b\uff09"],
    ]
    add_table(slide, gap_table, left=Inches(0.42), top=Inches(2.60), width=Inches(8.5),
              col_widths=[3.0, 1.0, 2.5, 2.0],
              header_font_pt=8, cell_font_pt=8, row_height=0.32)

    # Positive spiral diagram
    add_box(slide, "\u6b63\u306e\u30b9\u30d1\u30a4\u30e9\u30eb",
            Inches(9.30), Inches(1.90), Inches(3.5), Inches(0.30),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=9, bold=True)

    spiral_steps = [
        ("AI\u3092\u4f7f\u3046", GREEN_MINT),
        ("Gap\u767a\u751f", RGBColor(0xFD, 0xE8, 0xC8)),
        ("\u4fee\u6b63\u3059\u308b", GREEN_PALE),
        ("Gap\u304c\u6559\u5e2b\u30c7\u30fc\u30bf\u306b", ACCENT_ORANGE),
        ("AI\u304c\u8ce2\u304f\u306a\u308b", GREEN_MINT),
        ("\u3082\u3063\u3068\u4f7f\u3044\u305f\u304f\u306a\u308b", GREEN_PALE),
        ("\u3055\u3089\u306bGap\u84c4\u7a4d...", RGBColor(0xFD, 0xE8, 0xC8)),
    ]
    y_sp = 2.28
    for i, (label, color) in enumerate(spiral_steps):
        text_c = WHITE if color == ACCENT_ORANGE else CHARCOAL
        add_box(slide, label,
                Inches(9.30), Inches(y_sp), Inches(3.5), Inches(0.24),
                bg_color=color, text_color=text_c, font_size=7, bold=True)
        y_sp += 0.24
        if i < len(spiral_steps) - 1:
            add_arrow_down(slide, Inches(10.8), Inches(y_sp - 0.02), Inches(0.4), Inches(0.15))
            y_sp += 0.12

    # Extension tables
    add_box(slide, "\u62e1\u5f35: \u81ea\u52d5\u9032\u5316\uff08\u8ffd\u52a0\u958b\u767a\u306a\u3057\uff09",
            Inches(0.42), Inches(4.15), Inches(6.2), Inches(0.25),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    auto_evo = [
        ["\u4ed5\u7d44\u307f", "\u4f55\u304c\u8d77\u304d\u308b\u304b", "\u30bf\u30a4\u30df\u30f3\u30b0"],
        ["Decision Trace\u84c4\u7a4d", "\u627f\u8a8d\u30d1\u30bf\u30fc\u30f3 \u2192 \u30eb\u30fc\u30eb\u81ea\u52d5\u62bd\u51fa", "\u65e5\u5e38\u7684"],
        ["\u30c1\u30e3\u30c3\u30c8FB\u84c4\u7a4d", "\u7684\u5916\u308c\u691c\u51fa \u2192 \u30d7\u30ed\u30f3\u30d7\u30c8\u6539\u5584\u5019\u88dc", "\u65e5\u5e38\u7684"],
        ["\u6210\u7d04/\u5931\u6ce8\u84c4\u7a4d", "\u4e8b\u4f8b\u30d1\u30bf\u30fc\u30f3 \u2192 K5\u81ea\u52d5\u66f4\u65b0", "\u5546\u8ac7\u7d50\u679c\u78ba\u5b9a\u6642"],
    ]
    add_table(slide, auto_evo, left=Inches(0.42), top=Inches(4.45), width=Inches(6.2),
              col_widths=[2.2, 2.3, 1.0],
              header_font_pt=7, cell_font_pt=6.5, row_height=0.26)

    add_box(slide, "\u62e1\u5f35: \u4eba\u7684\u4ecb\u5165\u306b\u3088\u308b\u9032\u5316\uff08\u5b9a\u671f\u30ec\u30d3\u30e5\u30fc\uff09",
            Inches(6.90), Inches(4.15), Inches(5.9), Inches(0.25),
            bg_color=GREEN_DARK, text_color=WHITE, font_size=8, bold=True)

    human_evo = [
        ["\u4ed5\u7d44\u307f", "\u4f55\u304c\u8d77\u304d\u308b\u304b", "\u30bf\u30a4\u30df\u30f3\u30b0"],
        ["\u30a8\u30ad\u30b9\u30d1\u30fc\u30c8\u30d2\u30a2\u30ea\u30f3\u30b0", "\u30d9\u30c6\u30e9\u30f3\u306e\u7d4c\u9a13\u5247 \u2192 K1-K8\u521d\u671f\u69cb\u7bc9", "PoC\uff5eMVP\u521d\u671f"],
        ["PJ\u30c1\u30fc\u30e0\u30ec\u30d3\u30e5\u30fc", "\u84c4\u7a4d\u30d1\u30bf\u30fc\u30f3\u306e\u54c1\u8cea\u78ba\u8a8d \u2192 K\u6607\u683c", "\u56db\u534a\u671f\u3054\u3068"],
        ["\u7d44\u7e54\u6a2a\u65ad\u5206\u6790", "\u55b6\u696d\u5e97\u9593\u306e\u30d1\u30bf\u30fc\u30f3\u6bd4\u8f03 \u2192 \u5168\u793e\u30eb\u30fc\u30eb", "\u534a\u671f\u3054\u3068"],
    ]
    add_table(slide, human_evo, left=Inches(6.90), top=Inches(4.45), width=Inches(5.9),
              col_widths=[2.2, 2.5, 1.0],
              header_font_pt=7, cell_font_pt=6.5, row_height=0.26)

    # Save
    output_path = "output/15_tacit_knowledge_lifecycle.pptx"
    new_prs.save(output_path)
    print(f"Saved standard version: {output_path} ({len(slides)} slides)")
    return output_path


if __name__ == "__main__":
    build()
